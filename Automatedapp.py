import streamlit as st
import pandas as pd
from datetime import date
import time
import base64
import inflect
from openpyxl.utils import get_column_letter
from openpyxl.styles import Alignment
from openpyxl import load_workbook
import io
import numpy as np
import re
from PIL import Image
import matplotlib.pyplot as plt
# import spacy
import logging
import warnings
from nltk.corpus import stopwords
import nltk
import os
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Border, Side, Alignment, Font,PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils.dataframe import dataframe_to_rows # Add these imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches
from io import BytesIO
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator

import io, base64


# Streamlit app with a sidebar layout
st.set_page_config(layout="wide")

def strip_client_prefix(df):
    """Remove 'Client-' prefix from all column names and string values in 'Entity' column."""
    df = df.copy()
    df.columns = [c.replace("Client-", "") if isinstance(c, str) else c for c in df.columns]
    if "Entity" in df.columns:
        df["Entity"] = df["Entity"].astype(str).str.replace("Client-", "", regex=False)
    return df

# Function to process the Excel file
def process_excel(file):
    # Initialize Excel writer
    output = BytesIO()
    excel_writer = pd.ExcelWriter(output, engine='openpyxl')
    all_dframes = []
    sheet_results = {}

    # Iterate through each sheet in the uploaded file
    for sheet_name in pd.ExcelFile(file).sheet_names:
        data = pd.read_excel(file, sheet_name=sheet_name)

        # Convert 'unnamed 2' column to numeric and sort by 'unnamed 0' and 'unnamed 2'
        data['unnamed 2'] = pd.to_numeric(data['unnamed 2'], errors='coerce')
        sorted_data = data.sort_values(by=['unnamed 0', 'unnamed 2'], kind='mergesort')
        sorted_data.drop("unnamed 2", axis=1, inplace=True)
        sorted_data['Source'] = ""

        # Process different subsets of data
        df1 = sorted_data[sorted_data['unnamed 0'] == 'c'].drop(columns=["unnamed 0"] + sorted_data.columns[2:].tolist())
        df2 = sorted_data[sorted_data['unnamed 0'] == 'd'].drop(columns=["unnamed 0"] + sorted_data.columns[2:].tolist())
        df3 = sorted_data[sorted_data['unnamed 0'] == 'b'].drop(columns=sorted_data.columns[:2].tolist() + ['Source', 'unnamed 4'])

        # Reset indexes
        df1.reset_index(drop=True, inplace=True)
        df2.reset_index(drop=True, inplace=True)
        df3.reset_index(drop=True, inplace=True)

        # Combine dataframes
        result_1 = pd.concat([df3, df2, df1], axis=1, join='outer')
        result_1.rename({'unnamed 3': 'Headline', 'unnamed 1': 'Summary'}, axis=1, inplace=True)

        # Replace the column names
        s = result_1.columns.to_series()
        s.iloc[2] = 'Source'
        result_1.columns = s

        # Split 'Source' column
        split_data = result_1['Source'].str.split(',', expand=True)
        dframe = pd.concat([result_1, split_data], axis=1)
        dframe.drop('Source', axis=1, inplace=True)
        dframe.rename({0: 'Source', 1: 'Date', 2: 'Words', 3: 'Journalists'}, axis=1, inplace=True)
        dframe['Headline'] = dframe['Headline'].str.replace("Factiva Licensed Content", "").str.strip()

        # Add 'Entity' column
        dframe.insert(dframe.columns.get_loc('Headline'), 'Entity', sheet_name)

        # Replace specific words in 'Journalists' column with 'Bureau News'
        words_to_replace = ['Hans News Service', 'IANS', 'DH Web Desk', 'HT Entertainment Desk', 'Livemint', 
                            'Business Reporter', 'HT Brand Studio', 'Outlook Entertainment Desk', 'Outlook Sports Desk',
                            'DHNS', 'Express News Service', 'TIMES NEWS NETWORK', 'Staff Reporter', 'Affiliate Desk', 
                            'Best Buy', 'FE Bureau', 'HT News Desk', 'Mint SnapView', 'Our Bureau', 'TOI Sports Desk',
                            'express news service', '(English)', 'HT Correspondent', 'DC Correspondent', 'TOI Business Desk',
                            'India Today Bureau', 'HT Education Desk', 'PNS', 'Our Editorial', 'Sports Reporter',
                            'TOI News Desk', 'Legal Correspondent', 'The Quint', 'District Correspondent', 'etpanache',
                            'ens economic bureau', 'Team Herald', 'Equitymaster','Hans India','Motilal Oswal','Our Web Desk','TOI City Desk',
                            'HT Sports Desk','Team Agenda','TOI World Desk','HT US Desk','Pioneer News Service','HT Syndication','Outlook News Desk',
                            'Our Special Correspondent','TOI Tech Desk','ENS ECONOMIC BUREAU', 'N.E.W.S. Desk','BS Reporter','Team ET','Outlook Web Desk',
                            'Entertainment Web Desk','Outlook Brand Studio','Trending Desk','Bureau Newss','Our Web Correspondent','PR Content',
                            'BS REPORTER','Trending Desk','Team Lounge','Our Web Correspondent','Focus','HT Real Estate News','MintGenie Team',
                            'Our Correspondent','LM US Desk','Global Sports Desk','HT Trending Desk','Guest Post','HT Infotainment Desk','TEAM ET','AA Edit',
                            'Guest','Editorial','agencies','DC Web Desk']
        dframe['Journalists'] = dframe['Journalists'].replace(words_to_replace, 'Bureau News', regex=True)
        
        additional_replacements = ['@timesgroup.com', 'TNN']
        dframe['Journalists'] = dframe['Journalists'].replace(additional_replacements, '', regex=True)

        # Fill NaN or spaces in 'Journalists' column
        dframe['Journalists'] = dframe['Journalists'].apply(lambda x: 'Bureau News' if pd.isna(x) or x.isspace() else x)
        dframe['Journalists'] = dframe['Journalists'].str.lstrip()
        # Remove trailing city/brand keywords from Journalists column
        keywords_to_strip = [
            'Forbes India', 'Forbes', 'Mumbai', 'New Delhi', 'Delhi', 
            'Hyderabad', 'Bengaluru', 'Bengalore', 'Chennai', 
            '@timesofindia.com', 'Ahmedabad', 'Kolkata',' @timesofindia.com','Edited by','|',' |','Senior Journalist'
        ]

        # Build regex pattern to remove these words from anywhere in the string
        pattern = r'\s*\b(' + '|'.join(map(re.escape, keywords_to_strip)) + r')\b\s*'
        # First remove the slash between cities, then remove the keywords
        dframe['Journalists'] = dframe['Journalists'].str.replace(
            r'\s*(' + '|'.join(map(re.escape, keywords_to_strip)) + r')(\s*/\s*(' + '|'.join(map(re.escape, keywords_to_strip)) + r'))*', 
            '', regex=True
        ).str.strip()

        # Also clean columns 4 and 5 if they exist
        for col in [4, 5]:
            if col in dframe.columns:
                dframe[col] = dframe[col].str.replace(pattern, ' ', regex=True).str.strip()

        # Read additional data for merging
        data2 = pd.read_excel(r"FActiva Publications.xlsx")
        
        # Merge the current dataframe with additional data
        merged = pd.merge(dframe, data2, how='left', left_on=['Source'], right_on=['Source'])

        # Save the merged data to Excel with the sheet name
        #merged.to_excel(excel_writer, sheet_name=sheet_name, index=False)
        sheet_results[sheet_name] = merged
        
        # Append DataFrame to the list
        all_dframes.append(merged)
    
    # Combine all DataFrames into a single DataFrame
    combined_data = pd.concat(all_dframes, ignore_index=True)

    # Add a serial number column
    combined_data['sr no'] = combined_data.reset_index().index + 1

    # Rearrange columns to have 'sr no' before 'Entity'
    combined_data = combined_data[['sr no', 'Entity'] + [col for col in combined_data.columns if col not in ['sr no', 'Entity']]]

    # Save the combined data to a new sheet
    combined_data.to_excel(excel_writer, sheet_name='Combined_All_Sheets', index=False)

    for sheet_name, merged in sheet_results.items():
        merged.to_excel(excel_writer, sheet_name=sheet_name, index=False)
    
    # Show the processed dataframe in the web app
    st.write(combined_data)


    # Save and return the Excel file
    excel_writer.close()
    output.seek(0)
    return output
    
    
# Streamlit app setup
st.title("Print Excel File Processor & Merger")

# Upload file
uploaded_file = st.file_uploader("Upload an Excel file", type=["xlsx"])

# Process the file if uploaded
if uploaded_file is not None:
    processed_file = process_excel(uploaded_file)
    
    # Download button
    st.download_button(
        label="Download Processed Excel",
        data=processed_file,
        file_name="Processed_Excel.xlsx",

    )

# Function to extract entity name from file path
def extract_entity_name(file_path):
    base_name = os.path.basename(file_path)
    entity_name = base_name.split('_or_')[0].replace("_", " ").split('-')[0].strip()
    return entity_name

#
st.title('Online Excel File Merger & Entity Extractor (TW)')

uploaded_files_kalki = st.file_uploader(
    "Upload your TW Excel files",
    accept_multiple_files=True,
    type=['xlsx'],
    key="kalki_uploader"
)

if uploaded_files_kalki:
    kalki_df = pd.DataFrame()
    
    for uploaded_file in uploaded_files_kalki:
        df = pd.read_excel(uploaded_file)
        
        # Extract Entity from filename
        raw_name = uploaded_file.name
        base_name = raw_name.rsplit('.', 1)[0]
        entity_name = base_name.split('_')[0]
        
        df['Entity'] = entity_name
        cols = ['Entity'] + [col for col in df.columns if col != 'Entity']
        df = df[cols]
        
        kalki_df = pd.concat([kalki_df, df], ignore_index=True)

    st.write("Preview of merged data:")
    st.dataframe(kalki_df)

    # ====================== CREATE EXCEL WITH HYPERLINKS ======================
    output_kalki = BytesIO()
    
    with pd.ExcelWriter(output_kalki, engine='openpyxl') as writer:
        kalki_df.to_excel(writer, index=False, sheet_name='Merged_Data')
    
    # Load the workbook to add hyperlinks
    output_kalki.seek(0)
    wb = load_workbook(output_kalki)
    ws = wb.active

    # Find columns that likely contain URLs
    url_keywords = ['url', 'link', 'website', 'web', 'href']
    hyperlink_columns = []

    for col_idx, col_name in enumerate(kalki_df.columns, start=1):
        col_name_lower = str(col_name).lower()
        if any(keyword in col_name_lower for keyword in url_keywords):
            hyperlink_columns.append(col_idx)
        else:
            # Also check if first few values look like URLs
            sample_values = kalki_df[col_name].dropna().astype(str).head(5)
            if sample_values.str.startswith(('http://', 'https://')).any():
                hyperlink_columns.append(col_idx)

    # Apply hyperlinks
    for col_idx in hyperlink_columns:
        col_letter = get_column_letter(col_idx)
        for row in range(2, ws.max_row + 1):  # Skip header
            cell = ws[f"{col_letter}{row}"]
            url = str(cell.value).strip()
            if url.startswith(('http://', 'https://')):
                cell.hyperlink = url
                cell.font = Font(color="0000FF", underline="single")  # Blue + underline

    # Auto-adjust column widths
    for col in ws.columns:
        max_length = 0
        column = col[0].column_letter
        for cell in col:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column].width = adjusted_width

    # Save final file
    final_output = BytesIO()
    wb.save(final_output)
    final_output.seek(0)

    st.download_button(
        label="📥 Download Merged TW Excel (with Hyperlinks)",
        data=final_output.getvalue(),
        file_name='merged_TW_with_entity.xlsx',
        mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
# Web app title
st.title('Online Excel File Merger & Entity Extractor (Meltwater)')

# File uploader
uploaded_files = st.file_uploader("Upload your Excel files", accept_multiple_files=True, type=['xlsx'])

if uploaded_files:
    final_df = pd.DataFrame()
    
    # Loop through each uploaded file
    for uploaded_file in uploaded_files:
        df = pd.read_excel(uploaded_file)
        
        # Extract the entity name and add it as a new column
        entity_name = extract_entity_name(uploaded_file.name)
        df['Entity'] = entity_name
        
        # Concatenate all the dataframes
        final_df = pd.concat([final_df, df], ignore_index=True)
    
    # Process columns as required
    existing_columns = final_df.columns.tolist()
    influencer_index = existing_columns.index('Influencer')
    country_index = existing_columns.index('Country')
    
    new_order = (
        existing_columns[:influencer_index + 1] +  # All columns up to and including 'Influencer'
        ['Entity', 'Reach', 'Sentiment', 'Keywords', 'State', 'City', 'Engagement','Language'] +  # Adding new columns
        existing_columns[influencer_index + 1:country_index + 1]  # All columns between 'Influencer' and 'Country'
    )
    
    
    # Fill missing values in 'Influencer' column with 'Bureau News'
    final_df['Influencer'] = final_df['Influencer'].fillna('Bureau News')
    final_df['Date'] = (pd.to_datetime(final_df['Date'].str.strip(), format='%d-%b-%Y %I:%M%p',errors='coerce').dt.date.astype(str))  
    
    # Reorder the DataFrame
    final_df = final_df[new_order]
    # Read lookup file from backend
    lookup_df = pd.read_excel('Meltwater Publications Vlookup.xlsx')
    lookup_df = lookup_df[['Source', 'Publication Name', 'Publication Type']].drop_duplicates(subset='Source')

    # Merge on Source column
    final_df = final_df.merge(lookup_df, on='Source', how='left')

    # Reposition columns right after 'Source'
    cols = final_df.columns.tolist()
    cols.remove('Publication Name')
    cols.remove('Publication Type')
    source_index = cols.index('Source')
    cols.insert(source_index + 1, 'Publication Name')
    cols.insert(source_index + 2, 'Publication Type')
    final_df = final_df[cols]
    
    # Show the processed dataframe in the web app
    st.write(final_df)
    
    # Prepare Excel file in memory
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        final_df.to_excel(writer, index=False)
    
    # Convert buffer to bytes
    processed_data = output.getvalue()

    # Option to download the merged file
    st.download_button(
        label="Download Merged Excel",
        data=processed_data,
        file_name='merged_excel_with_entity.xlsx',
        mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Data preprocessing function (You can include your data preprocessing here)

# Function to create separate Excel sheets by Entity

def create_entity_sheets(data, writer):
    cols_to_drop = ["Keywords", "Engagement", "Language", "Country", "Exclusivity", "Topic"]

    if 'Entity' in data.columns:
        cols = ['Entity'] + [col for col in data.columns if col != 'Entity']
        data = data[cols]

    entities = data['Entity'].unique()

    for Entity in entities:
        entity_df = data[data['Entity'] == Entity].copy()
        entity_df['Date'] = entity_df['Date'].dt.date
        entity_df['Journalist'] = entity_df['Journalist'].str.replace(r"[\[\]']", "", regex=True)
        if "Entity" in entity_df.columns:
            entity_df["Entity"] = entity_df["Entity"].astype(str).str.replace("Client-", "", regex=False)

        # ── Drop sr no column if present ──────────────────────────────
        sr_cols = [c for c in entity_df.columns if str(c).strip().lower() in ("sr no", "sr no.", "sr_no", "srno")]
        entity_df.drop(columns=sr_cols, inplace=True, errors="ignore")

        # ── Move Date to 2nd column (right after Entity) ──────────────
        if "Date" in entity_df.columns and "Entity" in entity_df.columns:
            cols_order = ["Entity", "Date"] + [c for c in entity_df.columns if c not in ("Entity", "Date")]
            entity_df = entity_df[cols_order]

        entity_df.drop(columns=[col for col in cols_to_drop if col in entity_df.columns], inplace=True)

        clean_entity_name = str(Entity).replace("Client-", "")[:31]
        entity_df.to_excel(writer, sheet_name=clean_entity_name, index=False)
        worksheet = writer.sheets[clean_entity_name]

        # ── Identify special columns ──────────────────────────────────
        col_names = entity_df.columns.tolist()

        headline_cols = [c for c in col_names if any(
            kw in str(c).lower() for kw in ["headline", "title"]
        )]
        summary_cols = [c for c in col_names if any(
            kw in str(c).lower() for kw in ["summary", "opening text", "hit sentence", "article summary"]
        )]
        url_cols = [c for c in col_names if "url" in str(c).lower()]

        special_cols = set(headline_cols + summary_cols + url_cols)

        # ── Style constants ───────────────────────────────────────────
        hdr_font  = Font(bold=True, name="Calibri")
        data_font = Font(name="Calibri")

        ctr_nowrap = Alignment(horizontal="center", vertical="center", wrap_text=False)
        ctr_wrap   = Alignment(horizontal="center", vertical="center", wrap_text=True)
        top_ctr    = Alignment(horizontal="center", vertical="top",    wrap_text=True)

        # ── Header row: bold + center ─────────────────────────────────
        for cell in worksheet[1]:
            cell.font      = hdr_font
            cell.alignment = ctr_nowrap

        # ── Data rows ─────────────────────────────────────────────────
        for row in worksheet.iter_rows(min_row=2, max_row=worksheet.max_row):
            for cell in row:
                col_name = col_names[cell.column - 1] if cell.column - 1 < len(col_names) else ""
                cell.font = data_font

                if col_name in headline_cols:
                    cell.alignment = ctr_wrap          # wrap + center h + center v
                elif col_name in summary_cols:
                    cell.alignment = top_ctr           # wrap + center h + top v
                else:
                    cell.alignment = ctr_nowrap        # no wrap, center both

        # ── Uniform row height ────────────────────────────────────────
        for rn in range(1, worksheet.max_row + 1):
            worksheet.row_dimensions[rn].height = 30

        # ── Column widths ─────────────────────────────────────────────
        for col_0, col_name in enumerate(col_names):
            col_ltr   = get_column_letter(col_0 + 1)
            col_lower = str(col_name).lower()

            if col_name in headline_cols:
                worksheet.column_dimensions[col_ltr].width = 55

            elif col_name in summary_cols:
                worksheet.column_dimensions[col_ltr].width = 55

            elif col_name in url_cols:
                worksheet.column_dimensions[col_ltr].width = 35

            else:
                # Auto-fit: max of header length vs cell values, capped at 40
                max_len = len(str(col_name))
                for r in worksheet.iter_rows(
                    min_row=2, max_row=worksheet.max_row,
                    min_col=col_0 + 1, max_col=col_0 + 1
                ):
                    for cell in r:
                        try:
                            max_len = max(max_len, len(str(cell.value or "")))
                        except Exception:
                            pass
                worksheet.column_dimensions[col_ltr].width = min(float(max_len) + 2, 40)

        # ── URL hyperlinks ────────────────────────────────────────────
        for url_col in url_cols:
            col_index = col_names.index(url_col) + 1
            col_ltr   = get_column_letter(col_index)
            for rn in range(2, worksheet.max_row + 1):
                cell = worksheet[f"{col_ltr}{rn}"]
                if cell.value and isinstance(cell.value, str) and cell.value.startswith("http"):
                    cell.hyperlink = cell.value
                    cell.style     = "Hyperlink"
                    cell.alignment = ctr_nowrap

def add_entity_info(ws, entity_info, start_row):
    for i, line in enumerate(entity_info.split('\n'), start=1):
        cell = ws.cell(row=start_row + i - 1, column=1)
        cell.value = line
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))
#         cell.alignment = Alignment(horizontal='center')  # Merge and center for all lines
#         ws.merge_cells(start_row=start_row + i - 1, start_column=1, end_row=start_row + i, end_column=5)
        
        # Apply specific formatting for Source line
        if line.startswith('Source:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=3, start_column=1, end_row=3, end_column=5)
            cell.font = Font(color="000000",name="Gill Sans MT")
            
        # Apply specific formatting for Source line
        if line.startswith('Entity:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans MT", bold=True )
            cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
            
        # Apply specific formatting for Source line
        if line.startswith('Time Period of analysis:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans MT")
            
        # Apply specific formatting for Source line
        if line.startswith('News search:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=4, start_column=1, end_row=4, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans MT")
            
def add_styling_to_worksheet(ws, df, start_row, comment, highlight_last_row=False):
    # Apply table heading as comment
    cell = ws.cell(row=start_row, column=1)
    cell.value = comment
    cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
    cell.font = Font(color="000000", bold=True, name="Gill Sans MT")
    cell.alignment = Alignment(horizontal='center')
    ws.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=len(df.columns))
    thin_side = Side(border_style="thin", color="000000")
    for col in range(1, len(df.columns) + 1):
        c = ws.cell(row=start_row, column=col)
        if col == 1:
            c.border = Border(left=thin_side, top=thin_side, bottom=thin_side)
        elif col == len(df.columns):
            c.border = Border(right=thin_side, top=thin_side, bottom=thin_side)
        else:
            c.border = Border(top=thin_side, bottom=thin_side)
    
    # Increment the start row
    start_row += 1

    # Apply styling to column headers
    for col_idx, col_name in enumerate(df.columns, start=1):
        cell = ws.cell(row=start_row, column=col_idx)
        cell.value = col_name
        cell.font = Font(color="000000", bold=True, name="Gill Sans MT")
        cell.alignment = Alignment(horizontal='center')
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))  
        
    start_row += 1

    # Write DataFrame values with styling
    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=False), start=start_row):
        for c_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx)
            if isinstance(value, pd.Period):
                cell.value = value.strftime('%Y-%m') 
            else:
                cell.value = value
            cell.font = Font(name="Gill Sans MT")    
            cell.alignment = Alignment(horizontal='center')

    # Apply borders to all cells
    for row in ws.iter_rows(min_row=start_row, max_row=start_row + len(df) - 1, min_col=1, max_col=len(df.columns)):
        for cell in row:
            cell.border = Border(left=Side(border_style="thin", color="000000"),
                                 right=Side(border_style="thin", color="000000"),
                                 top=Side(border_style="thin", color="000000"),
                                 bottom=Side(border_style="thin", color="000000"))

    # ✅ Bold the last row if required
    if highlight_last_row:
        last_data_row = start_row + len(df) - 1
        for col_idx in range(1, len(df.columns) + 1):
            cell = ws.cell(row=last_data_row, column=col_idx)
            cell.font = Font(name="Gill Sans MT", bold=True)
def kalki_multiple_dfs(df_list, sheet_name, file_name, comments,
                       entity_info, highlight_refs=None):
    wb = Workbook()
    ws = wb.active
    current_row = 1
    add_entity_info(ws, entity_info, current_row)
    current_row += 6

    ref_ids = set(id(r) for r in (highlight_refs or []))

    for df, comment in zip(df_list, comments):
        # ── GUARD: skip empty DataFrames ──────────────────────────────
        if df.empty:
            continue
        # Bold last row if caller flagged this df, OR if last row has "total"
        highlight = (id(df) in ref_ids) or any(
            "total" in str(val).lower() for val in df.iloc[-1]
        )
        add_styling_to_worksheet(ws, df, current_row, comment,
                                 highlight_last_row=highlight)
        current_row += len(df) + 4

    wb.save(file_name)


def kalki_multiple_dfs1(df_list, sheet_name, wb, comments,
                        highlight_refs=None):
    ws = wb.create_sheet(title=sheet_name)
    current_row = 3
    ref_ids = set(id(r) for r in (highlight_refs or []))

    for df, comment in zip(df_list, comments):
        # ── GUARD: skip empty DataFrames ──────────────────────────────
        if df.empty:
            continue
        highlight = (id(df) in ref_ids) or any(
            "total" in str(val).lower() for val in df.iloc[-1]
        )
        add_styling_to_worksheet(ws, df, current_row, comment,
                                 highlight_last_row=highlight)
        current_row += len(df) + 4
def multiple_dfs(df_list, sheet_name, file_name, comments, entity_info):
    wb = Workbook()
    ws = wb.active
    current_row = 1

    # Add entity information to the first 4 rows
    add_entity_info(ws, entity_info, current_row)
    current_row += 6
    for df, comment in zip(df_list, comments):
        if df.empty:
            continue
        highlight = any("total" in str(val).lower() for val in df.iloc[-1])
        add_styling_to_worksheet(ws, df, current_row, comment, highlight_last_row=highlight)
        current_row += len(df) + 4
    wb.save(file_name)

def multiple_dfs1(df_list, sheet_name, wb, comments):
    ws = wb.create_sheet(title=sheet_name)
    current_row = 3

    for df, comment in zip(df_list, comments):
        if df.empty:
            continue
        highlight = any("total" in str(val).lower() for val in df.iloc[-1])
        add_styling_to_worksheet(ws, df, current_row, comment, highlight_last_row=highlight)
        current_row += len(df) + 4


def add_table_to_slide(slide, df, title, textbox_text):
    rows, cols = df.shape
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(14)
    max_table_height = Inches(5)
    total_height_needed = Inches(0.8 * (rows + 1))
    height = max_table_height if total_height_needed > max_table_height else total_height_needed

    # Add title shape (above the table)
    title_shape = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.2))
    title_frame = title_shape.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add the table
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = df.columns[i]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Gill Sans'
                run.font.size = Pt(15)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.values[i, j])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Gill Sans'
                    run.font.size = Pt(15)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Add a text box above the table (shared across all DataFrame slides)
    textbox_left = Inches(0.25)  # Adjust left positioning as needed
    textbox_right = Inches(0.25)
    textbox_top = Inches(0.8)  # Adjust top positioning as needed
    textbox_width = Inches(15.5)  # Adjust width
    textbox_height = Inches(2.1)  # Adjust height

    text_box = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = textbox_text  # The custom text box content for each slide
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(17)  # Adjust the font size as needed
#             run.font.bold = True
            run.font.name = 'Gill Sans'
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text

    # Add the image (footer logo) at the bottom of the slide
    # Add the image (footer logo) at the bottom of the slide
    left = Inches(0.0)
    top = prs.slide_height - Inches(1)
    slide.shapes.add_picture( img_path,left, top, height=Inches(1))  # Adjust as needed





# # Function to save multiple DataFrames in a single Excel sheet
# def multiple_dfs(df_list, sheets, file_name, spaces, comments):
#     writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
#     row = 2
#     for dataframe, comment in zip(df_list, comments):
#         pd.Series(comment).to_excel(writer, sheet_name=sheets, startrow=row,
#                                     startcol=1, index=False, header=False)
#         dataframe.to_excel(writer, sheet_name=sheets, startrow=row + 1, startcol=0)
#         row = row + len(dataframe.index) + spaces + 2
#     writer.close()


# # # Generate an image from the bar chart
# def generate_bar_chart(df):
#     # Filter out unwanted rows
#     df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
#     df = df[df["Entity"] != "Total"]
    
#     # Create the bar chart
#     fig, ax = plt.subplots(figsize=(12, 6))
#     x = range(len(df["Entity"]))  # Define x positions for the bars
#     bars = ax.bar(
#         x, 
#         df["News Count"], 
#         color="skyblue", 
#         edgecolor="black"
#     )
    
#     # Add data labels on top of the bars
#     for bar in bars:
#         height = bar.get_height()
#         ax.text(
#             bar.get_x() + bar.get_width() / 2, 
#             height, 
#             f"{height}", 
#             ha="center", 
#             va="bottom", 
#             fontsize=10
#         )
    
#     # Set chart title and axis labels
#     ax.set_title("Share of Voice (SOV)", fontsize=14)
#     ax.set_xlabel("Entity", fontsize=12)
#     ax.set_ylabel("News Count", fontsize=12)
    
#     # Customize x-axis ticks and labels
#     ax.set_xticks(x)
#     ax.set_xticklabels(df["Entity"], rotation=45, ha="right")
    
#     # Add gridlines for better readability
#     ax.grid(axis="y", linestyle="--", alpha=0.7)
    
   
      
    # # Save plot as image
    # img_path4 = "bar_chart.png"
    # fig.savefig(img_path4, dpi=300)
    # plt.close(fig)
    # return img_path4

def generate_bar_chart(df):
    # Remove 'Client-' prefix from 'Entity' column
    df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
    
    # Filter out unwanted rows
    df = df[df["Entity"] != "Total"]
    
    # Create the bar chart
    fig, ax = plt.subplots(figsize=(12, 6))  # Increase figure width for better label visibility
    x = range(len(df["Entity"]))  # Define x positions for the bars
    bars = ax.bar(
        x, 
        df["News Count"], 
        color="orange", 
        edgecolor="black"
    )
    
    # Add data labels on top of the bars without decimal
    for bar in bars:
        height = int(bar.get_height())  # Convert height to integer
        ax.text(
            bar.get_x() + bar.get_width() / 2, 
            height, 
            f"{height}", 
            ha="center", 
            va="bottom", 
            fontsize=12,
            fontweight="bold"
        )
    
    # Set chart title and axis labels
    # ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Entity", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")
    
    # Customize x-axis ticks and labels for better visibility
    ax.set_xticks(x)
    ax.set_xticklabels(df["Entity"], rotation=45, ha="right", fontsize=12,fontweight="bold")

    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black", which="major", width=1, labelrotation=0)
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")
    
    # Add gridlines for better readability
    ax.grid(axis="y", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path4 = "bar_chart.png"
    fig.savefig(img_path4, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path4
    
def add_image_to_slide(slide, img_path4):
    left = Inches(1)
    top = Inches(1)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path4, left, top, width=width, height=height)


def generate_line_graph(df):
    fig, ax = plt.subplots(figsize=(15, 5.6))
    
    # Exclude the 'Total' column and row for the graph
    filtered_df = df.loc[df['Date'] != 'Total'].copy()
    filtered_df = filtered_df.drop(columns=['Total'], errors='ignore')

    for entity in filtered_df.columns[1:]:  # Exclude the first column (Date)
        ax.plot(filtered_df['Date'].astype(str), filtered_df[entity], marker='o', label=entity)
        for x, y in zip(filtered_df['Date'].astype(str), filtered_df[entity]):
            ax.text(x, y, str(y), fontsize=10, ha='right', va='bottom',fontweight="bold")

    # Set labels and title
    ax.set_xlabel("Month", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")

    # Adjust legend position to avoid overlapping with the graph
    ax.legend(title="Entities", fontsize=10, bbox_to_anchor=(1.05, 1), loc='upper left')

    # Grid and other settings
    ax.grid(axis='y', linestyle='--', alpha=0.7)
    plt.xticks(rotation=45)

    # Use tight_layout to prevent clipping of elements
    plt.tight_layout()

    # Save plot as image
    img_path5 = "line_graph.png"
    fig.savefig(img_path5, dpi=300)
    plt.close(fig)
    return img_path5


def add_image_to_slide1(slide, img_path4):
    left = Inches(1)
    top = Inches(1)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path4, left, top, width=width, height=height)

# Generate bar chart
def generate_bar_pchart(df):
    # Remove 'Client-' prefix from column names
    df.columns = df.columns.str.replace("Client-", "", regex=False)
    
    # # Remove the 'Total' column if it exists
    # if 'Total' in df.columns:
    #     df = df.drop(columns=['Total'])

    # # Remove the 'Total' column if it exists
    # if 'GrandTotal' in df.rows:
    #     df = df.drop(columns=['GrandTotal'])

    # Remove 'Total' and 'GrandTotal' rows and columns
    df = df.loc[~df["Publication Type"].isin(["Total", "GrandTotal"])]
    df = df.drop(columns=["Total", "GrandTotal"], errors="ignore")

    # Plotting
    fig, ax = plt.subplots(figsize=(12, 6))  # Figure size
    bars = df.plot(kind='bar', ax=ax, stacked=False, width=0.8, cmap='Set3',edgecolor="black")  # Plot bars with colormap

    # Add data labels on top of the bars
    for container in ax.containers:
        ax.bar_label(container, fmt='%d', label_type='edge', fontsize=10, padding=3)
    
    # Set chart labels and title
    ax.set_xlabel("Publication Type", fontsize=12, fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12, fontweight="bold")
    # ax.set_title("Hospital Mentions by Publication", fontsize=14, fontweight="bold")
    
    # Customize x-axis labels for better readability
    ax.set_xticklabels(df["Publication Type"], rotation=45, ha="right", fontsize=10, fontweight="bold")
    
    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black")
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")

    # Add legend
    ax.legend(title="Hospitals", bbox_to_anchor=(1.05, 1), loc='upper left')
    
    # Save the plot
    img_path6 = "bar_chart.png"
    fig.savefig(img_path6, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path6
    
def add_image_to_slide2(slide, img_path6):
    left = Inches(1)
    top = Inches(1)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path6, left, top, width=width, height=height)
    
# Function to clean text
def clean_text(text):
    text = text.lower()  # Convert to lowercase
    text = re.sub(r"http\S+|www\S+|https\S+", '', text, flags=re.MULTILINE)  # Remove URLs
    text = re.sub(r'\@\w+|\#', '', text)  # Remove mentions and hashtags
    text = re.sub(r'[^a-zA-Z\s]', '', text)  # Remove non-alphabetic characters
    text = re.sub(r'\s+', ' ', text).strip()  # Remove extra whitespace
    return text
    
# Function to generate word cloud
def generate_word_cloud(df):
    text = ' '.join(df['Headline'].astype(str))
    text = clean_text(text)  # Clean the text
    stopwords = set(STOPWORDS)
    wordcloud = WordCloud(stopwords=stopwords, background_color="white" ,width=550,
        height=450,max_font_size=90, max_words=120,colormap='Set1',collocations=False).generate(text)
    
    # Plotting the word cloud
    fig, ax = plt.subplots(figsize=(6, 6), facecolor = 'black', edgecolor='black')
    ax.imshow(wordcloud, interpolation='bilinear')
    # ax.tight_layout(pad = 0) 
    ax.axis('off')
    
    # Save plot as image
    img_path11 = "wordcloud.png"
    fig.savefig(img_path11, dpi=300, bbox_inches='tight')
    plt.close(fig)
    
    return img_path11

    # # Example usage
    # img_path11 = generate_word_cloud(df)
    # print(f"Word cloud saved at: {img_path11}")

# Function to add image to slide (similar to the example you shared)
def add_image_to_slide11(slide, img_path11):
    from pptx.util import Inches
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(10)  # Adjust width
    height = Inches(6)  # Adjust height
    slide.shapes.add_picture(img_path11, left, top, width=width, height=height)


def top_10_dfs(df_list, file_name, comments, top_11_flags):
    writer = pd.ExcelWriter(file_name, engine='openpyxl')
    row = 2
    for dataframe, comment, top_11_flag in zip(df_list, comments, top_11_flags):
        if top_11_flag:
            top_df = dataframe.head(50)  # Select the top 11 rows for specific DataFrames
        else:
            top_df = dataframe  # Leave other DataFrames unchanged

        top_df.to_excel(writer, sheet_name="Top 10 Data", startrow=row, index=True)
        row += len(top_df) + 2  # Move the starting row down by len(top_df) + 2 rows

    # Create a "Report" sheet with all the DataFrames
    for dataframe, comment in zip(df_list, comments):
        dataframe.to_excel(writer, sheet_name="Report", startrow=row, index=True, header=True)
        row += len(dataframe) + 2  # Move the starting row down by len(dataframe) + 2 rows

    writer.close()    
    

# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("Data Insights/Tables Dashboard")

def format_pretty_date(d):
    day = d.day
    suffix = 'th' if 11 <= day <= 13 else {1: 'st', 2: 'nd', 3: 'rd'}.get(day % 10, 'th')
    return f"{day}{suffix} {d.strftime('%B %Y')}"

# Fixed min and max
min_date = date(2023, 1, 1)
max_date = date.today()

with st.sidebar:
    st.title("Enter Date Range here")
    st.markdown("**Select Date Range**")

    START_DATE = st.date_input("Start Date", value=min_date, min_value=min_date, max_value=max_date, key="start_date")
    END_DATE = st.date_input("End Date", value=max_date, min_value=min_date, max_value=max_date, key="end_date")
    start_date = format_pretty_date(START_DATE)
    end_date = format_pretty_date(END_DATE)

    date_selected = False

    if START_DATE > END_DATE:
        st.error("Start date cannot be after end date.")
    elif START_DATE == END_DATE:
        st.error("Start and end dates cannot be the same.")
    else:
        st.success(f"**Start:** {format_pretty_date(START_DATE)} | **End:** {format_pretty_date(END_DATE)}")
        date_selected = True
# Sidebar for file upload and download options
if date_selected:
    st.sidebar.write("## Provide Client's Industry")

    # --- INDUSTRY NAME (MANDATORY) ---
    industry_input = st.sidebar.text_input(
        "Industry Name*",
    )

    if not industry_input.strip():
        st.sidebar.warning("Please enter the Industry Name to proceed.")
        industry_provided = False
    else:
        industry_provided = True
        industry = industry_input.strip()
        
if date_selected and industry_provided :
    st.sidebar.write("## Upload TW Online file for tables")
    kalki_online_file = st.sidebar.file_uploader(
        "Upload TW Data File (Excel or CSV)",
        type=["xlsx", "csv"],
        key="kalki_online_tables_uploader",
    )
     
    if kalki_online_file:
        st.sidebar.write("TW File Uploaded Successfully!")
        kalki_data_raw = load_data(kalki_online_file)
     
        if kalki_data_raw is not None:
     
            # ── 1. BASIC CLEAN-UP ────────────────────────────────────────
            kdata = kalki_data_raw.copy()
            kdata.drop(columns=kdata.columns[20:], inplace=True, errors="ignore")
     
            # Standardise journalist/author column → "Author"
            for _src in ("Influencer", "Journalist"):
                if _src in kdata.columns:
                    kdata = kdata.rename(columns={_src: "Author"})
                    break
     
            # De-duplicate using whichever key columns exist
            for _subset in [
    ["Date", "Entity", "Title", "Publication"],
    ["Date", "Entity", "Headline", "Publication Name"],  # keep as fallback
]:
                if set(_subset).issubset(kdata.columns):
                    kdata.drop_duplicates(subset=_subset, keep="first",
                                          inplace=True, ignore_index=True)
     
            kfinaldata = kdata.copy()
            kfinaldata["Date"] = pd.to_datetime(kfinaldata["Date"]).dt.normalize()
            # For unique article counts (Total Unique Articles + Client %)
           
     
            # ── 2. CLIENT / COMPETITOR DETECTION ─────────────────────────
            k_client_col = next(
                (e for e in kfinaldata["Entity"].unique() if e.startswith("Client-")),
                None,
            )
     
            # ── 3. SOV TABLE ──────────────────────────────────────────────
            k_sov_raw = pd.crosstab(
                kfinaldata["Entity"],
                columns="News Count",
                values=kfinaldata["Entity"],
                aggfunc="count",
            ).round(0)
            k_sov_raw["% "] = (
    (k_sov_raw["News Count"] / k_sov_raw["News Count"].sum()) * 100
).round(0)
            k_sov_raw = k_sov_raw.sort_values("News Count", ascending=False)
            k_sov_raw.loc["Total"] = k_sov_raw.sum(numeric_only=True)

            k_Entity_SOV3 = pd.DataFrame(k_sov_raw.to_records()).round()
            k_Entity_SOV3["News Count"] = k_Entity_SOV3["News Count"].astype(int)

            # ── Fix: set Total % row to exactly 100% instead of summed rounded values ──
            k_Entity_SOV3.loc[k_Entity_SOV3["Entity"] == "Total", "% "] = 100

            k_Entity_SOV3["% "] = k_Entity_SOV3["% "].astype(int).astype(str) + "%"
     
            k_sov_order = k_Entity_SOV3["Entity"].tolist()
            k_sov_order_no_client = [
                e for e in k_sov_order
                if not e.startswith("Client-") and e != "Total"
            ]
            k_client_coldt = k_client_col  # shorthand alias
     
            # ── 4. MONTH-ON-MONTH TABLE ───────────────────────────────────
            k_sov_dt = pd.crosstab(
                kfinaldata["Date"].dt.to_period("M"),
                kfinaldata["Entity"],
                margins=True,
                margins_name="Total",
            )
            k_sov_dt1 = pd.DataFrame(k_sov_dt.to_records())
            k_mom_cols = (
                ["Date", k_client_coldt]
                + [e for e in k_sov_order_no_client if e in k_sov_dt1.columns]
                + (["Total"] if "Total" in k_sov_dt1.columns else [])
            )
            k_sov_dt11 = k_sov_dt1[k_mom_cols]
     
            # ── 5. PUBLICATION TABLE ──────────────────────────────────────
            # Column name: "Publication"  |  Total articles (not unique)  |  No Client %
            kfinaldata_non_exploded = kfinaldata.copy()
     
            # Rename "Publication Name" → "Publication" everywhere
            if "Publication Name" in kfinaldata_non_exploded.columns:
                kfinaldata_non_exploded = kfinaldata_non_exploded.rename(
                    columns={"Publication Name": "Publication"}
                )
            if "Publication Name" in kfinaldata.columns:
                kfinaldata = kfinaldata.rename(
                    columns={"Publication Name": "Publication"}
                )
            kfinaldata_unique = kfinaldata.copy()

            # Apply ALL matching dedup subsets, not just the first one
            for _usubset in [
                ["Date", "Title", "Publication", "Author"],
                ["Date", "Headline", "Publication", "Author"],
                ["Date", "Title", "Publication Name", "Author"],      # fallback if rename didn't apply
                ["Date", "Headline", "Publication Name", "Author"],   # fallback
            ]:
                if set(_usubset).issubset(kfinaldata_unique.columns):
                    kfinaldata_unique.drop_duplicates(
                        subset=_usubset, keep="first", inplace=True, ignore_index=True
                    )
                    break
            # Total UNIQUE articles per publication
            k_pub_unique_total = (
                kfinaldata_unique["Publication"]
                .value_counts()
                .reset_index()
            )
            k_pub_unique_total.columns = ["Publication", "Total Unique Articles"]

            # Entity-wise TOTAL articles (not unique)
            k_pub_cross = pd.crosstab(
                kfinaldata["Publication"],
                kfinaldata["Entity"],
            ).reset_index()

            k_pubs_table = k_pub_unique_total.merge(k_pub_cross, on="Publication", how="left")

            k_pub_cols = (
                ["Publication", k_client_coldt]
                + [e for e in k_sov_order_no_client if e in k_pubs_table.columns]
                + ["Total Unique Articles"]
            )
            k_pubs_table = k_pubs_table[k_pub_cols]
            k_pubs_table = k_pubs_table.sort_values("Total Unique Articles", ascending=False).round()
            k_pubs_table.loc["Total"] = k_pubs_table.sum(numeric_only=True)
            k_pubs_table["Publication"] = k_pubs_table["Publication"].fillna("Total")
            _k_num_pub = k_pubs_table.select_dtypes(include=["number"]).columns
            k_pubs_table[_k_num_pub] = k_pubs_table[_k_num_pub].astype(int)

            # Add Client %
            k_pubs_table["Client %"] = (
                (k_pubs_table[k_client_coldt] / k_pubs_table["Total Unique Articles"]) * 100
            ).round().astype(int)

            k_pubs_table20 = k_pubs_table.head(20).copy()
     
            # ── 6. AUTHOR TABLE ───────────────────────────────────────────
            # Column name: "Author"  |  Explode on comma  |  Total articles  |  No Client %
            # Explode on Author for TOTAL article counts (entity-wise columns)
            k_exploded = kfinaldata.copy()          # ← already correct, kfinaldata is the SOV source
            k_exploded["Author"] = (
    k_exploded["Author"].fillna("Bureau News").astype(str).str.split(",")
    .apply(lambda lst: [j.strip() for j in lst] if isinstance(lst, list) else ["Bureau News"])
)
            k_exploded = k_exploded.explode("Author").reset_index(drop=True)

            k_jr_cross = pd.crosstab(k_exploded["Author"], k_exploded["Entity"]).reset_index()

            # Explode on Author for UNIQUE article counts (Total Unique Articles column)
            k_exploded_unique = kfinaldata_unique.copy()
            k_exploded_unique["Author"] = (
    k_exploded_unique["Author"].fillna("Bureau News").astype(str).str.split(",")
    .apply(lambda lst: [j.strip() for j in lst] if isinstance(lst, list) else ["Bureau News"])
)
            k_exploded_unique = k_exploded_unique.explode("Author").reset_index(drop=True)

            k_unique_total = (
                k_exploded_unique["Author"].value_counts().reset_index()
            )
            k_unique_total.columns = ["Author", "Total Unique Articles"]

            # Publication lookup (first occurrence per author)
            k_pub_for_jour = (
                k_exploded[["Author", "Publication"]]
                .drop_duplicates(subset=["Author"], keep="first")
            )

            k_Jour_raw = pd.merge(k_jr_cross, k_pub_for_jour, on="Author", how="left")
            k_Jour_raw = pd.merge(k_Jour_raw, k_unique_total, on="Author", how="left")
            k_Jour_raw["Total Unique Articles"] = k_Jour_raw["Total Unique Articles"].fillna(0).astype(int)

            # Bureau News to bottom, then GrandTotal
            k_bn = k_Jour_raw[k_Jour_raw["Author"] == "Bureau News"]
            k_Jour_raw = k_Jour_raw[k_Jour_raw["Author"] != "Bureau News"]
            k_Jour_raw = pd.concat([k_Jour_raw, k_bn], ignore_index=True)
            k_Jour_raw.loc["GrandTotal"] = k_Jour_raw.sum(numeric_only=True)

            _k_int_cols = k_Jour_raw.columns.difference(["Author", "Publication"])
            k_Jour_raw[_k_int_cols] = k_Jour_raw[_k_int_cols].astype(int)
            k_Jour_raw.insert(1, "Publication", k_Jour_raw.pop("Publication"))

            # Add Client %
            k_Jour_raw["Client %"] = (
                (k_Jour_raw[k_client_coldt] / k_Jour_raw["Total Unique Articles"]) * 100
            ).round().fillna(0).astype(int)

            k_jour_cols = (
                ["Author", "Publication", k_client_coldt]
                + [e for e in k_sov_order_no_client if e in k_Jour_raw.columns]
                + ["Total Unique Articles", "Client %"]
            )
            k_Jour_table = k_Jour_raw[k_jour_cols].copy()

            # ── Sort by Total Unique Articles descending (exclude GrandTotal row) ──
            # REPLACE WITH:
            _k_gt_row = k_Jour_table.loc[k_Jour_table.index == "GrandTotal"]
            _k_bn_row = k_Jour_table.loc[
                k_Jour_table["Author"].astype(str).str.strip().str.lower() == "bureau news"
            ]
            _k_body   = k_Jour_table.loc[
                (k_Jour_table.index != "GrandTotal") &
                (~k_Jour_table["Author"].astype(str).str.strip().str.lower().isin(["bureau news"]))
            ]
            _k_body   = _k_body.sort_values("Total Unique Articles", ascending=False)
            k_Jour_table = pd.concat([_k_body, _k_bn_row, _k_gt_row])
            # ── FIX: Set Author label on GrandTotal row then override entity columns ──
            # The .loc["GrandTotal"] = sum() sets index label but NOT Author column value
            # So we must explicitly set it first, then find it by index label
            k_Jour_table.loc["GrandTotal", "Author"] = "GrandTotal"

            # Now build SOV lookup and override entity columns in GrandTotal row
            k_sov_lookup = dict(zip(k_Entity_SOV3["Entity"], k_Entity_SOV3["News Count"]))

            # Override each entity column with exact SOV count
            for _col in k_Jour_table.columns:
                if _col in k_sov_lookup and _col not in ["Total Unique Articles", "Client %", "Author", "Publication"]:
                    k_Jour_table.loc["GrandTotal", _col] = int(k_sov_lookup[_col])

            # Set Total Unique Articles for GrandTotal = total unique articles in dataset
            _k_grand_unique_total = int(
    k_Jour_table.loc[
        ~k_Jour_table.index.isin(["GrandTotal"])
        & ~k_Jour_table["Author"].astype(str).str.strip().str.lower().isin(["grandtotal"])
    , "Total Unique Articles"]
    .sum()
)
            k_Jour_table.loc["GrandTotal", "Total Unique Articles"] = _k_grand_unique_total

            # Recalculate Client % for GrandTotal
            _grand_client_val = k_Jour_table.loc["GrandTotal", k_client_coldt]
            if _k_grand_unique_total > 0:
                k_Jour_table.loc["GrandTotal", "Client %"] = int(
                    round((_grand_client_val / _k_grand_unique_total) * 100)
                )
            else:
                k_Jour_table.loc["GrandTotal", "Client %"] = 0

            # Refresh top-20 AFTER the fix
            k_Jour_table20 = k_Jour_table.head(20).copy()
            # ── FIX: Override GrandTotal row entity columns to match SOV exactly ──
            
           
            # ── 7. JOURN ON COMP, NOT CLIENT ─────────────────────────────
            k_client_cols_list = [
    c for c in k_Jour_table.columns if c.startswith("Client-")
]
            k_comp_cols_list = [
                c for c in k_Jour_table.columns
                if not c.startswith("Client-")
                and c not in ["Author", "Publication", "Total",
                              "Total Unique Articles", "Client %"]
                and k_Jour_table[c].dtype in ['int64', 'float64']  # only real entity columns
            ]
     
            k_Jour_Comp = k_Jour_table[
    k_Jour_table[k_client_cols_list].eq(0).any(axis=1)
    & ~k_Jour_table["Author"].astype(str).str.strip().str.lower().isin(["grandtotal", "bureau news"])
].sort_values("Total Unique Articles", ascending=False).head(10).copy()
            k_jcomp_cols = (
    ["Author", "Publication", k_client_coldt]
    + [e for e in k_sov_order_no_client if e in k_Jour_Comp.columns]
    + ["Total Unique Articles"]   # ← add this
)
            k_Jour_Comp = k_Jour_Comp[k_jcomp_cols]
     
            # ── 8. JOURN ON CLIENT, NOT COMP ─────────────────────────────
            k_Jour_filtered = k_Jour_table[
    ~k_Jour_table["Author"].astype(str).str.strip().str.lower().isin(
        ["bureau news", "grandtotal", "total"]
    )
]

            k_Jour_Client = k_Jour_filtered[
                (k_Jour_filtered[k_client_cols_list].gt(0).any(axis=1))
                & (k_Jour_filtered[k_comp_cols_list].eq(0).all(axis=1))
            ].head(10).copy()

            k_jclient_cols = (
                ["Author", "Publication", k_client_coldt]
                + [e for e in k_sov_order_no_client if e in k_Jour_Client.columns]
            )
            k_Jour_Client = k_Jour_Client[k_jclient_cols]
            

            #k_Jour_Client = k_Jour_Client[k_jclient_cols]
            k_Jour_table_display   = k_Jour_table.rename(columns={"Author": "Journalist"})
            k_Jour_table20_display = k_Jour_table20.rename(columns={"Author": "Journalist"})
            k_Jour_Comp_display    = k_Jour_Comp.rename(columns={"Author": "Journalist"})
            k_Jour_Client_display  = k_Jour_Client.rename(columns={"Author": "Journalist"})
     
            # ── 9. PREVIEW ────────────────────────────────────────────────
            
            k_preview_options = {
                "SOV Table":            strip_client_prefix(k_Entity_SOV3),
                "Month-on-Month":           strip_client_prefix(k_sov_dt11),
                "Publication Table":        strip_client_prefix(k_pubs_table),
                "Journalist Table":             strip_client_prefix(k_Jour_table),
                "Journ on Comp, not Client":strip_client_prefix(k_Jour_Comp),
                "Journ on Client, not Comp":strip_client_prefix(k_Jour_Client),
            }
            k_sel = st.selectbox(
                "Select DataFrame to Preview (Kalki):",
                list(k_preview_options.keys()),
                key="kalki_preview_select",
            )
            st.dataframe(k_preview_options[k_sel])
            st.sidebar.write("## Download Options")
            # ── 10. DOWNLOAD COMBINED EXCEL ───────────────────────────────
            st.sidebar.write("## Download TW Combined Excel")
            kalki_file_name = st.sidebar.text_input(
                "File name for TW Combined Excel",
                "TW_Combined_Excel.xlsx",
            )
            k_client_name_clean = (
                    k_client_col.replace("Client-", "") if k_client_col else ""
                )
     
            if st.sidebar.button("Download TW Combined Excel"):
     
                k_client_name_clean = (
                    k_client_col.replace("Client-", "") if k_client_col else ""
                )
                k_entity_info = (
                    f"Entity:{k_client_name_clean}\n"
                    f"Time Period of analysis: {start_date} to {end_date}\n"
                    "Source: (Online) Talkwalker All publications.\n"
                    "News search: All Articles: entity mentioned at least once in the article"
                )
     
                # DataFrames for the Report sheet — built after rename below
                k_comments_report = [
                    "SOV Table",
                    "Month-on-Month Table",
                    "Publication Table",
                    "Journalist Table",
                    "Journ-on Comp, not Client",
                    "Journ-on Client, not Comp",
                ]
                # ── Rename "Author" → "Journalist" in all display copies ─────

                k_Jour_table20_report = k_Jour_table20_display.reset_index(drop=True).copy()
                # DataFrames for the Report sheet (top-20 slices, renamed)
                k_dfs_report = [
                    strip_client_prefix(k_Entity_SOV3),
                    strip_client_prefix(k_sov_dt11),
                    strip_client_prefix(k_pubs_table20),
                    strip_client_prefix(k_Jour_table20_display),
                    strip_client_prefix(k_Jour_Comp_display),
                    strip_client_prefix(k_Jour_Client_display),
                ]

                # Pass these objects explicitly – their last rows get bolded
                k_highlight_report = [
                    k_Entity_SOV3,           # Total row
                    k_sov_dt11,              # Total row
                    
                ]

                # ── Build Report sheet ────────────────────────────────────
                k_excel_io = io.BytesIO()
                kalki_multiple_dfs(
                    k_dfs_report,
                    "Tables",
                    k_excel_io,
                    k_comments_report,
                    k_entity_info,
                    highlight_refs=k_highlight_report,
                )
                k_excel_io.seek(0)
                k_wb = load_workbook(k_excel_io)

                # ── Add All Pub-Jour sheet ────────────────────────────────
                k_pubs_all = k_pubs_table.copy()
                k_pubs_all.at[k_pubs_all.index[-1], "Publication"] = "Total"
                k_Jour_table_display.loc[k_Jour_table_display['Journalist'] == 'Bureau News', 'Publication'] = 'Across Publications'

                # ── Separate GrandTotal row before sorting ────────────────────────
                # REPLACE WITH:
                _k_all_gt   = k_Jour_table_display[
                    k_Jour_table_display["Journalist"].astype(str).str.strip().str.lower() == "grandtotal"
                ]
                _k_all_bn   = k_Jour_table_display[
                    k_Jour_table_display["Journalist"].astype(str).str.strip().str.lower() == "bureau news"
                ]
                _k_all_body = k_Jour_table_display[
                    ~k_Jour_table_display["Journalist"].astype(str).str.strip().str.lower().isin(
                        ["grandtotal", "bureau news"]
                    )
                ]
                _k_all_body = _k_all_body.sort_values("Total Unique Articles", ascending=False)
                k_Jour_all_display = pd.concat([_k_all_body, _k_all_bn, _k_all_gt]).reset_index(drop=True).copy()

                # Safety check — ensure GrandTotal label is present on last row
                if str(k_Jour_all_display.iloc[-1]["Journalist"]).strip().lower() != "grandtotal":
                    k_Jour_all_display.iloc[-1, k_Jour_all_display.columns.get_loc("Journalist")] = "GrandTotal"
                k_dfs_all      = [strip_client_prefix(k_pubs_all), strip_client_prefix(k_Jour_table_display)]
                k_comments_all = ["Publication Table",   "Journalist Table"]
                k_highlight_all = [strip_client_prefix(k_pubs_all), strip_client_prefix(k_Jour_table_display)]
     
                kalki_multiple_dfs1(
                    k_dfs_all,
                    "All Pub-Jour",
                    k_wb,
                    k_comments_all,
                    highlight_refs=k_highlight_all,
                )
     
                k_excel_io2 = io.BytesIO()
                k_wb.save(k_excel_io2)
                k_excel_io2.seek(0)
     
                # ── Add entity-wise sheets (raw uploaded data, split by Entity) ──
                # No renaming, no reformatting — data written exactly as uploaded.
                kraw_for_entity = kalki_data_raw.copy()
     
                with pd.ExcelWriter(
                    k_excel_io2, engine="openpyxl",
                    mode="a", if_sheet_exists="new"
                ) as k_writer:
                    for k_entity_name in kraw_for_entity["Entity"].unique():
                        k_entity_df = kraw_for_entity[
                            kraw_for_entity["Entity"] == k_entity_name
                        ].reset_index(drop=True)
                        # Strip "Client-" from Entity column values
                        if "Entity" in k_entity_df.columns:
                            k_entity_df["Entity"] = k_entity_df["Entity"].astype(str).str.replace("Client-", "", regex=False)
                        
                        # Strip "Client-" from sheet name
                        k_sheet_name = str(k_entity_name).replace("Client-", "")[:31]
                        if "Date" in k_entity_df.columns:
                            k_entity_df["Date"] = pd.to_datetime(k_entity_df["Date"]).dt.date
                        k_entity_df.to_excel(
                            k_writer, sheet_name=k_sheet_name, index=False
                        )

                        ws_k   = k_writer.sheets[k_sheet_name]
                        cols_k = k_entity_df.columns.tolist()

                        # ── Identify special columns (case-insensitive) ──────
                        def _col_letter_if(keywords):
                            """Return get_column_letter index for first matching col, else None."""
                            for i, c in enumerate(cols_k, 1):
                                if any(kw in str(c).lower() for kw in keywords):
                                    return i
                            return None

                        title_idx   = _col_letter_if(["title", "headline"])
                        summary_idx = _col_letter_if(["summary", "opening text",
                                                       "hit sentence", "article summary"])
                        url_idx     = _col_letter_if(["url"])

                        # Columns to SKIP during auto-width (0-based col positions)
                        skip_auto_idx = set()
                        for i, c in enumerate(cols_k):
                            c_lower = str(c).lower()
                            if any(kw in c_lower for kw in [
                                "title", "headline", "summary",
                                "opening text", "hit sentence",
                                "article summary", "url"
                            ]):
                                skip_auto_idx.add(i)

                        # ── Style objects ────────────────────────────────────
                        
                        hdr_font     = Font(bold=True, name="Calibri")
                        data_font    = Font(name="Calibri")
                        ctr          = Alignment(horizontal="center",
                                                 vertical="center",
                                                 wrap_text=False)
                        ctr_wrap     = Alignment(horizontal="center",
                                                 vertical="center",
                                                 wrap_text=True)

                        # ── Header row (row 1): yellow + bold + center ───────
                        for cell in ws_k[1]:
                            cell.font      = hdr_font
                            cell.alignment = ctr

                        # ── Data rows: center align; wrap only Title column ──
                        for row in ws_k.iter_rows(min_row=2, max_row=ws_k.max_row):
                            for cell in row:
                                cell.font      = data_font
                                cell.alignment = (
                                    ctr_wrap if cell.column == title_idx else ctr
                                )

                        # ── Uniform row height (all rows including header) ───
                        for rn in range(1, ws_k.max_row + 1):
                            ws_k.row_dimensions[rn].height = 30

                        # ── Column widths ────────────────────────────────────
                        for col_0, col_name in enumerate(cols_k):
                            col_ltr   = get_column_letter(col_0 + 1)
                            col_lower = str(col_name).lower()

                            if any(kw in col_lower for kw in ["title", "headline"]):
                                # Wide + wrap text already applied to data cells above
                                ws_k.column_dimensions[col_ltr].width = 55

                            elif any(kw in col_lower for kw in [
                                "summary", "opening text",
                                "hit sentence", "article summary"
                            ]):
                                # Keep narrow — do not auto-expand
                                ws_k.column_dimensions[col_ltr].width = 35

                            elif "url" in col_lower:
                                # Fixed width for URL column
                                ws_k.column_dimensions[col_ltr].width = 35

                            else:
                                # Auto-fit: max of header length vs all cell values
                                max_len = len(str(col_name))
                                for r in ws_k.iter_rows(
                                    min_row=2, max_row=ws_k.max_row,
                                    min_col=col_0 + 1, max_col=col_0 + 1
                                ):
                                    for cell in r:
                                        try:
                                            max_len = max(
                                                max_len,
                                                len(str(cell.value or ""))
                                            )
                                        except Exception:
                                            pass
                                # Cap at 40 to prevent absurdly wide columns
                                ws_k.column_dimensions[col_ltr].width = min(
                                    float(max_len) + 2, 40
                                )

                        # ── URL column: make each cell a clickable hyperlink ─
                        if url_idx:
                            url_ltr = get_column_letter(url_idx)
                            for rn in range(2, ws_k.max_row + 1):
                                cell = ws_k[f"{url_ltr}{rn}"]
                                if (cell.value
                                        and isinstance(cell.value, str)
                                        and cell.value.startswith("http")):
                                    cell.hyperlink = cell.value
                                    cell.style     = "Hyperlink"
                                    # Re-apply center after Hyperlink style resets alignment
                                    cell.alignment = ctr

                    k_writer.book.worksheets[0].title = "Report"
     
                # ── Re-order sheets ───────────────────────────────────────
                k_wb_final   = load_workbook(k_excel_io2)
                k_all_sheets = k_wb_final.sheetnames
                k_client_sheet = next(
    (s for s in k_all_sheets if s == k_client_name_clean), None
) or next(
    (s for s in k_all_sheets if s.startswith("Client-")), None
)
                k_ordered = ["Report", "All Pub-Jour"]
                if k_client_sheet:
                    k_ordered.append(k_client_sheet)
                k_ordered += [s for s in k_sov_order_no_client if s in k_all_sheets]
                k_ordered += [s for s in k_all_sheets if s not in k_ordered]
                k_wb_final._sheets = [k_wb_final[s] for s in k_ordered]
     
                k_final_io = io.BytesIO()
                k_wb_final.save(k_final_io)
                k_final_io.seek(0)
     
                k_b64  = base64.b64encode(k_final_io.read()).decode()
                k_href = (
                    f'<a href="data:application/vnd.openxmlformats-officedocument'
                    f'.spreadsheetml.sheet;base64,{k_b64}" '
                    f'download="{kalki_file_name}">Download TW Combined Excel</a>'
                )
                
                # ── Build Kalki presentation ────────────────────────────────────────
                st.sidebar.markdown(k_href, unsafe_allow_html=True)

            # ← This closing parenthesis ends the "Download Kalki Combined Excel" if-block
            # ── Build Kalki presentation ──────────────────────────────────────
            st.sidebar.write("## Download TW DataFrames as a PowerPoint File")
            k_pptx_file_name = st.sidebar.text_input(
                "Enter file name for TW PowerPoint",
                "TW_dataframes_presentation.pptx",
                key="kalki_pptx_name"
            )

            if st.sidebar.button("Download TW PowerPoint", key="kalki_pptx_button"):
                def add_kalki_table_to_slide(slide, df, title, textbox_text, prs_obj):
                    rows, cols = df.shape
                    left = Inches(0.8)
                    top = Inches(1.5)
                    width = Inches(14)
                    max_table_height = Inches(5)
                    total_height_needed = Inches(0.8 * (rows + 1))
                    height = max_table_height if total_height_needed > max_table_height else total_height_needed
                    title_shape = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.2))
                    title_frame = title_shape.text_frame
                    title_frame.text = title
                    for paragraph in title_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(28)
                            run.font.bold = True
                            run.font.name = 'Helvetica'
                            run.font.color.rgb = RGBColor(240, 127, 9)
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
                    for i in range(cols):
                        cell = table.cell(0, i)
                        cell.text = str(df.columns[i])
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = 'Gill Sans'
                                run.font.size = Pt(15)
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(0, 0, 0)
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    for i in range(rows):
                        for j in range(cols):
                            cell = table.cell(i + 1, j)
                            cell.text = str(df.values[i, j])
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = 'Gill Sans'
                                    run.font.size = Pt(15)
                            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    
                    left_logo = Inches(0.0)
                    top_logo = prs_obj.slide_height - Inches(1)
                    if os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, left_logo, top_logo, height=Inches(1))

                                # ── img paths (reuse from Online/Print section) ─────────────────
                img_path  = r"New logo snip.png"
                img_path1 = r"New Templete main slide.png"

                # ── CREATE PRESENTATION ─────────────────────────────────────────
                k_prs = Presentation()
                k_prs.slide_width = Inches(16)
                k_prs.slide_height = Inches(9)

                # Slide 1 — Title
                slide_layout = k_prs.slide_layouts[0]
                slide = k_prs.slides.add_slide(slide_layout)
                if os.path.exists(img_path1):
                    slide.shapes.add_picture(img_path1, Inches(0), Inches(0), width=k_prs.slide_width, height=k_prs.slide_height)
                tb = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), Inches(15), Inches(1))
                tf = tb.text_frame
                tf.text = f"{k_client_name_clean}\nNews Analysis\nBy Media Research & Analytics Team"
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(50)
                        run.font.bold = True
                        run.font.name = 'Helvetica'
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.alignment = PP_ALIGN.LEFT

                # Slide 2 — Parameters
                slide_layout = k_prs.slide_layouts[6]
                slide = k_prs.slides.add_slide(slide_layout)
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(0.0), k_prs.slide_height - Inches(1), height=Inches(1))
                for shape in slide.placeholders:
                    if shape.has_text_frame:
                        shape.text_frame.clear()
                hdr = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.7))
                hdr.text_frame.text = "Parameters"
                for paragraph in hdr.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(30)
                        run.font.bold = True
                        run.font.name = 'Helvetica'
                        run.font.color.rgb = RGBColor(240, 127, 9)
                        paragraph.alignment = PP_ALIGN.CENTER
                tp = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(14), Inches(0.5))
                tp.text_frame.text = f"Time Period : {start_date} to {end_date}"
                tp.text_frame.paragraphs[0].font.size = Pt(24)
                tp.text_frame.paragraphs[0].font.name = 'Gill Sans'
                src = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))
                src.text_frame.word_wrap = True

                p_src = src.text_frame.add_paragraph()
                p_src.text = "Source: (Online) Kalki All publications."
                p_src.font.size = Pt(24)
                p_src.font.name = 'Gill Sans'

                # Add this new paragraph
                p_src2 = src.text_frame.add_paragraph()
                p_src2.text = "(Print) Factiva, 25 leading publications including Financial and mainline newspapers, Business and General Magazines."
                p_src2.font.size = Pt(24)
                p_src2.font.name = 'Gill Sans'
                ns = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(10), Inches(0.75))
                ns.text_frame.word_wrap = True
                p_ns = ns.text_frame.add_paragraph()
                p_ns.text = "News Search : All Articles: entity mentioned at least once in the article"
                p_ns.font.size = Pt(24)
                p_ns.font.name = 'Gill Sans'

                # Slide 3 — "Online Media" divider
                slide_layout = k_prs.slide_layouts[0]
                slide = k_prs.slides.add_slide(slide_layout)
                if os.path.exists(img_path1):
                    slide.shapes.add_picture(img_path1, Inches(0), Inches(0), width=k_prs.slide_width, height=k_prs.slide_height)
                tb2 = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), Inches(15), Inches(1))
                tb2.text_frame.text = "Online Media"
                for paragraph in tb2.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(50)
                        run.font.bold = True
                        run.font.name = 'Helvetica'
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.alignment = PP_ALIGN.LEFT

                # ── Prepare table copies ────────────────────────────────────────
                k_Entity_SOV3_ppt = k_Entity_SOV3.copy()
                k_sov_dt11_ppt    = k_sov_dt11.copy()
                k_pubs_ppt        = k_pubs_table.head(10).copy()
                _kn = k_pubs_ppt.select_dtypes(include=['number']).columns
                k_pubs_ppt[_kn]   = k_pubs_ppt[_kn].astype(int)
                k_Jour_ppt        = k_Jour_table20_display.head(10).copy()
                k_Jour_Comp_ppt   = k_Jour_Comp_display.copy()
                k_Jour_Client_ppt = k_Jour_Client_display.copy()

                k_dfs_ppt = [
    strip_client_prefix(k_Entity_SOV3_ppt),
    strip_client_prefix(k_sov_dt11_ppt),
    strip_client_prefix(k_pubs_ppt),
    strip_client_prefix(k_Jour_ppt),
    strip_client_prefix(k_Jour_Comp_ppt),
    strip_client_prefix(k_Jour_Client_ppt),
]
                k_table_titles_ppt = [
                    f'SOV Table of {k_client_name_clean} and competition',
                    f'Month-on-Month Table of {k_client_name_clean} and competition',
                    f'Publication Table on {k_client_name_clean} and competition',
                    f'Journalist writing on {k_client_name_clean} and competition',
                    f'Journalists writing on Comp and not on {k_client_name_clean}',
                    f'Journalists writing on {k_client_name_clean} and not on Comp',
                ]
                k_textbox_texts = [
    "• ",
    "• ",
    "• ",
    "• ",
    "• ",
    "• ",
]

                # ── Add data slides ─────────────────────────────────────────────
                for i, (df, title) in enumerate(zip(k_dfs_ppt, k_table_titles_ppt)):
                    slide = k_prs.slides.add_slide(k_prs.slide_layouts[6])
                    add_kalki_table_to_slide(slide, df, title, k_textbox_texts[i], k_prs)
                    if i == 0:
                        k_img4 = generate_bar_chart(k_Entity_SOV3_ppt.copy())
                        add_image_to_slide(slide, k_img4)
                    if i == 1:
                        # ── Add a NEW slide immediately after MOM table with line graph ──
                        k_img5 = generate_line_graph(k_sov_dt11_ppt.copy())
                        line_slide = k_prs.slides.add_slide(k_prs.slide_layouts[6])
                        if os.path.exists(img_path):
                            line_slide.shapes.add_picture(
                                img_path, Inches(0.0),
                                k_prs.slide_height - Inches(1), height=Inches(1)
                            )
                        line_title = line_slide.shapes.add_textbox(
    Inches(0.8), Inches(0.2), Inches(14), Inches(0.2)
)
                        line_title.text_frame.text = f'Month-on-Month Graph of {k_client_col} and competition'
                        for paragraph in line_title.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(28)
                                run.font.bold = True
                                run.font.name = 'Helvetica'
                                run.font.color.rgb = RGBColor(240, 127, 9)
                        line_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        add_image_to_slide1(line_slide, k_img5)

                # ── Save and use st.download_button (NOT markdown href) ─────────
                k_pptx_output = io.BytesIO()
                k_prs.save(k_pptx_output)
                k_pptx_output.seek(0)

                st.sidebar.download_button(
                    label="⬇️ Click here to Download TW PPT",
                    data=k_pptx_output,
                    file_name=k_pptx_file_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="kalki_pptx_download"
                )
            # ── Kalki Grok Prompts ────────────────────────────────────────
            st.sidebar.write("## Download TW Grok Prompts (.docx)")

            if st.sidebar.button("Download TW Grok Prompts", key="kalki_grok_button"):
                from docx import Document
                from docx.shared import RGBColor as DocxRGBColor, Pt as DocxPt

                # ── Build competitors string from Kalki entities ──────────
                k_competitors = [
                    e for e in kfinaldata["Entity"].unique()
                    if not e.startswith("Client-")
                ]
                if len(k_competitors) > 1:
                    k_competitors_str = ", ".join(k_competitors[:-1]) + f" and {k_competitors[-1]}"
                elif len(k_competitors) == 1:
                    k_competitors_str = k_competitors[0]
                else:
                    k_competitors_str = "None"

                k_ind = industry  # reuse industry from sidebar input
                k_c   = k_client_name_clean
                k_s   = start_date
                k_e   = end_date

                # ── Color palette ─────────────────────────────────────────
                CLIENT_COLOR     = DocxRGBColor(0, 102, 255)
                DATE_COLOR       = DocxRGBColor(0, 128, 0)
                COMPETITOR_COLOR = DocxRGBColor(255, 140, 0)
                INDUSTRY_COLOR   = DocxRGBColor(128, 0, 128)
                PUB_COLOR        = DocxRGBColor(255, 20, 147)
                JOURNALIST_COLOR = DocxRGBColor(225, 167, 63)
                BLACK            = DocxRGBColor(0, 0, 0)

                def k_add_run(p, txt, color=BLACK, bold=False):
                    r = p.add_run(txt)
                    r.font.color.rgb = color
                    r.bold = bold
                    return p

                k_doc = Document()

                # ── Color legend ──────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "Color code: ", bold=True)
                k_add_run(p, "Client name, ",      CLIENT_COLOR,     bold=True)
                k_add_run(p, "Dates, ",            DATE_COLOR,       bold=True)
                k_add_run(p, "Competitor name, ",  COMPETITOR_COLOR, bold=True)
                k_add_run(p, "Industry name, ",    INDUSTRY_COLOR,   bold=True)
                k_add_run(p, "Publications writing on Industry, ", PUB_COLOR,        bold=True)
                k_add_run(p, "Journalists writing on Industry",    JOURNALIST_COLOR, bold=True)

                p = k_doc.add_paragraph()
                p.add_run(
                    "I work in Media Research Team at a PR Company, I will be sharing the below "
                    "Qualitative insights with the PR professionals. Please keep this in mind and "
                    "provide insights accordingly."
                )

                # ── Requirements ─────────────────────────────────────────
                p = k_doc.add_paragraph()
                p.add_run("Satisfy the below requirements :").bold = True

                p = k_doc.add_paragraph()
                p.add_run("Prompt 1 follow the below requirements").bold = True
                for req in [
                    "Please give topicwise/bucketwise paragraph with topic/bucket highlighted please be very much elaborative as possible",
                    "Max 2 sentences should be in a line then move into next line and follow the same thing consider them as a point and the points should always be elaborative and not in one liner",
                    "In each topicwise/bucketwise paragraph and in each and every point in the paragraph, the content should be very much elaborative as possible and there should be atleast 5-6 such points without losing relevant news in each topic and those should be elaborative",
                ]:
                    k_doc.add_paragraph(req, style='List Bullet')

                p = k_doc.add_paragraph()
                r = p.add_run("NOTE : "); r.bold = True
                p.add_run("Don't provide insight for prompt 3")

                p = k_doc.add_paragraph()
                p.add_run("Prompt 4,5 follow the below requirements").bold = True
                k_doc.add_paragraph(
                    "Please give the insights in tabular format, don't mention no. of articles anywhere) "
                    "Give one elaborative paragraph before creating a table.",
                    style='List Bullet'
                )

                p = k_doc.add_paragraph()
                p.add_run("Prompt 6,7 follow the below requirements").bold = True
                k_doc.add_paragraph(
                    "Give five or more critiques pointwise (be elaborative as much as possible) with the "
                    "critiques pointer headers highlighted (just have 2-3 elaborative pointers in each critiques). "
                    "Give one paragraph elaborative description before giving the pointers for critiques.",
                    style='List Bullet'
                )

                p = k_doc.add_paragraph()
                p.add_run("Prompt 9,10 follow the below requirements").bold = True
                p = k_doc.add_paragraph(style='List Bullet')
                p.add_run(
                    "Give paragraph wise for each Publications/Journalist(along with publication name) mentioned here "
                    "and highlight the Publication name/Journalist, please be very much elaborative don't be generic "
                    "relate it with the news released by these publications/journalist on the "
                )
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry. (")
                r = p.add_run("Note : "); r.bold = True
                p.add_run("Publications/Journalist(along with publication name) should be highlighted)")

                p = k_doc.add_paragraph(style='List Bullet')
                p.add_run("Don't give me the number of news written by these publications/journalists just the content that has been written on the ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry.")

                k_doc.add_paragraph(
                    "Additionally, follow the same requirements given for Prompt 1 excluding the first point in it "
                    "(i.e Please give topicwise/bucketwise paragraph with topic/bucket)",
                    style='List Bullet'
                )

                k_doc.add_paragraph()

                p = k_doc.add_paragraph()
                r = p.add_run("For Prompts 1-10 "); r.bold = True
                p.add_run("please note that : do not consider press release from the companys website or any social media platform")

                p = k_doc.add_paragraph()
                p.add_run("Before giving insight for each prompt mention the title like Conversations on Client company, Topicwise Conversations on Client company,… etc (and replace Client Company with ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(" and Industry with ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(") with formatting ### and no bold formatting and topics/buckets with bold formatting **, please follow the formatting strictly don't use bold formatting in the content of any of the buckets/topics")

                k_doc.add_paragraph()

                # ── Prompt 1 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "1) Conversations on Client company - ", bold=True)
                p.add_run("Could you Summarize the news articles from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(" for ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(
                    "? Please summarize as many topics as possible but do not consider press release from the "
                    "companys website or any social media platform. Only summarize the articles from print and "
                    "online news platforms."
                )

                # ── Prompt 2 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "2) Topicwise Conversation on Client Company - ", bold=True)
                p.add_run("Could you Summarize the news articles from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(" for ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(
                    ". Please summarize the news articles as per the following categories. I am giving you the buckets. "
                    "Please arrange the news as per their content in the relevant buckets and summarize that news. "
                    "Only summarize the news articles from print and online news platforms. The buckets are as follows: "
                    "Financial Performance, Product and Services, Social Good (includes CSR, ESG, Philanthropy, Environment), "
                    "Employee Engagement (includes hiring, resignation, layoffs, training, skilling, employee benefits, appraisals...), "
                    "Business Strategy (include growth, mergers, future, market share...), Vision and Leadership (Interviews, "
                    "interaction, thought leadership, authored articles...), Legal and Regulatory, Tech & innovation, "
                    "Stock related (stock recommendations, stock movements). ("
                )
                r = p.add_run("Note : "); r.bold = True
                p.add_run(
                    "Please give topicwise/bucketwise paragraph with topic/bucket highlighted please be very much "
                    "elaborative as possible. Max 2 sentences should be in a line then move into next line and follow "
                    "the same thing consider them as a point and the points should always be elaborative and not in one liner. "
                    "In each topicwise/bucketwise paragraph and in each and every point in the paragraph, the content should "
                    "be very much elaborative as possible and there should be atleast 5-6 such points without losing relevant "
                    "news in each topic and those should be elaborative)"
                )

                # ── Prompt 3 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "3) Topicwise Conversation on Competitor Company - ", bold=True)
                p.add_run("Could you Summarize the news articles from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(" for ")
                k_add_run(p, k_competitors_str, COMPETITOR_COLOR)
                p.add_run(
                    "? Give Topicwise conversation entity wise, one after the other follow the bucket structure for each "
                    "of the entities, want separate separate topic wise conversation for each entity. Please summarize the "
                    "news articles as per the following categories. I am giving you the buckets. Please arrange the news as "
                    "per their content in the relevant buckets and summarize that news. Only summarize the news articles from "
                    "print and online news platforms. The buckets are as follows: Financial Performance, Product and Services, "
                    "Social Good (includes CSR, ESG, Philanthropy, Environment), Employee Engagement (includes hiring, "
                    "resignation, layoffs, training, skilling, employee benefits, appraisals...), Business Strategy (include "
                    "growth, mergers, future, market share...), Vision and Leadership (Interviews, interaction, thought "
                    "leadership, authored articles...), Legal and Regulatory, Tech & innovation, Stock related (stock "
                    "recommendations, stock movements). ("
                )
                r = p.add_run("Note : "); r.bold = True
                p.add_run("Please give topicwise/bucketwise paragraph with topic/bucket highlighted please be very much elaborative as possible. Max 2 sentences should be in a line then move into next line and follow the same thing consider them as a point and the points should always be elaborative and not in one liner. In each topicwise/bucketwise paragraph and in each and every point in the paragraph, the content should be very much elaborative as possible and there should be atleast 5-6 such points without losing relevant news in each topic and those should be elaborative) Give it with the Header i.e Topicwise Conversation on Competitor Company with formatting ### and competitor name ")
                k_add_run(p, k_competitors_str, COMPETITOR_COLOR)
                p.add_run(" with formatting ## with no bold formatting and buckets with bold formatting **, please follow the formatting strictly don't use bold formatting in the content of any of the buckets/topics")

                # ── Prompt 4 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "4) Month – on – Month Insights - ", bold=True)
                p.add_run("Give me month on month breakdown of news coverage for ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(" from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(
                    ". Give me details of events that have lead to a spike in the media coverage and if no news is "
                    "present for a particular month mention that in the table itself."
                )

                # ── Prompt 5 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "5) Unique Conversations by Competitors - ", bold=True)
                p.add_run("What factors contributed to ")
                k_add_run(p, k_competitors_str, COMPETITOR_COLOR)
                p.add_run(" higher media coverage in the ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry between ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(" compared to ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(
                    "? Identify the unique conversation topics (topics where "
                )
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(
                    " is not mentioned) where these companies were mentioned in the Headline or the lead para or if "
                    "they were mentioned atleast twice in the article that have driven the higher media coverage. ("
                )
                r = p.add_run("Note : "); r.bold = True
                p.add_run(
                    "give the insights in tabular format with Column name : Company, Unique Conversation Topics be "
                    "elaborative relating it with the news, (for each row there should be just 1 company name all its "
                    "unique conversation should be there beside it in the Unique Conversation Topics column) Give one "
                    "elaborative paragraph before creating a table."
                )

                # ── Prompt 6 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "6) Reputational Risks for Client - ", bold=True)
                p.add_run("What is the online news media saying about ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(" from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run("? Give five or more critiques.")

                # ── Prompt 7 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "7) Reputational Risks for Industry - ", bold=True)
                p.add_run("What is the online news media saying about ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run("? Give five or more critiques.")

                # ── Prompt 8 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "8) Industry Snapshot - ", bold=True)
                p.add_run("What is the online news media conversation in the ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(
                    " and identify the companies who are a part of these conversations. ("
                )
                r = p.add_run("Note : "); r.bold = True
                p.add_run(
                    "give the insights in tabular format and be elaborative relating it with the news, give company "
                    "wise insight for this, don't put more than 1 company in one row) Give one elaborative paragraph "
                    "before creating a table. You can also consider companies who are not the competitors but are "
                    "present online news media conversation in the "
                )
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry.")

                # ── Prompt 9 ──────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "9) Publications writing on Industry – ", bold=True)
                p.add_run("Which Indian publications (min 3) have frequently written on ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(" and what are the conversations about the ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry AND which ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" companies have been mentioned by them?")

                # ── Prompt 10 ─────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "10) Journalist writing on Industry – ", bold=True)
                p.add_run("Which Indian journalists name (min 3) have (mention their Publication name too) frequently written on ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry from ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run(" and what are the conversations about the ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" industry AND which ")
                k_add_run(p, k_ind, INDUSTRY_COLOR)
                p.add_run(" companies have been mentioned by them?")

                # ── Prompt 11 ─────────────────────────────────────────────
                p = k_doc.add_paragraph()
                k_add_run(p, "11) X Insights - ", bold=True)
                p.add_run("What is being said on Twitter X about ")
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(" between ")
                k_add_run(p, k_s, DATE_COLOR)
                p.add_run(" to ")
                k_add_run(p, k_e, DATE_COLOR)
                p.add_run("? (")
                r = p.add_run("Note : "); r.bold = True
                p.add_run(
                    "Please be very much elaborative, should be majorly on users conversation around "
                )
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(
                    ", if possible highlight the user whose content about "
                )
                k_add_run(p, k_c, CLIENT_COLOR)
                p.add_run(
                    " has larger influence among the twitter audience and give breakdown of Positive Discussions, "
                    "Criticisms and Complaints, Neutral/Informational Mentions, and highlight it, please be as "
                    "elaborative as possible)"
                )

                # ── Final NOTE ────────────────────────────────────────────
                p = k_doc.add_paragraph()
                r = p.add_run("NOTE : "); r.bold = True
                p.add_run("Don't provide insight for prompt 3")

                # ── Save & download ───────────────────────────────────────
                k_grok_buffer = io.BytesIO()
                k_doc.save(k_grok_buffer)
                k_grok_buffer.seek(0)
                k_grok_b64 = base64.b64encode(k_grok_buffer.read()).decode()
                k_grok_href = (
                    f'<a href="data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;'
                    f'base64,{k_grok_b64}" download="TW_Grok_Prompts_{k_c}.docx">'
                    f'Download TW Grok Prompts (.docx)</a>'
                )
                st.sidebar.markdown(k_grok_href, unsafe_allow_html=True)
                # ── Build presentation ────────────────────────────────
if date_selected and industry_provided :# File Upload Section
    st.sidebar.write("## Upload an Online or Print file for tables")
    file = st.sidebar.file_uploader("Upload Data File (Excel or CSV)", type=["xlsx", "csv"])
    if file:
        st.sidebar.write("File Uploaded Successfully!")
        data = load_data(file)
        if data is not None:
            # Data Preview Section (optional)
            # st.write("## Data Preview")
            # st.write(data)
            # Data preprocessing
            data.drop(columns =data.columns[20:], inplace=True)
            data = data.rename({'Influencer': 'Journalist'}, axis=1)
            all_duplicates = pd.DataFrame()
        # data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        # Check if specific columns exist before dropping duplicates
            if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
                data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
            if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
                data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
            if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
                data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
                
            finaldata = data
            raw_data_for_pubtype = data.copy()  # save before dedup
            finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()
            competitors = [ent for ent in finaldata['Entity'].unique() if not ent.startswith("Client-")]
            if len(competitors) > 1:
                competitors_str = ", ".join(competitors[:-1]) + f" and {competitors[-1]}"
            elif len(competitors) == 1:
                competitors_str = competitors[0]
            else:
                competitors_str = "None"

            # Date formatting (Month Year)


            # Industry (from user input)
            industry = industry_input.strip()
           # Share of Voice (SOV) Calculation
            En_sov = pd.crosstab(finaldata['Entity'], columns='News Count', values=finaldata['Entity'], aggfunc='count').round(0)
            En_sov.sort_values('News Count', ascending=False)
            En_sov['% '] = ((En_sov['News Count'] / En_sov['News Count'].sum()) * 100).round(0)
            Sov_table = En_sov.sort_values(by='News Count', ascending=False)
            Sov_table.loc['Total'] = Sov_table.sum(numeric_only=True, axis=0)
            Entity_SOV1 = Sov_table
            Entity_SOV3 = pd.DataFrame(Entity_SOV1.to_records()).round()
            Entity_SOV3['News Count'] = Entity_SOV3['News Count'].astype(int)
            Entity_SOV3['% '] = Entity_SOV3['% '].astype(int)
            Entity_SOV3['% '] = Entity_SOV3['% '].astype(str) + '%'
        # Entity_SOV3 = pd.DataFrame(Entity_SOV3.to_records())

        # # Plot the bar graph
        # plt.figure(figsize=(10, 6))
        # bars = plt.bar(
        #   Entity_SOV3['Entity'], 
        #   Entity_SOV3['News Count'], 
        #   color='skyblue', 
        #   edgecolor='black'
        #    )
        # # Add labels on top of each bar
        # for bar in bars:
        #     height = bar.get_height()
        #     plt.text(bar.get_x() + bar.get_width() / 2, height, f'{height}', ha='center', va='bottom', fontsize=10)

        # # Customize the graph
        # plt.title("Share of Voice (SOV)", fontsize=14)
        # plt.xlabel("Entity", fontsize=12)
        # plt.ylabel("News Count", fontsize=12)
        # plt.grid(axis='y', linestyle='--', alpha=0.7)
        # plt.tight_layout()


        
        #News Count Total 
            total_news_count = int(Entity_SOV3.loc[Entity_SOV3["Entity"] == "Total", "News Count"].values[0])
# Additional MOM DataFrames
            finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()
            sov_dt = pd.crosstab((finaldata['Date'].dt.to_period('M')), finaldata['Entity'], margins=True, margins_name='Total')
            sov_dt1 = pd.DataFrame(sov_dt.to_records())
            client_columndt = [ent for ent in finaldata['Entity'].unique() if ent.startswith("Client-")][0]
            sov_order = Entity_SOV3['Entity'].tolist()
            sov_order_no_client = [ent for ent in sov_order if not ent.startswith("Client-") and ent != "Total"]
            ordered_cols = (['Date', client_columndt] +[ent for ent in sov_order_no_client if ent in sov_dt1.columns] + (['Total'] if 'Total' in sov_dt1.columns else []))
            for_entity_data = finaldata.copy()
           # Reorder the columns
            sov_dt11 = sov_dt1[ordered_cols]
            selected_columndt = sov_dt1[["Date", client_columndt]]
            selected_columndt = selected_columndt.iloc[:-1]
            selected_columndt = selected_columndt.sort_values(by=client_columndt, ascending=False)
            # Extract the top 3 publications and their counts
            topdt_1 = selected_columndt.iloc[0:1]  # First publication
            # topc_2 = selected_columndt.iloc[1:2]  # Second publication
            # topc_3 = selected_columndt.iloc[2:3]  # Third publication
            # Save them in separate DataFrames
            df_topdt1 = topdt_1.reset_index(drop=True)
            # df_topc2 = topc_2.reset_index(drop=True)
# df_topc3 = topc_3.reset_index(drop=True)
# Extract publication name and count for the top 3
            topdt_1_name = df_topdt1.iloc[0]["Date"]
            topdt_1_count = df_topdt1.iloc[0][client_columndt]

        # topc_2_name = df_topc2.iloc[0]["Publication Name"]
        # topc_2_count = df_topc2.iloc[0][client_column]

        
        #Publication Name
            finaldata_non_exploded = finaldata.copy()
            finaldata['Journalist'] = (finaldata['Journalist'].astype(str).str.split(',').apply(lambda x: [j.strip() for j in x]))
            finaldata = finaldata.explode('Journalist').reset_index(drop=True)
    
            jr_tab = pd.crosstab(finaldata['Journalist'], finaldata['Entity'])
            jr_tab = jr_tab.reset_index(level=0)
            newdata = finaldata[['Journalist', 'Publication Name']]
            Journalist_Table = pd.merge(jr_tab, newdata, how='inner', left_on=['Journalist'], right_on=['Journalist'])
            Journalist_Table.drop_duplicates(subset=['Journalist'], keep='first', inplace=True)
            valid_columns = Journalist_Table.select_dtypes(include='number').columns
            Journalist_Table['Total'] = Journalist_Table[valid_columns].sum(axis=1)
            Jour_table = Journalist_Table.sort_values('Total', ascending=False).round()
            bn_row = Jour_table.loc[Jour_table['Journalist'] == 'Bureau News']
            Jour_table = Jour_table[Jour_table['Journalist'] != 'Bureau News']
            Jour_table = pd.concat([Jour_table, bn_row], ignore_index=True)
    #         Jour_table = Journalist_Table.reset_index(drop=True)
            Jour_table.loc['GrandTotal'] = Jour_table.sum(numeric_only=True, axis=0)
            columns_to_convert = Jour_table.columns.difference(['Journalist', 'Publication Name'])
            Jour_table[columns_to_convert] = Jour_table[columns_to_convert].astype(int)
            Jour_table.insert(1, 'Publication Name', Jour_table.pop('Publication Name'))
            Jour_table1 = Jour_table.head(10)
    
            #UNIQUE_ARTICLES_WRITTEN_BY_JOURNALIST
            finaldatau = finaldata.copy()
            finaldatau.drop_duplicates(subset=['Date','Headline','Publication Name','Journalist'], keep='first', inplace=True, ignore_index=True)
        
            jr_tabu = finaldatau['Journalist'].value_counts().reset_index()
            jr_tabu.columns = ['Journalist', 'Total']
            newdatau = finaldatau[['Journalist', 'Publication Name']]
            Journalist_Tableu = pd.merge(jr_tabu, newdatau, how='inner', left_on=['Journalist'], right_on=['Journalist'])
            Journalist_Tableu.drop_duplicates(subset=['Journalist'], keep='first', inplace=True)
            Jour_tableu = Journalist_Tableu.sort_values('Total', ascending=False).round()
            bn_rowu = Jour_tableu.loc[Jour_tableu['Journalist'] == 'Bureau News']
            Jour_tableu = Jour_tableu[Jour_tableu['Journalist'] != 'Bureau News']
            Jour_tableu = pd.concat([Jour_tableu, bn_rowu], ignore_index=True)
            Jour_tableu.loc['Total'] = Jour_tableu.sum(numeric_only=True, axis=0)
            Jour_tableu.loc['Total', 'Journalist'] = 'Total'   # ← add this line
            columns_to_convert = Jour_tableu.columns.difference(['Journalist', 'Publication Name'])
            Jour_tableu[columns_to_convert] = Jour_tableu[columns_to_convert].astype(int)
            Jour_tableu.insert(1, 'Publication Name', Jour_tableu.pop('Publication Name'))
            Jour_tableu.loc[Jour_tableu['Journalist'] == 'Bureau News', 'Publication Name'] = 'Across Publications'
            numeric_cols_u = Jour_tableu.select_dtypes(include='number').columns
            cols_to_drop_u = [col for col in numeric_cols_u if col != 'Total']
            Jour_tableu = Jour_tableu.drop(columns=cols_to_drop_u)
            numeric_cols_t = Jour_table.select_dtypes(include='number').columns
            numeric_cols_t = [col for col in numeric_cols_t if col != 'Total']
            cols_to_use_t = ['Journalist'] + numeric_cols_t
            Jour_table_subset = Jour_table[cols_to_use_t]
            Unique_Articles = Jour_tableu.merge(Jour_table_subset, on='Journalist', how='left')
            cols = list(Unique_Articles.columns)
            cols.remove('Total')
            cols.append('Total')
            Unique_Articles = Unique_Articles[cols]
            bn_index = Unique_Articles[Unique_Articles['Journalist'] == 'Bureau News' ].index
            Unique_Articles.loc[bn_index+1, 'Journalist'] = 'Total'
            ordered_cols = ['Journalist', 'Publication Name', client_columndt] + [ent for ent in sov_order_no_client if ent in Unique_Articles.columns] + (['Total'] if 'Total' in Unique_Articles.columns else [])
            Unique_Articles = Unique_Articles[ordered_cols]
            Unique_Articles['Client %'] = (
    (Unique_Articles[client_columndt] / Unique_Articles['Total']) * 100
).replace([float('inf'), float('-inf')], 0).fillna(0).round().astype(int)

            # ── FIX: force the Total row to match SOV exactly ──────────────────
            sov_lookup = dict(zip(Entity_SOV3['Entity'], Entity_SOV3['News Count']))
            total_idx = Unique_Articles[Unique_Articles['Journalist'] == 'Total'].index
            
            for _col in Unique_Articles.columns:
                if _col in sov_lookup and _col not in ['Total', 'Client %', 'Journalist', 'Publication Name']:
                    Unique_Articles.loc[total_idx, _col] = int(sov_lookup[_col])
            
            # Total Unique Articles = sum of unique articles across all journalists (excluding the Total row itself)
            grand_unique_total = int(
                Unique_Articles.loc[Unique_Articles['Journalist'] != 'Total', 'Total'].sum()
            )
            Unique_Articles.loc[total_idx, 'Total'] = grand_unique_total
            
            # Recalculate Client % for the Total row using the corrected numbers
            grand_client_val = Unique_Articles.loc[total_idx, client_columndt].values[0]
            if grand_unique_total > 0:
                Unique_Articles.loc[total_idx, 'Client %'] = int(round((grand_client_val / grand_unique_total) * 100))
            else:
                Unique_Articles.loc[total_idx, 'Client %'] = 0
            pub_table1 = pd.crosstab(finaldata_non_exploded['Publication Name'], finaldata_non_exploded['Entity'])
            pub_table1 = pub_table1.reset_index(level=0)
            finaldatauq = finaldata_non_exploded.copy()
            finaldatauq.drop_duplicates(subset=['Date','Headline','Publication Name','Journalist'], keep='first', inplace=True, ignore_index=True)
            pub_table = finaldatauq['Publication Name'].value_counts().reset_index()
            pub_table.columns = ['Publication Name', 'Total']
            # Get all numeric columns in Jour_tableu
            numeric_cols_u = pub_table.select_dtypes(include='number').columns
            # Determine which numeric columns to drop (all except 'Total')
            cols_to_drop_u = [col for col in numeric_cols_u if col != 'Total']
            pub_table = pub_table.drop(columns=cols_to_drop_u)
            numeric_cols_t = pub_table1.select_dtypes(include='number').columns
            numeric_cols_t = [col for col in numeric_cols_t if col != 'Total']
            cols_to_use_t = ['Publication Name'] + numeric_cols_t
            pub_table_subset = pub_table1[cols_to_use_t]
            Unique_pub = pub_table.merge(pub_table_subset, on='Publication Name', how='left')
            cols = list(Unique_pub.columns)
            cols.remove('Total')
            cols.append('Total')
            pubs_table = Unique_pub[cols]

            ordered_cols = ['Publication Name', client_columndt] + [ent for ent in sov_order_no_client if ent in pubs_table.columns] + (['Total'] if 'Total' in pubs_table.columns else [])
            pubs_table = pubs_table[ordered_cols]
            pubs_table = pubs_table.sort_values('Total', ascending=False).round()

            pubs_table.loc['Total'] = pubs_table.sum(numeric_only=True, axis=0)
            pubs_table.loc['Total', 'Publication Name'] = 'Total'
            
            # ── Override entity columns in Total row to match SOV exactly ──
            sov_lookup_pub = dict(zip(Entity_SOV3['Entity'], Entity_SOV3['News Count']))
            for _col in pubs_table.columns:
                if _col in sov_lookup_pub:
                    pubs_table.loc['Total', _col] = int(sov_lookup_pub[_col])
            
            # Override Total in Total row to match SOV grand total
            pubs_table.loc['Total', 'Total'] = int(Entity_SOV3.loc[Entity_SOV3['Entity'] == 'Total', 'News Count'].values[0])
            
            # ── Recalculate Client % for ALL rows (Total row now uses corrected SOV values) ──
            pubs_table['Client %'] = (
                (pubs_table[client_columndt] / pubs_table['Total']) * 100
            ).round().fillna(0).astype(int)
            
            numeric_columns = pubs_table.select_dtypes(include=['number']).columns
            pubs_table[numeric_columns] = pubs_table[numeric_columns].astype(int)
            pubs_table1 = pubs_table.head(10)
            pubs_table2O=  pubs_table.head(20)
            pubs_table2O =pubs_table2O.rename(columns= {'Total': 'Total Unique Articles'})



            # Extract the top 3 publications and their counts
            top_1 = pubs_table1.iloc[0:1]  # First publication
            top_2 = pubs_table1.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(pubs_table1) > 2:
                top_3 = pubs_table1.iloc[2:3]  # Third publication
                df_top3 = top_3.reset_index(drop=True)
                top_3_name = df_top3.iloc[0]["Publication Name"]
                top_3_count = df_top3.iloc[0]["Total"]
            else:
                top_3_name = ""
                top_3_count = 0  # You can assign any default value for count
    
            # Save the first two publications in separate DataFrames
            df_top1 = top_1.reset_index(drop=True)
            df_top2 = top_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            top_1_name = df_top1.iloc[0]["Publication Name"]
            top_1_count = df_top1.iloc[0]["Total"]
    
            top_2_name = df_top2.iloc[0]["Publication Name"]
            top_2_count = df_top2.iloc[0]["Total"]
            
    
            # # Extract the top 3 publications and their counts
            # top_1 = pubs_table1.iloc[0:1]  # First publication
            # top_2 = pubs_table1.iloc[1:2]  # Second publication
            # top_3 = pubs_table1.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_top1 = top_1.reset_index(drop=True)
            # df_top2 = top_2.reset_index(drop=True)
            # df_top3 = top_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # top_1_name = df_top1.iloc[0]["Publication Name"]
            # top_1_count = df_top1.iloc[0]["Total"]
    
            # top_2_name = df_top2.iloc[0]["Publication Name"]
            # top_2_count = df_top2.iloc[0]["Total"]
    
            # top_3_name = df_top3.iloc[0]["Publication Name"]
            # top_3_count = df_top3.iloc[0]["Total"]
        
            # Dynamically identify the client column
            client_column = [col for col in pubs_table1.columns if col.startswith("Client-")][0]
    
            # Select the "Publication Name" column and the dynamically identified client column
            selected_columns = pubs_table1[["Publication Name", client_column]]
            
            selected_columns = selected_columns.iloc[:-1]
            selected_columns = selected_columns.sort_values(by=client_column, ascending=False)
    
            # Extract the top 3 publications and their counts
            topc_1 = selected_columns.iloc[0:1]  # First publication
            topc_2 = selected_columns.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(selected_columns) > 2:
                topc_3 = selected_columns.iloc[2:3]  # Third publication
                df_topc3 = topc_3.reset_index(drop=True)
                topc_3_name = df_topc3.iloc[0]["Publication Name"]
                topc_3_count = df_topc3.iloc[0][client_column]
            else:
                topc_3_name = ""
                topc_3_count = 0  # You can assign any default value for count

            # topc_3 = selected_columns.iloc[2:3]  # Third publication
            pubs_table_trimmed = pubs_table.iloc[:-1]
            top10_pub_sum = pubs_table_trimmed[client_column].sort_values(ascending=False).head(10).sum()
            client_sov_count = int(Entity_SOV3.loc[Entity_SOV3["Entity"] == client_column, "News Count"].values[0])
            top10_pub_perc = int(round(( top10_pub_sum / client_sov_count) * 100))

            client_sov = Unique_Articles.loc[Unique_Articles['Journalist'] == 'Total',client_column].values[0]
            bn_match = Unique_Articles.loc[Unique_Articles['Journalist'] == 'Bureau News', client_column]
            bureau_articles = bn_match.values[0] if not bn_match.empty else 0
            individual_articles = client_sov - bureau_articles
            client_sov = pd.to_numeric(client_sov, errors='coerce')
            client_sov = client_sov if pd.notna(client_sov) and client_sov != 0 else 0
            
            if client_sov > 0:
                bureau_percentage = int(round((bureau_articles / client_sov) * 100, 0))
                individual_percentage = int(round((individual_articles / client_sov) * 100, 0))
            else:
                bureau_percentage = 0
                individual_percentage = 0
            filtered_df = Unique_Articles[~Unique_Articles['Journalist'].isin(['Total', 'Bureau News'])]
            total_journalists = len(filtered_df)
            total_articles = filtered_df[filtered_df['Total'].notna() & (filtered_df['Total'] > 0)]['Total'].sum()
            non_zero_journalists = filtered_df[filtered_df[client_column] > 0].shape[0]
            articles_for_client = filtered_df[filtered_df[client_column] > 0][client_column].sum()
            #client_journalist_percentage =  int(round((non_zero_journalists/ total_journalists) * 100,0))
            #engage_with = total_journalists-non_zero_journalists
        
    
            # Save them in separate DataFrames
            df_topc1 = topc_1.reset_index(drop=True)
            df_topc2 = topc_2.reset_index(drop=True)
            # df_topc3 = topc_3.reset_index(drop=True)
    
            # Extract publication name and count for the top 3
            topc_1_name = df_topc1.iloc[0]["Publication Name"]
            topc_1_count = df_topc1.iloc[0][client_column]
    
            topc_2_name = df_topc2.iloc[0]["Publication Name"]
            topc_2_count = df_topc2.iloc[0][client_column]
    
            # topc_3_name = df_topc3.iloc[0]["Publication Name"]
            # topc_3_count = df_topc3.iloc[0][client_column]


            PP = pd.crosstab(finaldata_non_exploded['Publication Name'], finaldata_non_exploded['Publication Type'])
            PP['Total'] = PP.sum(axis=1)
            PP_table = PP.sort_values('Total', ascending=False).round()
            PP_table.loc['GrandTotal'] = PP_table.sum(numeric_only=True, axis=0)
            
            #Publication Name & Entity Table
            PT_Entity = pd.crosstab(raw_data_for_pubtype['Publication Type'], raw_data_for_pubtype['Entity'])
            PT_Entity['Total'] = PT_Entity.sum(axis=1)
            PType_Entity = PT_Entity.sort_values('Total', ascending=False).round()
            PType_Entity.loc['Total'] = PType_Entity.sum(numeric_only=True, axis=0)
            PType_Entity = pd.DataFrame(PType_Entity.to_records())
            ordered_cols = ['Publication Type', client_columndt] + [ent for ent in sov_order_no_client if ent in PType_Entity.columns] + (['Total'] if 'Total' in PType_Entity.columns else [])
            PType_Entity = PType_Entity[ordered_cols]
            excluded_df = PType_Entity.sort_values(by=client_columndt, ascending=False)
            p = inflect.engine()
            publication_types = excluded_df['Publication Type'].unique()[3:].tolist()
            publication_types_str = p.join(publication_types)

           
            # Extract the top 3 publications and their counts
            topt_1 = PType_Entity.iloc[0:1]  # First publication
            topt_2 = PType_Entity.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(PType_Entity) > 2:
                topt_3 = PType_Entity.iloc[2:3]  # Third publication
                df_topt3 = topt_3.reset_index(drop=True)
                topt_3_name = df_topt3.iloc[0]["Publication Type"]
                topt_3_count = df_topt3.iloc[0]["Total"]
            else:
                topt_3_name = ""
                topt_3_count = 0  # You can assign any default value for count
    
            # Save the first two publications in separate DataFrames
            df_topt1 = topt_1.reset_index(drop=True)
            df_topt2 = topt_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            topt_1_name = df_topt1.iloc[0]["Publication Type"]
            topt_1_count = df_topt1.iloc[0]["Total"]

            topt_2_name = df_topt2.iloc[0]["Publication Type"]
            topt_2_count = df_topt2.iloc[0]["Total"]


        # # Extract the top 3 publications and their counts
        # topt_1 = PType_Entity.iloc[0:1]  # First publication
        # topt_2 = PType_Entity.iloc[1:2]  # Second publication
        # topt_3 = PType_Entity.iloc[2:3]  # Third publication

        # # Save them in separate DataFrames
        # df_topt1 = topt_1.reset_index(drop=True)
        # df_topt2 = topt_2.reset_index(drop=True)
        # df_topt3 = topt_3.reset_index(drop=True)

        # # Extract publication name and count for the top 3
        # topt_1_name = df_topt1.iloc[0]["Publication Type"]
        # topt_1_count = df_topt1.iloc[0]["Total"]

        # topt_2_name = df_topt2.iloc[0]["Publication Type"]
        # topt_2_count = df_topt2.iloc[0]["Total"]

        # topt_3_name = df_topt3.iloc[0]["Publication Type"]
        # topt_3_count = df_topt3.iloc[0]["Total"]

        # Dynamically identify the client column
            client_columnp = [col for col in PType_Entity.columns if col.startswith("Client-")][0]
    
            # Select the "Publication Name" column and the dynamically identified client column
            selected_columnp = PType_Entity[["Publication Type", client_columnp]]
            
            selected_columnp = selected_columnp.iloc[:-1]
            selected_columnp = selected_columnp.sort_values(by=client_columnp, ascending=False)
    
            # Extract the top 3 publications and their counts
            topp_1 = selected_columnp.iloc[0:1]  # First publication
            topp_2 = selected_columnp.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(selected_columnp) > 2:
                topp_3 = selected_columnp.iloc[2:3]  # Third publication
                df_topp3 = topp_3.reset_index(drop=True)
                topp_3_name = df_topp3.iloc[0]["Publication Type"]
                topp_3_count = df_topp3.iloc[0][client_column]
            else:
                topp_3_name = ""
                topp_3_count = 0  # You can assign any default value for count
    
            
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # Save them in separate DataFrames
            df_topp1 = topp_1.reset_index(drop=True)
            df_topp2 = topp_2.reset_index(drop=True)
            # df_topc3 = topc_3.reset_index(drop=True)
    
            # Extract publication name and count for the top 3
            topp_1_name = df_topp1.iloc[0]["Publication Type"]
            topp_1_count = df_topp1.iloc[0][client_column]
    
            #topp_2_name = df_topp2.iloc[0]["Publication Type"]
            #topp_2_count = df_topp2.iloc[0][client_column]
    
    
            # # Extract the top 3 publications and their counts
            # topp_1 = selected_columnp.iloc[0:1]  # First publication
            # topp_2 = selected_columnp.iloc[1:2]  # Second publication
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topp1 = topp_1.reset_index(drop=True)
            # df_topp2 = topp_2.reset_index(drop=True)
            # df_topp3 = topp_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topp_1_name = df_topp1.iloc[0]["Publication Type"]
            # topp_1_count = df_topp1.iloc[0][client_column]
    
            # topp_2_name = df_topp2.iloc[0]["Publication Type"]
            # topp_2_count = df_topp2.iloc[0][client_column]
    
            # topp_3_name = df_topp3.iloc[0]["Publication Type"]
            # topp_3_count = df_topp3.iloc[0][client_column]
    
            # Journalist Table
            Unique_Articles1O = Unique_Articles.head(10)
            Unique_Articles2O = Unique_Articles.head(20)
            Unique_Articles2O = Unique_Articles2O.rename(columns={'Total': 'Total Unique Articles'})
    
            # Extract the top 3 publications and their counts
            topj_1 =  Unique_Articles1O.iloc[0:1]  # First publication
            topj_2 =  Unique_Articles1O.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(Unique_Articles1O) > 2:
                topj_3 = Unique_Articles1O.iloc[2:3]  # Third publication
                df_topj3 = topj_3.reset_index(drop=True)
                topj_3_name = df_topj3.iloc[0]["Journalist"]
                topj_3_count = df_topj3.iloc[0]["Total"]
            else:
                topj_3_name = ""
                topj_3_count = 0  # You can assign any default value for count
            Unique_Articles1O = Unique_Articles1O.rename(columns={'Total': 'Total Unique Articles'})
            Unique_Articles = Unique_Articles.rename(columns={'Total': 'Total Unique Articles'})
            client_col_jour = [col for col in Unique_Articles.columns if col.startswith("Client-")][0]
            low_mention_jour = Unique_Articles[
                (Unique_Articles[client_col_jour].isin([0, 1])) &
                (~Unique_Articles['Journalist'].astype(str).str.strip().str.lower().eq('total'))
            ]
            filtered_jour = low_mention_jour[low_mention_jour['Journalist'] != 'Bureau News']
            top3_jour = filtered_jour.sort_values('Total Unique Articles', ascending=False).head(3)

            journalists_list = [f"{row['Journalist']} ({row['Publication Name']})" for _, row in top3_jour.iterrows()]

            if len(journalists_list) > 1:
                journalists_str = ", ".join(journalists_list[:-1]) + f" and {journalists_list[-1]}"
            elif len(journalists_list) == 1:
                journalists_str = journalists_list[0]
            else:
                journalists_str = "None"
            # Save the first two publications in separate DataFrames
            df_topj1 = topj_1.reset_index(drop=True)
            df_topj2 = topj_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            topj_1_name = df_topj1.iloc[0]["Journalist"]
            topj_1_count = df_topj1.iloc[0]["Total"]
    
            topj_2_name = df_topj2.iloc[0]["Journalist"]
            topj_2_count = df_topj2.iloc[0]["Total"]
    
    
            # # Extract the top 3 publications and their counts
            # topj_1 = Jour_table1.iloc[0:1]  # First publication
            # topj_2 = Jour_table1.iloc[1:2]  # Second publication
            # topj_3 = Jour_table1.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topj1 = topj_1.reset_index(drop=True)
            # df_topj2 = topj_2.reset_index(drop=True)
            # df_topj3 = topj_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topj_1_name = df_topj1.iloc[0]["Journalist"]
            # topj_1_count = df_topj1.iloc[0]["Total"]
    
            # topj_2_name = df_topj2.iloc[0]["Journalist"]
            # topj_2_count = df_topj2.iloc[0]["Total"]
    
            # topj_3_name = df_topj3.iloc[0]["Journalist"]
            # topj_3_count = df_topj3.iloc[0]["Total"]
    
            # Extract the top 3 publications and their counts
            topjt_1 = Jour_table1.iloc[0:1]  # First publication
            topjt_2 = Jour_table1.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(Jour_table1) > 2:
                topjt_3 = Jour_table1.iloc[2:3]  # Third publication
                df_topjt3 = topjt_3.reset_index(drop=True)
                topjt_3_name = df_topjt3.iloc[0]["Publication Name"]
                topjt_3_count = df_topjt3.iloc[0]["Total"]
            else:
                topjt_3_name = ""
                topjt_3_count = 0  # You can assign any default value for count
    
            # Save the first two publications in separate DataFrames
            df_topjt1 = topjt_1.reset_index(drop=True)
            df_topjt2 = topjt_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            topjt_1_name = df_topjt1.iloc[0]["Publication Name"]
            topjt_1_count = df_topjt1.iloc[0]["Total"]
    
            topjt_2_name = df_topjt2.iloc[0]["Publication Name"]
            topjt_2_count = df_topjt2.iloc[0]["Total"]
    
            # # Extract the top 3 publications and their counts
            # topjt_1 = Jour_table1.iloc[0:1]  # First publication
            # topjt_2 = Jour_table1.iloc[1:2]  # Second publication
            # topjt_3 = Jour_table1.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topjt1 = topjt_1.reset_index(drop=True)
            # df_topjt2 = topjt_2.reset_index(drop=True)
            # df_topjt3 = topjt_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topjt_1_name = df_topjt1.iloc[0]["Publication Name"]
            # # top_1_count = df_topjt1.iloc[0]["Total"]
    
            # topjt_2_name = df_topjt2.iloc[0]["Publication Name"]
            # # top_2_count = df_topjt2.iloc[0]["Total"]
    
            # topjt_3_name = df_topjt3.iloc[0]["Publication Name"]
            # # top_3_count = df_topjt3.iloc[0]["Total"]
    
            # Dynamically identify the client column
            # client_columns = [col for col in Jour_table1.columns if isinstance(col, str) and col.startswith("Client-")]
            # if client_columns:
            #     client_columns = client_columns[0]
            # else:
            #     raise ValueError("No columns starting with 'Client-' were found.")
            client_columns = [col for col in Unique_Articles1O.columns if col.startswith("Client-")][0]
            Unique_Articles1O = Unique_Articles1O.rename(columns={'Total': 'Total Unique Articles'})
            filtered_df = Unique_Articles[~Unique_Articles['Journalist'].isin(['Bureau News', 'Total'])].sort_values(by=client_columns, ascending = False)
            #journalist_name1 = filtered_df.iloc[0]["Journalist"]
            #publication_name1 = filtered_df.iloc[0]["Publication Name"]
            #client_count1 = filtered_df.iloc[0][client_column]
            #if len(filtered_df)>=2:
                #journalist_name2 = filtered_df.iloc[1]["Journalist"]
                #publication_name2 = filtered_df.iloc[1]["Publication Name"]
                #client_count2 = filtered_df.iloc[1][client_column]
                #if len(filtered_df)>=3:
                    #journalist_name3 = filtered_df.iloc[2]["Journalist"]
                    #publication_name3 = filtered_df.iloc[2]["Publication Name"]
                    #client_count3 = filtered_df.iloc[2][client_column]
                #else:
                    #journalist_name3 = ""
                    #publication_name3 = ""
                    #client_count3 = 0
            #else:
                #journalist_name2 = ""
                #publication_name2 = ""
                #client_count2 = 0
                #journalist_name3 = ""
                #publication_name3 = ""
                #client_count3 = 0
          
    
            # Find columns containing the word 'Client'
            client_columns = [col for col in Unique_Articles.columns if 'Client' in col]
            # Filter the dataframe where any 'Client' column has 0
            filtered_df = Unique_Articles[Unique_Articles[client_columns].eq(0).any(axis=1)]
    
            # Find the column with "Client" in its name
            # Dynamically identify the client column
            # client_columns = [col for col in Jour_table1.columns if isinstance(col, str) and col.startswith("Client")]
            # if client_columns:
            #     client_columns = client_columns[0]
            # else:
            #     raise ValueError("No columns starting with 'Client' were found.")
            client_columns = [col for col in Unique_Articles.columns if col.startswith('Client-') and col != 'Total Unique Articles'and col != 'Client %']
            competitor_columns = [col for col in Unique_Articles.columns if not col.startswith('Client-') and col != 'Total Unique Articles' and col!='Journalist' and col!= 'Publication Name' and col != 'Client %' ]
            filtered_df1 = Unique_Articles[(Unique_Articles[client_columns].gt(0).any(axis=1)) & (Unique_Articles[competitor_columns].eq(0).all(axis=1))  # All competitor columns should have zero values
]

            Jour_Comp = filtered_df.head(10)
            ordered_cols = ['Journalist', 'Publication Name',client_columndt] + [ent for ent in sov_order_no_client if ent in  Jour_Comp.columns] + (['Total Unique Articles'] if 'Total Unique Articles' in Jour_Comp.columns else [])
            
            Jour_Comp= Jour_Comp[ordered_cols]
    
            Jour_Client = filtered_df1.head(10)
            ordered_cols = ['Journalist', 'Publication Name',client_columndt] + [ent for ent in sov_order_no_client if ent in   Jour_Client.columns] 
            Jour_Client= Jour_Client[ordered_cols]
    
            # Extract the top 3 publications and their counts
            topjc_1 = Jour_Comp.iloc[0:1]  # First publication
            topjc_2 = Jour_Comp.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(Jour_Comp) > 2:
                topjc_3 = Jour_Comp.iloc[2:3]  # Third publication
                df_topjc3 = topjc_3.reset_index(drop=True)
                topjc_3_name = df_topjc3.iloc[0]["Journalist"]
                topjc_3_count = df_topjc3.iloc[0]['Total Unique Articles']
            else:
                topjc_3_name = ""
                topjc_3_count = 0  # You can assign any default value for count
    
            
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # Save them in separate DataFrames
            df_topjc1 = topjc_1.reset_index(drop=True)
            df_topjc2 = topjc_2.reset_index(drop=True)
            # df_topjc3 = topjc_3.reset_index(drop=True)
            # Initialize variables with default values
            topjc_1_name = "N/A"
            topjc_1_count = 0
            topjc_2_name = "N/A"
            topjc_2_count = 0

            # Extract publication name and count for the top 1, if available
            if not df_topjc1.empty:
                topjc_1_name = df_topjc1.iloc[0]["Journalist"]
                topjc_1_count = df_topjc1.iloc[0]['Total Unique Articles']

            # Extract publication name and count for the top 2, if available
            if not df_topjc2.empty:
                topjc_2_name = df_topjc2.iloc[0]["Journalist"]
                topjc_2_count = df_topjc2.iloc[0]['Total Unique Articles']
            
            # # Extract the top 3 journalits writing in comp and not on client and their counts
            # topjc_1 = Jour_Comp.iloc[0:1]  # First publication
            # topjc_2 = Jour_Comp.iloc[1:2]  # Second publication
            # topjc_3 = Jour_Comp.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topjc1 = topjc_1.reset_index(drop=True)
            # df_topjc2 = topjc_2.reset_index(drop=True)
            # df_topjc3 = topjc_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topjc_1_name = df_topjc1.iloc[0]["Journalist"]
            # topjc_1_count = df_topjc1.iloc[0]["Total"]
    
            # topjc_2_name = df_topjc2.iloc[0]["Journalist"]
            # topjc_2_count = df_topjc2.iloc[0]["Total"]
    
            # topjc_3_name = df_topjc3.iloc[0]["Journalist"]
            # topjc_3_count = df_topjc3.iloc[0]["Total"]
    
            # Extract the top 3 publications and their counts
            topjp_1 = Jour_Comp.iloc[0:1]  # First publication
            topjp_2 = Jour_Comp.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(Jour_Comp) > 2:
                topjp_3 = Jour_Comp.iloc[2:3]  # Third publication
                df_topjp3 = topjp_3.reset_index(drop=True)
                topjp_3_name = df_topjp3.iloc[0]["Publication Name"]
                topjp_3_count = df_topjp3.iloc[0]['Total Unique Articles']
            else:
                topjp_3_name = ""
                topjp_3_count = 0  # You can assign any default value for count
    
            
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # Save them in separate DataFrames
            df_topjp1 = topjp_1.reset_index(drop=True)
            df_topjp2 = topjp_2.reset_index(drop=True)
            # df_topjc3 = topjc_3.reset_index(drop=True)
            try:
                # Extract publication name and count for the top 3
                topjp_1_name = df_topjp1.iloc[0]["Publication Name"] if not df_topjp1.empty else "N/A"
                topjp_1_count = df_topjp1.iloc[0][client_column] if not df_topjp1.empty else 0

                topjp_2_name = df_topjp2.iloc[0]["Publication Name"] if not df_topjp2.empty else "N/A"
                topjp_2_count = df_topjp2.iloc[0][client_column] if not df_topjp2.empty else 0

            except IndexError:
                # Handle the case where the DataFrame is empty or index is out of bounds
                topjp_1_name = "N/A"
                topjp_1_count = 0
                topjp_2_name = "N/A"
                topjp_2_count = 0

            if len(Jour_Client)>=1:
                journalist_client1 = Jour_Client.iloc[0]["Journalist"]
                publication_client1 = Jour_Client.iloc[0]["Publication Name"]
                jour_client1 = Jour_Client.iloc[0][client_column]
                if len(Jour_Client)>=2:
                    journalist_client2 = Jour_Client.iloc[1]["Journalist"]
                    publication_client2 = Jour_Client.iloc[1]["Publication Name"]
                    jour_client2 = Jour_Client.iloc[1][client_column]
                    if len(Jour_Client)>=3:
                        journalist_client3 = Jour_Client.iloc[2]["Journalist"]
                        publication_client3 = Jour_Client.iloc[2]["Publication Name"]
                        jour_client3 = Jour_Client.iloc[2][client_column]
                    else:
                        journalist_client3 = ""
                        publication_client3 = ""
                        jour_client3 = 0
                else:
                    journalist_client2 = ""
                    publication_client2 = ""
                    jour_client2 = 0
                    journalist_client3 = ""
                    publication_client3 = ""
                    jour_client3 = 0
            else:
                journalist_client1 = ""
                publication_client1 = ""
                jour_client1 = 0
                journalist_client2 = ""
                publication_client2 = ""
                jour_client2 = 0
                journalist_client3 = ""
                publication_client3 = ""
                jour_client3 = 0
                
                
            pubs_table = pubs_table.rename(columns={'Total': 'Total Unique Articles'})

            # 2. Find the correct client column (Client-Tenet, Client-XYZ, etc.)
            client_col_pub = [col for col in pubs_table.columns if col.startswith("Client-")][0]

            # 3. CRITICAL: Remove the "Total" row + clean Publication Name in one go
            pubs_clean = (
                pubs_table
                .drop(pubs_table.index[-1])           # exclude Total row
                .copy()
            )

            # 4. Make sure Publication Name is always a clean string (this kills the float/NaN error)
            pubs_clean['Publication Name'] = (
                pubs_clean['Publication Name']
                .fillna('Unknown Publication')   # replace NaN
                .astype(str)                     # force everything to string
                .str.strip()                     # remove extra spaces/tabs
            )

            # 5. Filter publications where our client has 0 or 1 mention
            low_mention_pubs = pubs_clean[pubs_clean[client_col_pub].isin([0, 1])]

            # 6. Get top 3 by total unique articles
            top3_pubs = (
                low_mention_pubs
                .sort_values('Total Unique Articles', ascending=False)
                .head(3)
            )

            # 7. Extract publication names as clean strings
            top3_pubs_list = top3_pubs['Publication Name'].tolist()   # already strings!

            # 8. Build the perfect English string
            if len(top3_pubs_list) >= 2:
                publications_str = f"{', '.join(top3_pubs_list[:-1])} and {top3_pubs_list[-1]}"
            elif len(top3_pubs_list) == 1:
                publications_str = top3_pubs_list[0]
            else:
                publications_str = "None"
           
    
    
    
            # # Extract the top 3 publications and their counts
            # topjp_1 = Jour_Comp.iloc[0:1]  # First publication
            # topjp_2 = Jour_Comp.iloc[1:2]  # Second publication
            # topjp_3 = Jour_Comp.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topjp1 = topjp_1.reset_index(drop=True)
            # df_topjp2 = topjp_2.reset_index(drop=True)
            # df_topjp3 = topjp_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topjp_1_name = df_topjp1.iloc[0]["Publication Name"]
            # # top_1_count = df_topjp1.iloc[0]["Total"]
    
            # topjp_2_name = df_topjp2.iloc[0]["Publication Name"]
            # # top_2_count = df_topjp2.iloc[0]["Total"]
    
            # topjp_3_name = df_topjp3.iloc[0]["Publication Name"]
            # # top_3_count = df_topjp3.iloc[0]["Total"]
    
            
            
            # Remove square brackets and single quotes from the 'Journalist' column
            finaldata['Journalist'] = finaldata['Journalist'].str.replace(r"^\['(.+)'\]$", r"\1", regex=True)
            # Fill missing values in 'Influencer' column with 'Bureau News'
            # data['Journalist'] = data['Journalist'].fillna('Bureau News')
    
            # # Function to classify news exclusivity and topic
            # def classify_exclusivity(row):
            #     entity_name = row['Entity']
            #     if entity_name.lower() in row['Headline'].lower():
            #         return "Exclusive"
            #     else:
            #         return "Not Exclusive"
    
            def classify_exclusivity(row):
                entity_name = row['Entity']
                headline = row['Headline']
                
                # Ensure both entity_name and headline are strings
                if isinstance(entity_name, float) or isinstance(headline, float):
                    return "Not Exclusive"
                if str(entity_name).lower() in str(headline).lower():
                    return "Exclusive"
                else:
                    return "Not Exclusive"
    
    
    
            
    
            #finaldata['Exclusivity'] = finaldata.apply(classify_exclusivity, axis=1)
    
            # # Define a dictionary of keywords for each entity
            # entity_keywords = {
            #     'Amazon': ['Amazon', 'Amazons', 'amazon'],
            #     # Add other entities and keywords here
            # }
    
            # def qualify_entity(row):
            #     entity_name = row['Entity']
            #     text = row['Headline']
            #     if entity_name in entity_keywords:
            #         keywords = entity_keywords[entity_name]
            #         if any(keyword in text for keyword in keywords):
            #             return "Qualified"
            #     return "Not Qualified"
    
            # finaldata['Qualification'] = finaldata.apply(qualify_entity, axis=1)
    
            # Topic classification
            topic_mapping = {
                  'Merger': ['merger', 'merges'],
                    
                  'Acquire': ['acquire', 'acquisition', 'acquires'],
                    
                  'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
                    'Partnership': ['IPO','ipo'],
                   'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                        'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
                    
                   'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
                    
                  'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
                    
                  'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
                'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
                    
                   'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
                    
                   'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
                    
                    'Awards & Recognition': ['award', 'awards'],
                    
                    'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
                
                'Sale - Offers - Discounts' : ['sale','offers','discount','discounts','discounted']
            }
    
            def classify_topic(headline):
                for topic, words in topic_mapping.items():
                    if any(word in headline.lower() for word in words):
                        return topic
                return 'Other'
    
            #finaldata['Topic'] = finaldata['Headline'].apply(classify_topic)
    
            # Filter or select the row for which you need the client name
            filtered_rows = data[data["Entity"].str.contains("Client-", na=False)]
        
            # Check if any rows match and select the first one
            if not filtered_rows.empty:
                selected_row = filtered_rows.iloc[0]  # Get the first matching row
                entity = selected_row["Entity"]
                # Extract the brand name from the "Entity" column (after "Client-")
                client_name = entity.split("Client-")[-1]
            else:
                client_name = "Unknown Client"
    
            # Extract the brand name from the "Entity" column (after "Client-" if present)
            client_name = entity.split("Client-")[-1]
    
            dfs = [
    strip_client_prefix(Entity_SOV3),
    strip_client_prefix(sov_dt11),
    strip_client_prefix(pubs_table1),
    strip_client_prefix(Unique_Articles1O),
    strip_client_prefix(PType_Entity),
    strip_client_prefix(Jour_Comp),
    strip_client_prefix(Jour_Client),
]
            comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table','Pub Type and Entity Table','Journ-on Comp, not Client','Journ-on Client, not Comp']
    
            # Sidebar for download options
            st.sidebar.write("## Download Options")
            
            st.sidebar.write("## Download Combined Excel")
            file_name_all = st.sidebar.text_input("Enter file name for Combined Excel", "Combined Excel.xlsx")
            if st.sidebar.button("Download Combined Excel"):
                dfs = [
    strip_client_prefix(Entity_SOV3),
    strip_client_prefix(sov_dt11),
    strip_client_prefix(pubs_table2O),
    strip_client_prefix(Unique_Articles2O),
    strip_client_prefix(PType_Entity),
    strip_client_prefix(Jour_Comp),
    strip_client_prefix(Jour_Client),
]
                comments =['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table','Pub Type and Entity Table','Journ-on Comp, not Client','Journ-on Client, not Comp']
                        
                entity_info = f"""Entity:{client_name}
Time Period of analysis: {start_date} to {end_date}
Source: (Print) Factiva, 25 leading publications including Financial and mainline newspapers, Business and General Magazines..
News search: All Articles: entity mentioned at least once in the article"""
                excel_io_all = io.BytesIO()
                w1 = multiple_dfs(dfs, 'Tables', excel_io_all, comments, entity_info)
                excel_io_all.seek(0)
                wb = load_workbook(excel_io_all)
                pubs_table.at[pubs_table.index[-1], 'Publication Name'] = 'Total'
        
                dfs1 = [strip_client_prefix(pubs_table), strip_client_prefix(Unique_Articles)]
                comments1 = ['Publication Table', 'Journalist Table']
                multiple_dfs1(dfs1, 'All Pub-Jour', wb, comments1)  # <-- this writes directly into base workbook
                
                excel_io_2 = io.BytesIO()
                wb.save(excel_io_2)
                excel_io_2.seek(0)
                with pd.ExcelWriter(excel_io_2, engine='openpyxl', mode='a', if_sheet_exists='new') as writer:
                    create_entity_sheets(for_entity_data, writer)
                    writer.book.worksheets[0].title = "Report"
                wb_final = load_workbook(excel_io_2)
                all_sheets = wb_final.sheetnames
                client_sheet = next((s for s in all_sheets if s == client_name), None) or \
               next((s for s in all_sheets if s.startswith("Client-")), None)
                ordered_sheets = ['Report', 'All Pub-Jour']
                if client_sheet:
                    ordered_sheets.append(client_sheet)
                ordered_sheets += [s for s in sov_order_no_client if s in all_sheets]
                remaining_sheets = [s for s in all_sheets if s not in ordered_sheets]
                ordered_sheets += remaining_sheets
                wb_final._sheets = [wb_final[s] for s in ordered_sheets]
                final_io = io.BytesIO()
                wb_final.save(final_io)
                final_io.seek(0)
                combined_data = final_io.read()
                b64_all = base64.b64encode(combined_data).decode()
                href_all = (f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;'f'base64,{b64_all}" download="{file_name_all}">Download Combined Excel</a>')
                st.sidebar.markdown(href_all, unsafe_allow_html=True)
                
            st.write("## Preview Selected DataFrame")
            dataframes_to_download = {
"Entity_SOV1": strip_client_prefix(Entity_SOV3),
"Data": data,
"Finaldata": finaldata,
"Month-on-Month": strip_client_prefix(sov_dt11),
"Publication Table": strip_client_prefix(pubs_table),
"Journalist Table": strip_client_prefix(Unique_Articles),
"Publication Type Table with Entity": strip_client_prefix(PType_Entity),
"Entity-wise Sheets": finaldata,
"Journalist writing on Comp not on Client": strip_client_prefix(Jour_Comp),
"Journalist writing on Client & not on Comp": strip_client_prefix(Jour_Client)
}
            selected_dataframe = st.selectbox("Select DataFrame to Preview:", list(dataframes_to_download.keys()))
            st.dataframe(dataframes_to_download[selected_dataframe])

        
     # Load the image files
            img_path = r"New logo snip.png"
            img_path1 = r"New Templete main slide.png"
        
            # Create a new PowerPoint presentation with widescreen dimensions
            prs = Presentation()               
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
        
            # Add the first slide with the image
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path1, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
            # Add the text box above the image
            textbox_left = Inches(0.5)  # Adjust the left position as needed
            textbox_top = Inches(5)   # Adjust the top position as needed
            textbox_width = Inches(15)  # Adjust the width as needed
            textbox_height = Inches(1)  # Adjust the height as needed
        
            
        
            text_box = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), textbox_width, textbox_height)
            text_frame = text_box.text_frame
            text_frame.text = (f"{client_name}\n"
                              "News Analysis\n"
                            "By Media Research & Analytics Team")
        
                
            # Set font size to 30 and make the text bold and white
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(50)
                    run.font.bold = True
        #           run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White color
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            # Add title slide after the first slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
            left = Inches(0.0)  # Adjust the left position as needed
            top = prs.slide_height - Inches(1)  # Adjust the top position as needed
            slide.shapes.add_picture(img_path, left, top, height=Inches(1))  # Adjust the height as needed 
        
                
            # Clear existing placeholders
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.clear()  # Clear existing text frames
        
            # Set title text and format for Parameters slide
            header_text = "Parameters"
            header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.7))
            header_frame = header_shape.text_frame
            header_frame.text = header_text
        
            for paragraph in header_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = header_text
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
            # Add Time Period text
            time_period_text = f"""Time Period : {start_date} to {end_date}"""
            time_period_shape = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(14), Inches(0.5))
            time_period_frame = time_period_shape.text_frame
            time_period_frame.text = time_period_text
            # time_period_frame.paragraphs[0].font.bold = True
            time_period_frame.paragraphs[0].font.size = Pt(24)
            time_period_frame.paragraphs[0].font.name = 'Gill Sans'
        
        
           # Add Source text
            source_shape = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))
            source_frame = source_shape.text_frame
            source_frame.word_wrap = True  

            # First line (Online)
            p1 = source_frame.add_paragraph()
            p1.text = "Source: (Online) Kalki All publications."
            p1.font.size = Pt(24)
            p1.font.name = 'Gill Sans'

            # Second line (Print / Factiva)
            p2 = source_frame.add_paragraph()
            p2.text = "(Print) Factiva, 25 leading publications including Financial and mainline newspapers, Business and General Magazines."
            p2.font.size = Pt(24)
            p2.font.name = 'Gill Sans'
        
            # Add News Search text
            news_search_text = "News Search : All Articles: entity mentioned at least once in the article "
            news_search_shape = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(10), Inches(0.75))  # Adjusted width and height
            news_search_frame = news_search_shape.text_frame
            news_search_frame.word_wrap = True  # Enable text wrapping
            p2 = news_search_frame.add_paragraph()  # Create a paragraph for text
            p2.text = news_search_text  # Set the text
        
            # Set font properties after text is added
            # p2.font.bold = True
            p2.font.size = Pt(24)
            p2.font.name = 'Gill Sans'  # Changed to Arial for compatibility
                
            # Add the first slide with the image
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path1, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
            # Add the text box above the image
            textbox_left = Inches(0.5)  # Adjust the left position as needed
            textbox_top = Inches(5)   # Adjust the top position as needed
            textbox_width = Inches(15)  # Adjust the width as needed
            textbox_height = Inches(1)  # Adjust the height as needed
        
            text_box = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), textbox_width, textbox_height)
            text_frame = text_box.text_frame
            text_frame.text = "Online Media/Print Media"
        
            # Set font size to 30 and make the text bold and white
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(50)
                    run.font.bold = True
        #           run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White color
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
            # Add title slide after the first slide
            #slide_layout = prs.slide_layouts[6]
            #slide = prs.slides.add_slide(slide_layout)
        
            #left = Inches(0.0)  # Adjust the left position as needed
            #top = prs.slide_height - Inches(1)  # Adjust the top position as needed
            #slide.shapes.add_picture(img_path, left, top, height=Inches(1))  # Adjust the height as needed 
                 
            # Clear existing placeholders
            #for shape in slide.placeholders:
                #if shape.has_text_frame:
                    #shape.text_frame.clear()  # Clear existing text frames
        
            # Set title text and format for Parameters slide
            #header_text = "Inferences and Recommendations"
            #header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(14), Inches(0.7))
            #header_frame = header_shape.text_frame
            #header_frame.text = header_text
            #for paragraph in header_frame.paragraphs:
                #for run in paragraph.runs:
                    #run.text = header_text
                    #run.font.size = Pt(30)
                    #run.font.bold = True
                    #run.font.name = 'Helvetica'
                    #run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    #paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    #paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  
        
        
            # Add SOV text
            #sov_text = ("Share of Voice :")
            #sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
            #sov_text_frame = sov_text_shape.text_frame
            #sov_text_frame.word_wrap = True
            #sov_text_frame.clear()  # Clear any default paragraph
        
            #p = sov_text_frame.add_paragraph()
            #p.text = "Share of Voice :"
            #p.font.size = Pt(20)
            #p.font.name = 'Gill Sans'
            #p.font.underline = True
            #p.font.bold = True
        
            #sov_text 
        
        
            #p = sov_text_frame.add_paragraph()
            #p.text 
            #p.font.size = Pt(18)
            #p.font.name = 'Gill Sans'
        
        #     # Add Source text
        #     source_text = ("Publications :")
        #     source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.8), Inches(14), Inches(1))
        #     source_frame = source_shape.text_frame
        #     source_frame.word_wrap = True
        #     source_frame.clear()  # Clear any default paragraph
        #     p = source_frame.add_paragraph()
        #     p.text = "Publications :"
        #     p.font.size = Pt(20)
        #     p.font.name = 'Gill Sans'
        #     p.font.underline = True
        #     p.font.bold = True
        
        
        #     source_text = (
        #     f"•The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} articles, followed by {top_2_name} with {top_2_count} articles, and {top_3_name} with {top_3_count} articles.\n"
        # f"•Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
        # f"•The top 10 publications writing articles on {client_name} contribute {top10_pub_perc}% (which is {top10_pub_sum} of the total {client_sov_count}  articles).\n" 
        # )
        #     source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.1), Inches(14), Inches(1))
        #     source_frame = source_shape.text_frame
        #     source_frame.word_wrap = True
        #     source_frame.clear()  # Clear any default paragraph
        #     p = source_frame.add_paragraph()
        #     p.text = (
        #     f"•The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} articles, followed by {top_2_name} with {top_2_count} articles, and {top_3_name} with {top_3_count} articles.\n"
        # f"•Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
        # f"•The top 10 publications writing articles on {client_name} contribute {top10_pub_perc}% (which is {top10_pub_sum} of the total {client_sov_count} articles).\n" 
        # )
        #     p.font.size = Pt(18)
        #     p.font.name = 'Gill Sans'

              # Add News Search text
            #news_search_text = ("Publication Types :" )
            #news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.6), Inches(14), Inches(0.5))
            #news_search_frame = news_search_shape.text_frame
            #news_search_frame.word_wrap = True
            #news_search_frame.clear()  # Clear any default paragraph
            #p = news_search_frame.add_paragraph()
            #p.text = "Publication Type :"
            #p.font.size = Pt(20)
            #p.font.name = 'Gill Sans'
            #p.font.underline = True
            #p.font.bold = True
        
         
            #news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.0), Inches(14), Inches(0.5))
            #news_search_frame = news_search_shape.text_frame
            #news_search_frame.word_wrap = True
            #news_search_frame.clear()  # Clear any default paragraph
            #p = news_search_frame.add_paragraph()
            #
            #p.font.size = Pt(18)
            #p.font.name = 'Gill Sans'
        
            # Add title slide after the first slide
            #slide_layout = prs.slide_layouts[6]
            #slide = prs.slides.add_slide(slide_layout)
        
        
            # Clear existing placeholders
            #for shape in slide.placeholders:
                #if shape.has_text_frame:
                    #shape.text_frame.clear()  # Clear existing text frames
        
        
            # Set title text and format for Parameters slide
            #header_text = "Inferences and Recommendations"
            #header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
            #header_frame = header_shape.text_frame
            #header_frame.text = header_text 
            #for paragraph in header_frame.paragraphs:
                #for run in paragraph.runs:
                    #run.text = header_text
                    #run.font.size = Pt(30)
                    #run.font.bold = True
                    #run.font.name = 'Helvetica'
                    #run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    #paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    #paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
        
            # Add News Search text
            #news_search_text = ("Journalists :")
            #news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
            #news_search_frame = news_search_shape.text_frame
            #news_search_frame.word_wrap = True
            #news_search_frame.clear()  # Clear any default paragraph
            #p = news_search_frame.add_paragraph()
           #p.text = "Journalists :"
            #p.font.size = Pt(20)
            #p.font.name = 'Gill Sans'
            #p.font.underline = True
            #p.font.bold = True
        
            # Add News Search text
          
            #news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
            #news_search_frame = news_search_shape.text_frame
            #news_search_frame.word_wrap = True
            #news_search_frame.clear()  # Clear any default paragraph
            #p = news_search_frame.add_paragraph()
          
            #p.font.size = Pt(18)
            #p.font.name = 'Gill Sans'

           
        
            # Sidebar for PowerPoint download settings
            st.sidebar.write("## Download All DataFrames as a PowerPoint File")
            pptx_file_name = st.sidebar.text_input("Enter file name for PowerPoint", "dataframes_presentation.pptx")
            if st.sidebar.button("Download PowerPoint"):
                # List of DataFrames to save
                pubs_table1 = pubs_table.head(10)
                numeric_columns = pubs_table1.select_dtypes(include=['number']).columns
                pubs_table1[numeric_columns] = pubs_table1[numeric_columns].astype(int)
                Jour_table1 = Jour_table.head(10)
                dfs = [
    strip_client_prefix(Entity_SOV3),
    strip_client_prefix(sov_dt11),
    strip_client_prefix(pubs_table1),
    strip_client_prefix(Jour_table1),
    strip_client_prefix(PType_Entity),
    strip_client_prefix(Jour_Comp),
    strip_client_prefix(Jour_Client),
]
                table_titles = [f'SOV Table of {client_name} and competition', f'Month-on-Month Table of {client_name} and competition', f'Publication Table on {client_name} and competition', f'Journalist writing on {client_name} and competition',
                            f'Publication Types writing on {client_name} and competition',f'Journalists writing on Comp and not on {client_name}', f'Journalists writing on {client_name} and not on Comp'
                            ]
                textbox_text = [" "," "," "," "," "," "," "]
              
        
                # Loop through each DataFrame and create a new slide with a table
                for i, (df, title) in enumerate(zip(dfs, table_titles)):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_table_to_slide(slide, df, title, textbox_text[i])
                    if i == 0:
                        img_path4 = generate_bar_chart(dfs[0])  # Generate chart from first DataFrame
                        add_image_to_slide(slide, img_path4)
                    if i == 1:  
                        img_path5 = generate_line_graph(sov_dt1)  # Generate chart from first DataFrame
                        add_image_to_slide1(slide, img_path5)
                    if i == 4:  
                        img_path6 = generate_bar_pchart(dfs[4])  # Generate chart from first DataFrame
                        add_image_to_slide2(slide, img_path6)
                    if i == 6:
                        wordcloud_path = generate_word_cloud(finaldata)  # Generate word cloud from DataFrame
                        add_image_to_slide11(slide, wordcloud_path)
                      
                # Save presentation to BytesIO for download
                pptx_output = io.BytesIO()
                prs.save(pptx_output)
                pptx_output.seek(0)
                st.sidebar.download_button(
                    label="Download PowerPoint Presentation",
                    data=pptx_output,
                    file_name=pptx_file_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        #         # Download All DataFrames as a Single Excel Sheet
        #         st.sidebar.write("## Download All DataFrames as a Single Excel Sheet")
        #         file_name_all = st.sidebar.text_input("Enter file name for all DataFrames", "all_dataframes.xlsx")
        # #         download_options = st.sidebar.selectbox("Select Download Option:", [ "Complete Dataframes"])
                
        #         if st.sidebar.button("Download All DataFrames"):
        #             # List of DataFrames to save
        #             dfs = [Entity_SOV1, sov_dt, pubs_table, Jour_table, PType_Entity, PP_table, ppe1]
        #             comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
        #                         'Pub Type and Entity Table', 'Pub Type and Pub Name Table',
        #                         'PubType PubName and Entity Table']
                    
        #             excel_path_all = os.path.join(download_path, file_name_all)
        #             multiple_dfs(dfs, 'Tables', excel_path_all, 2, comments)
        #             st.sidebar.write(f"All DataFrames saved at {excel_path_all}")
        
        #         # Loop through each dataframe and create a new slide for each one
        #         for i, (df, title) in enumerate(zip(dfrs, table_titles)):
        #             slide = prs.slides.add_slide(prs.slide_layouts[6])
        #             add_table_to_slide(slide, df, title, textbox_text[i])
                        
            #
            # ------------------------------------------------------------------
            # === DOWNLOAD GROK PROMPTS (.docx) ===
            st.sidebar.write("## Download Grok Prompts (.docx)")
            from docx import Document
            from docx.shared import RGBColor, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            import io
            import base64

            if st.sidebar.button("Download Prompts"):
               
                # === COLOR PALETTE (NO RED) ===
                CLIENT_NAME_COLOR = RGBColor(0, 102, 255) # Blue for Client Name
                DATE_COLOR = RGBColor(0, 128, 0) # Green for Dates
                COMPETITOR_COLOR = RGBColor(255, 140, 0) # Orange for Competitors
                INDUSTRY_COLOR = RGBColor(128, 0, 128) # Purple for Industry
                PUB_COLOR = RGBColor(255, 20, 147) # Deep Pink for Publications
                JOURNALIST_COLOR = RGBColor(225, 167, 63) # Gold for Journalists
                BLACK = RGBColor(0, 0, 0)
                
                def add_run(p, txt, color=BLACK, bold=False):
                    """Add text with color and optional bold"""
                    r = p.add_run(txt)
                    r.font.color.rgb = color
                    r.bold = bold
                    return p
                
                doc = Document()
                
                # Short variables
                s = start_date
                e = end_date
                c = client_name
                comp = competitors_str
                ind = industry
                
                # === COLOR LEGEND ===
                p = doc.add_paragraph()
                add_run(p, "Color code: ", bold=True)
                add_run(p, "Client name, ", CLIENT_NAME_COLOR, bold=True)
                add_run(p, "Dates, ", DATE_COLOR, bold=True)
                add_run(p, "Competitor name, ", COMPETITOR_COLOR, bold=True)
                add_run(p, "Industry name, ", INDUSTRY_COLOR, bold=True)
                add_run(p, "Publications writing on Industry", PUB_COLOR, bold=True)
                add_run(p, "Journalists writing on Industry", JOURNALIST_COLOR, bold=True)

                p = doc.add_paragraph()
                p.add_run("I work in Media Research Team at a PR Company, I will be sharing the below Qualitative insights with the PR professionals. Please keep this in mind and provide insights accordingly.")         
                
                # === REQUIREMENTS SECTION ===
                p = doc.add_paragraph()
                p.add_run("Satisfy the below requirements :").bold = True
                
                # Prompt 1,2 requirements
                p = doc.add_paragraph()
                p.add_run("Prompt 1 follow the below requirements").bold = True
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("Please give topicwise/bucketwise paragraph with topic/bucket highlighted please be very much elaborative as possible")
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("Max 2 sentences should be in a line then move into next line and follow the same thing consider them as a point and the points should always be elaborative and not in one liner")
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("In each topicwise/bucketwise paragraph and in each and every point in the paragraph, the content should be very much  elaborative as possible and there should be atleast 5-6 such points without losing relevant news in each topic and those should be elaborative")
                
                # NOTE for prompt 3
                p = doc.add_paragraph()
                r = p.add_run("NOTE : ")
                r.bold = True
                p.add_run("Don't provide insight for prompt 3")
                
                # Prompt 4,5 requirements
                p = doc.add_paragraph()
                p.add_run("Prompt 4,5 follow the below requirements").bold = True
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("Please give the insights in tabular format, don't mention no. of articles anywhere) Give one elaborative paragraph before creating a table.")
                
                # Prompt 6,7 requirements
                p = doc.add_paragraph()
                p.add_run("Prompt 6,7 follow the below requirements").bold = True
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("Give five or more critiques pointwise (be elaborative as much as possible) with the critiques pointer headers highlighted (just have 2-3 elaborative pointers in each critiques). Give one paragraph elaborative description before giving the pointers for critiques.")
                
                # Prompt 9,10 requirements
                p = doc.add_paragraph()
                p.add_run("Prompt 9,10 follow the below requirements").bold = True
                
                p = doc.add_paragraph(style='List Bullet')
                r = p.add_run("Give paragraph wise for each Publications/Journalist(along with publication name) mentioned here and highlight the Publication name/Journalist,please be very much elaborative don't be generic relate it with the news released by these publications/journalist on the ")
                add_run(p, ind, INDUSTRY_COLOR)
                p.add_run(" industry. (")
                r = p.add_run("Note : ")
                r.bold = True
                p.add_run("Publications/Journalist(along with publication name) should be highlighted)")
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("Don't give me the number of news written by these publications/journalists  just the content that has been written on the ")
                add_run(p, ind, INDUSTRY_COLOR)
                p.add_run(" industry.")
                
                p = doc.add_paragraph(style='List Bullet')
                p.add_run("Additionally, follow the same requirements given for Prompt 1 excluding the first point in it (i.e Please give topicwise/bucketwise paragraph with topic/bucket)")
                
                doc.add_paragraph()

                p = doc.add_paragraph()
                r = p.add_run("For Prompts 1-10 ")
                r.bold = True
                p.add_run("please note that : do not consider press release from the companys website or any social media platform")                # Add formatting instruction before prompts

                p = doc.add_paragraph()
                p.add_run("Before giving insight for each prompt mention the title like Conversations on Client company, Topicwise Conversations on Client company,… etc (and replace Client Company with ")
                add_run(p, c, CLIENT_NAME_COLOR)
                p.add_run(" and Industry with ")
                add_run(p, ind, INDUSTRY_COLOR)
                p.add_run(")  with formatting ### and no bold formatting and topics/buckets with bold formatting **,please follow the formatting strictly don't use bold formatting in the content of any of the buckets/topics")
                
                doc.add_paragraph()
                
                # === PROMPT 1 ===
                p = doc.add_paragraph()
                add_run(p, "1) Conversations on Client company - ", bold=True)
                add_run(p, "Could you Summarize the news articles from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, "? Please summarize as many topics as possible but do not consider press release from the companys website or any social media platform. Only summarize the articles from print and online news platforms.")
                
                # === PROMPT 2 ===
                p = doc.add_paragraph()
                add_run(p, "2) Topicwise Conversation on Client Company - ", bold=True)
                add_run(p, "Could you Summarize the news articles from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, ". Please summarize the news articles as per the following categories. I am giving you the buckets. Please arrange the news as per their content in the relevant buckets and summarize that news. Only summarize the news articles from print and online news platforms. The buckets are as follows: Financial Performance, Product and Services, Social Good (includes CSR, ESG, Philanthropy, Environment), Employee Engagement (includes hiring, resignation, layoffs, training, skilling, employee benefits, appraisals...), Business Strategy (include growth, mergers, future, market share...), Vision and Leadership (Interviews, interaction, thought leadership, authored articles...), Legal and Regulatory, Tech & innovation, Stock related (stock recommendations, stock movements).( ")
                r = p.add_run("Note : ")
                r.bold = True
                add_run(p, "Please give topicwise/bucketwise paragraph with topic/bucket highlighted please be very much elaborative as possible. Max 2 sentences should be in a line then move into next line and follow the same thing consider them as a point and the points should always be elaborative and not in one liner. In each topicwise/bucketwise paragraph and in each and every point in the paragraph, the content should be very much  elaborative as possible and there should be atleast 5-6 such points without losing relevant news in each topic and those should be elaborative)")
                
                # === PROMPT 3 ===
                p = doc.add_paragraph()
                add_run(p, "3) Topicwise Conversation on Competitor Company - ", bold=True)
                add_run(p, "Could you Summarize the news articles from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " for ")
                add_run(p, comp, COMPETITOR_COLOR)
                add_run(p, "? Give Topicwise conversation entity wise, one after the other follow the bucket structure for each of the entities, want separate separate topic wise conversation for each entity. Please summarize the news articles as per the following categories. I am giving you the buckets. Please arrange the news as per their content in the relevant buckets and summarize that news. Only summarize the news articles from print and online news platforms. The buckets are as follows: Financial Performance, Product and Services, Social Good (includes CSR, ESG, Philanthropy, Environment), Employee Engagement (includes hiring, resignation, layoffs, training, skilling, employee benefits, appraisals...), Business Strategy (include growth, mergers, future, market share...), Vision and Leadership (Interviews, interaction, thought leadership, authored articles...), Legal and Regulatory, Tech & innovation, Stock related (stock recommendations, stock movements). (")
                r = p.add_run("Note : ")
                r.bold = True
                add_run(p, "Please give topicwise/bucketwise paragraph with topic/bucket highlighted please be very much elaborative as possible. Max 2 sentences should be in a line then move into next line and follow the same thing consider them as a point and the points should always be elaborative and not in one liner. In each topicwise/bucketwise paragraph and in each and every point in the paragraph, the content should be very much  elaborative as possible and there should be atleast 5-6 such points without losing relevant news in each topic and those should be elaborative) Give it with the Header i.e Topicwise Conversation on Competitor Company with formatting ### and competitor name ")
                add_run(p, comp, COMPETITOR_COLOR)
                add_run(p, "  with formatting ## with no bold formatting and buckets with bold formatting **, please follow the formatting strictly don't use bold formatting in the content of any of the buckets/topics")
                
                # === PROMPT 4 ===
                p = doc.add_paragraph()
                add_run(p, "4) Month – on – Month Insights - ", bold=True)
                add_run(p, "Give me month on month breakdown of news coverage for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, ". Give me details of events that have lead to a spike in the media coverage and if no news is present for a particular month mention that in the table itself.")
                r = p                
                # === PROMPT 5 ===
                p = doc.add_paragraph()
                add_run(p, "5) Unique Conversations by Competitors - ", bold=True)
                add_run(p, "What factors contributed to ")
                add_run(p, comp, COMPETITOR_COLOR)
                add_run(p, " higher media coverage in the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry between ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " compared to ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, "? Identify the unique conversation topics (topics where ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " is not mentioned) where these companies were mentioned in the Headline or the lead para or if they were mentioned atleast twice in the article that have driven the higher media coverage.  (")
                r = p.add_run("Note : ")
                r.bold = True
                add_run(p, "give the insights in tabular format with Column name : Company, Unique Conversation Topics be elaborative relating it with the news, (for each row there should be just 1 company name all its unique conversation should be there beside it in the Unique Conversation Topics column) Give one elaborative paragraph before creating a table.")
                
                # === PROMPT 6 ===
                p = doc.add_paragraph()
                add_run(p, "6) Reputational Risks for Client - ", bold=True)
                add_run(p, "What is the online news media saying about ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, "? Give five or more critiques.")
                
                # === PROMPT 7 ===
                p = doc.add_paragraph()
                add_run(p, "7) Reputational Risks for Industry - ", bold=True)
                add_run(p, "What is the online news media saying about ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, "? Give five or more critiques.")
                
                # === PROMPT 8 ===
                p = doc.add_paragraph()
                add_run(p, "8) Industry Snapshot - ", bold=True)
                add_run(p, "What is the online news media conversation in the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " and identify the companies who are a part of these conversations.(")
                r = p.add_run("Note : ")
                r.bold = True
                add_run(p, "give the insights in tabular format and be elaborative relating it with the news, give company wise insight for this, don't put more than 1 company in one row) Give one elaborative paragraph before creating a table. You can also consider companies who are not the competitors but are present online news media conversation in the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry.")
                
                # === PROMPT 9 ===
                p = doc.add_paragraph()
                add_run(p, "9) Publications writing on Industry – ", bold=True)
                add_run(p, "Which Indian publications (min 3) have frequently written on  ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " and what are the conversations about the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry AND which ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " companies have been mentioned by them?")
                
                # === PROMPT 10 ===
                p = doc.add_paragraph()
                add_run(p, "10) Journalist writing on Industry – ", bold=True)
                add_run(p, "Which Indian journalists name (min 3) have (mention their Publication name too) frequently written on ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " and what are the conversations about the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry AND which ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " companies have been mentioned by them?")
                
                # === PROMPT 11 ===
                p = doc.add_paragraph()
                add_run(p, "11) X Insights - ", bold=True)
                add_run(p, "What is being said on Twitter X about ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " between ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, "? (")
                r = p.add_run("Note : ")
                r.bold = True
                add_run(p, "Please be very much elaborative , should be majorly on users conversation around ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, ", if possible highlight the user whose content about ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " has larger influence among the twitter audience and give breakdown of Positive Discussions, Criticisms and Complaints, Neutral/Informational Mentions, and highlight it, please be as elaborative as possible)")
                
                # Final NOTE
                p = doc.add_paragraph()
                r = p.add_run("NOTE : ")
                r.bold = True
                p.add_run("Don't provide insight for prompt 3")
                
                # === SAVE & DOWNLOAD ===
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                b64 = base64.b64encode(buffer.read()).decode()
                href = f'<a href="data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,{b64}" download="Grok_Prompts_{client_name}.docx">Download Grok Prompts (.docx)</a>'
                st.sidebar.markdown(href, unsafe_allow_html=True)
    
    # Continue with the rest of your prompts (2-12) as before...
            # Download selected DataFrame
            st.sidebar.write("## Download Selected DataFrame")
            dataframes_to_download = {
                    "Entity_SOV1": Entity_SOV3,
                    "Data": data,
                    "Finaldata": finaldata,
                    "Month-on-Month":sov_dt11,
                    "Publication Table":pubs_table,
                   "Journalist Table": Unique_Articles,
                    # "Publication Type and Name Table":PP_table,
                    "Publication Type Table with Entity":PType_Entity,
                    # "Publication type,Publication Name and Entity Table":ppe1,
                    "Entity-wise Sheets": finaldata,                            # Add this option to download entity-wise sheets
                    "Journalist writing on Comp & not on Client" : Jour_Comp, 
                    "Journalist writing on Client & not on Comp" : Jour_Client
                }
            selected_dataframe = st.sidebar.selectbox("Select DataFrame:", list(dataframes_to_download.keys()))
            if st.sidebar.button("Download Selected DataFrame"):
                if selected_dataframe in dataframes_to_download:
                    selected_df = dataframes_to_download[selected_dataframe]
                    excel_io_selected = io.BytesIO()
                    with pd.ExcelWriter(excel_io_selected, engine="openpyxl") as writer:
                        selected_df.to_excel(writer, index=True)
                        excel_io_selected.seek(0)
                        b64_selected = base64.b64encode(excel_io_selected.read()).decode()
                        href_selected = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_selected}" download="{selected_dataframe}.xlsx">Download {selected_dataframe} Excel</a>'
                        st.sidebar.markdown(href_selected, unsafe_allow_html=True) 
                        
            st.sidebar.write("## Download Data")            
            download_formats = st.sidebar.selectbox("Select format:", ["Excel", "CSV", "Excel (Entity Sheets)"])
    
            if st.sidebar.button("Download Data"):
                if download_formats == "Excel":
                    # Download all DataFrames as a single Excel file
                    excel_io = io.BytesIO()
                    with pd.ExcelWriter(excel_io, engine="openpyxl") as writer:
                        for df, comment in zip(dfs, comments):
                            df.to_excel(writer, sheet_name=comment, index=False)
                    excel_io.seek(0)
                    b64_data = base64.b64encode(excel_io.read()).decode()
                    href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="data.xlsx">Download Excel</a>'
                    st.sidebar.markdown(href_data, unsafe_allow_html=True)
    
                elif download_formats == "CSV":
                    # Download all DataFrames as CSV
                    csv_io = io.StringIO()
                    for df in dfs:
                        df.to_csv(csv_io, index=False)
                    csv_io.seek(0)
                    b64_data = base64.b64encode(csv_io.read().encode()).decode()
                    href_data = f'<a href="data:text/csv;base64,{b64_data}" download="data.csv">Download CSV</a>'
                    st.sidebar.markdown(href_data, unsafe_allow_html=True)
    
                elif download_formats == "Excel (Entity Sheets)":
                    # Download DataFrames as Excel with separate sheets by entity
                    excel_io = io.BytesIO()
                    with pd.ExcelWriter(excel_io, engine="openpyxl") as writer:
                        create_entity_sheets(finaldata, writer)
                    excel_io.seek(0)
                    b64_data = base64.b64encode(excel_io.read()).decode()
                    href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="entity_sheets.xlsx">Download Entity Sheets</a>'
                    st.sidebar.markdown(href_data, unsafe_allow_html=True)
    else:
        st.sidebar.write("No file uploaded yet")


from wordcloud import WordCloud
from PIL import Image
from fuzzywuzzy import fuzz
import matplotlib.pyplot as plt
# import gensim
# import spacy
# import pyLDAvis.gensim_models
# from gensim.utils import simple_preprocess
# from gensim.models import CoherenceModel
from pprint import pprint
import logging
import warnings
warnings.filterwarnings("ignore", category=UserWarning, module="xlsxwriter")

from nltk.corpus import stopwords
# import gensim.corpora as corpora
from io import BytesIO
import nltk

# Download NLTK stopwords
nltk.download('stopwords')

# Set up logging
logging.basicConfig(format='%(asctime)s : %(levelname)s : %(message)s', level=logging.ERROR)

# Ignore warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

# Initialize Spacy 'en' model
# nlp = spacy.load('en_core_web_sm', disable=['parser', 'ner'])

stop_words = set(stopwords.words('english'))
# Define a function to clean the text


# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("SimilarNews , Wordcloud and Topic Explorer")

st.sidebar.markdown(
    "<hr style='border-top: 5px dotted black; margin-top: 100px; margin-bottom: 10px;'/>",
    unsafe_allow_html=True
)

st.sidebar.title("**Other Analysis :**")

# --- Callback functions to update the latest file indicator ---
def update_latest_wordcloud():
    st.session_state.latest = "wordcloud"

def update_latest_similarity():
    st.session_state.latest = "similarity"

def update_latest_qualitative():
    st.session_state.latest = "qualitative"

# --- Qualitative Report Section ---
# File uploader for WordCloud with an on_change callback
st.sidebar.markdown(
    "<p style='font-size: 20px; font-weight: bold;'>Qualitative Report</p>",
    unsafe_allow_html=True
)

uploaded_docx = st.sidebar.file_uploader(
    "Upload Word Document (.docx)",
    type=["docx"],
    key="qualitative_docx_uploader",
    on_change=update_latest_qualitative
)
if uploaded_docx is not None:
    st.session_state['docx_bytes'] = uploaded_docx.read()
# ── Shared helper functions (used by both PPT and Word generators) ──
def is_table_start(line):
    stripped = line.strip()
    return stripped.startswith("|") and "|" in stripped[1:]

def parse_markdown_table(lines):
    rows = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if any(cells):
                rows.append(cells)
    if len(rows) > 1 and all(set(c) <= {"-", ":", " "} for c in rows[1]):
        rows.pop(1)
    return rows

def remove_bold_markers(text):
    return text.replace("**", "")

def is_topic_heading(text):
    stripped = text.strip()
    return stripped.startswith("**") and stripped.endswith("**")

def is_competitor_marker(text):
    stripped = text.strip()
    return stripped.startswith("##")
if uploaded_docx is not None:
    if st.sidebar.button("Generate PPT Report"):
        with st.spinner("Generating PowerPoint presentation..."):
            try:
                import io
                from docx import Document
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                
                # ========================== CONFIGURATION ==========================
                TEMPLATE_IMAGE = "Template Image.png"
                FIRST_SLIDE_IMAGE = "First slide.png"
                DISCLAIMER_IMAGE = "Disclaimer.png"
                
                SLIDE_WIDTH = Inches(13.333)
                SLIDE_HEIGHT = Inches(7.5)
                
                LOGO_SIZE = Inches(1.0)
                LOGO_LEFT = Inches(0.0)
                LOGO_BOTTOM_MARGIN = Inches(0.0)
                
                HEADER_TOP = Inches(0.15)
                HEADER_HEIGHT = Inches(0.65)
                
                CONTENT_LEFT = Inches(0.4)
                CONTENT_TOP = Inches(0.50)
                CONTENT_WIDTH = SLIDE_WIDTH - Inches(0.8)
                
                MAX_CONTENT_HEIGHT_PT = 460
                
                ORANGE = RGBColor(255, 140, 0)
                
                FONT_HEADING = 18
                FONT_COMPETITOR = 22
                FONT_NORMAL = 14
                
                # =====================================================================
                
                
                
                def estimate_line_height(text, font_size, is_heading=False, is_competitor=False):
                    chars_per_line = 100
                    lines_needed = max(1, len(text) / chars_per_line)
                    spacing = 1.25
                    height = lines_needed * font_size * spacing
                    
                    if is_heading:
                        height += 12
                    if is_competitor:
                        height += 20
                    
                    return height
                
                # Load uploaded DOCX
                # ✅ New
                doc = Document(io.BytesIO(st.session_state['docx_bytes']))
                full_text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
                
                sections = []
                current_title = None
                current_body = []
                
                for line in full_text.splitlines():
                    stripped = line.strip()
                    if stripped.startswith("###"):
                        if current_title:
                            sections.append({"title": current_title, "body": current_body})
                        current_title = stripped[3:].strip()
                        current_body = []
                    elif current_title and stripped:
                        current_body.append(line)
                
                if current_title:
                    sections.append({"title": current_title, "body": current_body})
                
                # Extract client name
                if sections:
                    first_title = sections[0]["title"].strip()
                    
                    prefixes_to_remove = [
                        "Conversations on ",
                        "Topicwise Conversation on ",
                        "Exclusive Conversation on ",
                        "Month – on – Month Insights ",
                    ]
                    
                    client_name = first_title
                    for prefix in prefixes_to_remove:
                        if client_name.startswith(prefix):
                            client_name = client_name[len(prefix):].strip()
                            break
                    
                    if not client_name:
                        client_name = "Client"
                    
                    output_filename = f"{client_name} Qualitative Insights Report.pptx"
                else:
                    client_name = "Client"
                    output_filename = "Qualitative Insights Report.pptx"
                
                # CREATE POWERPOINT
                prs = Presentation()
                prs.slide_width = SLIDE_WIDTH
                prs.slide_height = SLIDE_HEIGHT
                blank_layout = prs.slide_layouts[6]
                
                # Add Custom First Slide
                if os.path.exists(FIRST_SLIDE_IMAGE):
                    first_slide = prs.slides.add_slide(blank_layout)
                    first_slide.shapes.add_picture(FIRST_SLIDE_IMAGE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
                    
                    upper_text = f"{client_name} Insights"
                    upper_tb = first_slide.shapes.add_textbox(Inches(1.5), Inches(0.5), SLIDE_WIDTH - Inches(1.8), Inches(1.2))
                    tf_upper = upper_tb.text_frame
                    tf_upper.word_wrap = True
                    tf_upper.auto_size = True
                    p_upper = tf_upper.add_paragraph()
                    p_upper.text = upper_text
                    p_upper.alignment = PP_ALIGN.LEFT
                    p_upper.font.name = 'Helvetica'
                    p_upper.font.size = Pt(41)
                    p_upper.font.bold = True
                    p_upper.font.color.rgb = RGBColor(255, 255, 255)
                    
                    lower_text = "By Media Research & Analytics Team"
                    lower_tb = first_slide.shapes.add_textbox(Inches(1.5), Inches(1.3), SLIDE_WIDTH - Inches(1.8), Inches(0.8))
                    tf_lower = lower_tb.text_frame
                    p_lower = tf_lower.add_paragraph()
                    p_lower.text = lower_text
                    p_lower.alignment = PP_ALIGN.LEFT
                    p_lower.font.name = 'Helvetica'
                    p_lower.font.size = Pt(41)
                    p_lower.font.bold = True
                    p_lower.font.color.rgb = RGBColor(255, 255, 255)
                    
                    if os.path.exists(DISCLAIMER_IMAGE):
                        disclaimer_left = Inches(1.2)
                        disclaimer_top = Inches(6.5)
                        disclaimer_width = SLIDE_WIDTH - Inches(1.2)
                        first_slide.shapes.add_picture(DISCLAIMER_IMAGE, disclaimer_left, disclaimer_top, width=disclaimer_width)
                
                # Add Content Slides
                for section_idx, section in enumerate(sections, 1):
                    title = section["title"]
                    lines = section["body"]
                    
                    is_competitor_section = "competitor" in title.lower()
                    max_slides = 999 if is_competitor_section else 5
                    
                    current_slide_idx = 1
                    current_y_offset = 0
                    
                    def add_new_slide(title_text, slide_number):
                        slide = prs.slides.add_slide(blank_layout)
                        bg = slide.background
                        bg.fill.solid()
                        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        
                        if os.path.exists(TEMPLATE_IMAGE):
                            slide.shapes.add_picture(TEMPLATE_IMAGE, LOGO_LEFT, SLIDE_HEIGHT - LOGO_SIZE - LOGO_BOTTOM_MARGIN, LOGO_SIZE)
                        
                        tb = slide.shapes.add_textbox(Inches(0.4), HEADER_TOP, SLIDE_WIDTH - Inches(0.8), HEADER_HEIGHT)
                        p = tb.text_frame.paragraphs[0]
                        p.text = title_text
                        p.font.name = 'Calibri (Headings)'
                        p.font.size = Pt(24)
                        p.font.color.rgb = ORANGE
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                        
                        return slide
                    
                    slide = add_new_slide(title, current_slide_idx)
                    content_box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Pt(MAX_CONTENT_HEIGHT_PT))
                    tf = content_box.text_frame
                    tf.word_wrap = True
                    tf.auto_size = False
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    
                    i = 0
                    while i < len(lines):
                        line = lines[i]
                        
                        if is_table_start(line):
                            table_lines = []
                            while i < len(lines) and is_table_start(lines[i]):
                                table_lines.append(lines[i])
                                i += 1
                            
                            table_data = parse_markdown_table(table_lines)
                            if not table_data:
                                continue
                            
                            rows, cols = len(table_data), len(table_data[0])
                            est_height_pt = len(table_data) * 24 + 40
                            
                            if current_y_offset + est_height_pt > MAX_CONTENT_HEIGHT_PT and current_slide_idx < max_slides:
                                current_slide_idx += 1
                                slide = add_new_slide(title, current_slide_idx)
                                current_y_offset = 0
                                content_box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Pt(MAX_CONTENT_HEIGHT_PT))
                                tf = content_box.text_frame
                                tf.word_wrap = True
                                tf.auto_size = False
                                tf.vertical_anchor = MSO_ANCHOR.TOP
                            
                            left = CONTENT_LEFT
                            top = CONTENT_TOP + Pt(current_y_offset)
                            tbl_shape = slide.shapes.add_table(rows, cols, left, top, CONTENT_WIDTH, Pt(est_height_pt))
                            tbl = tbl_shape.table
                            
                            for r_idx, row in enumerate(table_data):
                                for c_idx, cell_text in enumerate(row):
                                    cell = tbl.cell(r_idx, c_idx)
                                    cell.text = cell_text
                                    para = cell.text_frame.paragraphs[0]
                                    para.font.size = Pt(10)
                                    para.alignment = PP_ALIGN.LEFT
                                    cell.vertical_anchor = MSO_ANCHOR.TOP
                                    
                                    cell.fill.solid()
                                    if r_idx == 0:
                                        cell.fill.fore_color.rgb = ORANGE
                                        para.font.bold = True
                                        para.font.color.rgb = RGBColor(0, 0, 0)
                                    else:
                                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                                    
                                    cell.margin_left = Pt(5)
                                    cell.margin_right = Pt(5)
                                    cell.margin_top = Pt(5)
                                    cell.margin_bottom = Pt(5)
                            
                            current_y_offset += (tbl_shape.height.pt + 20)
                            continue
                        
                        is_head = is_topic_heading(line)
                        is_comp = is_competitor_marker(line)
                        
                        if is_comp:
                            clean_text = line.strip()[2:].strip()
                        else:
                            clean_text = remove_bold_markers(line).strip()
                        
                        if not clean_text:
                            i += 1
                            continue
                        
                        f_size = FONT_COMPETITOR if is_comp else (FONT_HEADING if is_head else FONT_NORMAL)
                        line_h = estimate_line_height(clean_text, f_size, is_head, is_comp)
                        
                        if is_comp and current_y_offset > 10:
                            current_slide_idx += 1
                            slide = add_new_slide(title, current_slide_idx)
                            current_y_offset = 0
                            content_box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Pt(MAX_CONTENT_HEIGHT_PT))
                            tf = content_box.text_frame
                            tf.word_wrap = True
                            tf.auto_size = False
                            tf.vertical_anchor = MSO_ANCHOR.TOP
                        elif current_y_offset + line_h > MAX_CONTENT_HEIGHT_PT and current_slide_idx < max_slides:
                            current_slide_idx += 1
                            slide = add_new_slide(title, current_slide_idx)
                            current_y_offset = 0
                            content_box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Pt(MAX_CONTENT_HEIGHT_PT))
                            tf = content_box.text_frame
                            tf.word_wrap = True
                            tf.auto_size = False
                            tf.vertical_anchor = MSO_ANCHOR.TOP
                        
                        p = tf.add_paragraph()
                        p.text = clean_text
                        p.word_wrap = True
                        p.font.name = 'Calibri'
                        p.font.size = Pt(f_size)
                        
                        if is_comp:
                            p.font.bold = True
                            p.font.color.rgb = ORANGE
                            p.space_before = Pt(20)  # Space BEFORE competitor header
                            p.space_after = Pt(0)    # No space after - goes directly to content
                        elif is_head:
                            p.font.bold = True
                            p.space_before = Pt(12)  # Space BEFORE topic/bucket header
                            p.space_after = Pt(0)    # No space after - goes directly to content
                        else:
                            # Normal points - no extra spacing
                            p.space_before = Pt(0)
                            p.space_after = Pt(0)
                        current_y_offset += line_h
                        i += 1
                
                # Thank You Slide
                thank_you_slide = prs.slides.add_slide(blank_layout)
                
                if os.path.exists(TEMPLATE_IMAGE):
                    thank_you_slide.shapes.add_picture(TEMPLATE_IMAGE, LOGO_LEFT, SLIDE_HEIGHT - LOGO_SIZE - LOGO_BOTTOM_MARGIN, LOGO_SIZE)
                
                thank_you_tb = thank_you_slide.shapes.add_textbox(Inches(0.3), Inches(2.0), SLIDE_WIDTH - Inches(1.0), Inches(2.0))
                tf_thank = thank_you_tb.text_frame
                tf_thank.word_wrap = True
                tf_thank.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_thank = tf_thank.add_paragraph()
                p_thank.text = "Thank You"
                p_thank.alignment = PP_ALIGN.CENTER
                p_thank.font.name = 'Helvetica'
                p_thank.font.size = Pt(33)
                p_thank.font.bold = True
                p_thank.font.color.rgb = ORANGE
                
                # Save to buffer
                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                
                # Store in session state for download
                st.session_state['ppt_file'] = ppt_buffer.getvalue()
                st.session_state['ppt_filename'] = output_filename
                st.sidebar.success("✅ PPT generated successfully!")
                
            except Exception as e:
                st.sidebar.error(f"Error generating PPT: {str(e)}")

# Show download link if PPT was generated
if 'ppt_file' in st.session_state and 'ppt_filename' in st.session_state:
    b64 = base64.b64encode(st.session_state['ppt_file']).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{st.session_state["ppt_filename"]}">Download {st.session_state["ppt_filename"]}</a>'
    st.sidebar.markdown(href, unsafe_allow_html=True)
# ========================== WORD DOC DOWNLOAD SECTION ==========================
if 'docx_bytes' in st.session_state:
    if st.sidebar.button("Generate Word Report"):
        with st.spinner("Generating Word document..."):
            try:
                import io
                from docx import Document as DocxDocument
                from docx.shared import Pt, RGBColor, Inches
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                from docx.oxml.ns import qn
                from docx.oxml import OxmlElement

                ORANGE_HEX = "FF8C00"

                def set_cell_bg(cell, hex_color):
                    """Set table cell background color."""
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    shd = OxmlElement('w:shd')
                    shd.set(qn('w:val'), 'clear')
                    shd.set(qn('w:color'), 'auto')
                    shd.set(qn('w:fill'), hex_color)
                    tcPr.append(shd)

                def add_paragraph(doc, text, bold=False, font_size=12,
                                  center=False, color_hex=None, space_before=0, space_after=6):
                    """Add a formatted paragraph to the document."""
                    p = doc.add_paragraph()
                    p.paragraph_format.space_before = Pt(space_before)
                    p.paragraph_format.space_after = Pt(space_after)
                    if center:
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = p.add_run(text)
                    run.bold = bold
                    run.font.size = Pt(font_size)
                    run.font.name = 'Calibri'
                    if color_hex:
                        run.font.color.rgb = RGBColor.from_string(color_hex)
                    return p

                # Re-read the uploaded file (reset seek position)
                doc_input = DocxDocument(io.BytesIO(st.session_state['docx_bytes']))
                full_text = "\n".join(para.text for para in doc_input.paragraphs if para.text.strip())

                # Parse sections (same logic as PPT)
                sections = []
                current_title = None
                current_body = []

                for line in full_text.splitlines():
                    stripped = line.strip()
                    if stripped.startswith("###"):
                        if current_title:
                            sections.append({"title": current_title, "body": current_body})
                        current_title = stripped[3:].strip()
                        current_body = []
                    elif current_title and stripped:
                        current_body.append(line)

                if current_title:
                    sections.append({"title": current_title, "body": current_body})

                # Extract client name (same logic as PPT)
                if sections:
                    first_title = sections[0]["title"].strip()
                    prefixes_to_remove = [
                        "Conversations on ",
                        "Topicwise Conversation on ",
                        "Exclusive Conversation on ",
                        "Month – on – Month Insights ",
                    ]
                    client_name = first_title
                    for prefix in prefixes_to_remove:
                        if client_name.startswith(prefix):
                            client_name = client_name[len(prefix):].strip()
                            break
                    if not client_name:
                        client_name = "Client"
                    word_filename = f"{client_name} Qualitative Insights Report.docx"
                else:
                    client_name = "Client"
                    word_filename = "Qualitative Insights Report.docx"

                # CREATE WORD DOCUMENT
                word_doc = DocxDocument()

                # Set page margins
                for section in word_doc.sections:
                    section.top_margin = Inches(1)
                    section.bottom_margin = Inches(1)
                    section.left_margin = Inches(1)
                    section.right_margin = Inches(1)

                # Cover / Title
                add_paragraph(
                    word_doc,
                    f"{client_name} Qualitative Insights Report",
                    bold=True, font_size=24, center=True,
                    color_hex=ORANGE_HEX, space_before=20, space_after=6
                )
                add_paragraph(
                    word_doc,
                    "By Media Research & Analytics Team",
                    bold=True, font_size=14, center=True, space_after=20
                )
                word_doc.add_page_break()

                # Process each section
                for section in sections:
                    title = section["title"]
                    lines = section["body"]

                    # Section title — center aligned, bold, orange
                    add_paragraph(
                        word_doc, title,
                        bold=True, font_size=18, center=True,
                        color_hex=ORANGE_HEX, space_before=12, space_after=8
                    )

                    i = 0
                    while i < len(lines):
                        line = lines[i]

                        # ── Table block ──────────────────────────────────────────
                        if is_table_start(line):
                            table_lines = []
                            while i < len(lines) and is_table_start(lines[i]):
                                table_lines.append(lines[i])
                                i += 1

                            table_data = parse_markdown_table(table_lines)
                            if not table_data:
                                continue

                            rows_count = len(table_data)
                            cols_count = len(table_data[0])

                            tbl = word_doc.add_table(rows=rows_count, cols=cols_count)
                            tbl.style = 'Table Grid'

                            # Evenly distribute column widths across available width
                            available_width = Inches(6.5)  # 8.5" page - 2" margins
                            col_width = available_width / cols_count

                            for r_idx, row_data in enumerate(table_data):
                                row = tbl.rows[r_idx]
                                for c_idx, cell_text in enumerate(row_data):
                                    cell = row.cells[c_idx]
                                    cell.width = col_width
                                    cell.text = cell_text

                                    # Style header row
                                    para = cell.paragraphs[0]
                                    run = para.add_run() if not para.runs else para.runs[0]
                                    run.font.size = Pt(10)
                                    run.font.name = 'Calibri'

                                    if r_idx == 0:
                                        run.bold = True
                                        set_cell_bg(cell, ORANGE_HEX)
                                    else:
                                        set_cell_bg(cell, "FFFFFF")

                            word_doc.add_paragraph()  # spacing after table
                            continue

                        # ── Competitor heading (## marker) ───────────────────────
                        if is_competitor_marker(line):
                            clean_text = line.strip()[2:].strip()
                            if clean_text:
                                add_paragraph(
                                    word_doc, clean_text,
                                    bold=True, font_size=16,
                                    color_hex=ORANGE_HEX, space_before=16, space_after=4
                                )
                            i += 1
                            continue

                        # ── Topic / bucket heading (**bold**) ────────────────────
                        if is_topic_heading(line):
                            clean_text = remove_bold_markers(line).strip()
                            if clean_text:
                                add_paragraph(
                                    word_doc, clean_text,
                                    bold=True, font_size=13, space_before=10, space_after=3
                                )
                            i += 1
                            continue

                        # ── Normal body line ─────────────────────────────────────
                        clean_text = remove_bold_markers(line).strip()
                        if clean_text:
                            add_paragraph(
                                word_doc, clean_text,
                                font_size=11, space_before=0, space_after=3
                            )
                        i += 1

                    word_doc.add_page_break()

                # Thank You page
                add_paragraph(
                    word_doc, "Thank You",
                    bold=True, font_size=28, center=True,
                    color_hex=ORANGE_HEX, space_before=40
                )

                # Save to buffer
                word_buffer = io.BytesIO()
                word_doc.save(word_buffer)
                word_buffer.seek(0)

                st.session_state['word_file'] = word_buffer.getvalue()
                st.session_state['word_filename'] = word_filename
                st.sidebar.success("✅ Word doc generated successfully!")

            except Exception as e:
                st.sidebar.error(f"Error generating Word doc: {str(e)}")

# Download button for Word doc
if 'word_file' in st.session_state:
    st.sidebar.download_button(
        label="⬇️ Download Word Report",
        data=st.session_state['word_file'],
        file_name=st.session_state['word_filename'],
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )

# File uploader for WordCloud with an on_change callback
st.sidebar.markdown(
    "<p style='font-size: 20px; font-weight: bold;'>WordCloud</p>",
    unsafe_allow_html=True
)
file = st.sidebar.file_uploader(
    "Upload Excel File for WordCloud", 
    type=["xlsx"], 
    key="wordcloud_uploader", 
    on_change=update_latest_wordcloud
)

# File uploader for Similarity News with an on_change callback
st.sidebar.markdown(
    "<p style='font-size: 20px; font-weight: bold;'>Similar News</p>",
    unsafe_allow_html=True
)
file1 = st.sidebar.file_uploader(
    "Upload Excel File for Similarity News", 
    type=["xlsx"], 
    key="similarity_uploader", 
    on_change=update_latest_similarity
)


    
# Process WordCloud only if the first file is uploaded and the second file is not uploaded
if file and not file1:
    st.sidebar.write("File Uploaded Successfully!")
    
    # Importing Dataset for WordCloud
    data = pd.read_excel(file)
    
    # Define a function to clean the text
    def clean(text):
        text = text.lower()
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        return text

    # Cleaning the text in the Headline column
    data['Cleaned_Headline'] = (
    data['Headline']
    .fillna('')                    # Replace NaN with empty string
    .astype(str)                   # Ensure all are strings
    .apply(clean)
)

    # Define a function to clean the text (more thorough)
    def cleaned(text):
        # Removes all special characters and numericals leaving the alphabets
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub(r'[[0-9]*]', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        text = re.sub('[\\n]', ' ', text)
        #text = re.sub(r'\b\w{1,3}\b', '', text)
        # removing apostrophes
        text = re.sub("'s", '', str(text))
        # removing hyphens
        text = re.sub("-", ' ', str(text))
        text = re.sub("— ", '', str(text))
        # removing quotation marks
        text = re.sub('\"', '', str(text))
        # removing any reference to outside text
        text = re.sub("[\(\[].*?[\)\]]", "", str(text))
        tokens = text.split()
        filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
        return " ".join(filtered_tokens)

    # Cleaning the text in the review column
    # Full version with str.cat()
    if {'Headline', 'Opening Text', 'Hit Sentence'}.issubset(data.columns):
        data['Text'] = (
            data['Headline']
            .fillna('')
            .astype(str)
            .str.cat(
                [data['Opening Text'].fillna('').astype(str),
                 data['Hit Sentence'].fillna('').astype(str)],
                sep=' '
            )
        )

    # Fallback version
    else:
        data['Text'] = (
            data['Headline']
            .fillna('')
            .astype(str)
            .str.cat(
                data.get('Summary', pd.Series('')).fillna('').astype(str),
                sep=' '
            )
        )

    data['Text'] = data['Text'].str.replace(r'\s+', ' ', regex=True).str.strip()
    data['Text'] = data['Text'].apply(cleaned)
    data.head()    

    # Define the 'entities' variable outside of the conditional blocks
    entities = list(data['Entity'].unique())

    # Define an empty 'wordclouds' dictionary
    wordclouds = {}
    st.sidebar.subheader("Word Cloud Parameters")
    st.sidebar.title("Word Clouds")
    wordcloud_entity = st.sidebar.selectbox("Select Entity for Word Cloud", entities)

    # Custom Stop Words Section
    st.sidebar.title("Custom Stop Words")
        
    custom_stopwords = st.sidebar.text_area("Enter custom stop words (comma-separated)", "")
    custom_stopwords = [word.strip() for word in custom_stopwords.split(',')]
    
    # Widget to adjust word cloud parameters
    wordcloud_size_height = st.slider("Select Word Cloud Size Height", 100, 1000, 400, step=50, key="wordcloud_height")
    wordcloud_size_width = st.slider("Select Word Cloud Size Width", 100, 1000, 400, step=50, key="wordcloud_width")
    wordcloud_max_words = st.slider("Select Max Number of Words", 10, 500, 50)
    
    if wordcloud_entity:
        st.header(f"Word Cloud for Entity: {wordcloud_entity}")
        # Generate Word Cloud with custom stop words removed
        cleaned_headlines = ' '.join(data[data['Entity'] == wordcloud_entity]['Text'])
    
        if custom_stopwords:
            for word in custom_stopwords:
                cleaned_headlines = cleaned_headlines.replace(word, '')
    
        wordcloud_image = WordCloud(
            background_color="white", 
            width=wordcloud_size_width, 
            height=wordcloud_size_height, 
            max_font_size=80, 
            max_words=wordcloud_max_words,
            colormap='Set1', 
            contour_color='black', 
            contour_width=2, 
            collocations=False
        ).generate(cleaned_headlines)
        
        # Create entity_data for the selected entity
        entity_data = data[data['Entity'] == wordcloud_entity]
    
        # Resize the word cloud image using PIL
        img = Image.fromarray(np.array(wordcloud_image))
        img = img.resize((wordcloud_size_width, wordcloud_size_height))
        
        # Add the entity to the wordclouds dictionary
        wordclouds[wordcloud_entity] = (wordcloud_image, entity_data)
    
        # Display the resized word cloud image in Streamlit
        st.image(img, caption=f"Word Cloud for Entity: {wordcloud_entity}")
    
    # Word Cloud Interaction
    if wordcloud_entity:
        # Get the selected entity's word cloud
        entity_wordcloud, entity_data = wordclouds.get(wordcloud_entity, (None, None))  # Use .get() to handle missing keys gracefully
        if entity_wordcloud is None:
            st.warning(f"No word cloud found for '{wordcloud_entity}'")
        else:
            words = list(entity_wordcloud.words_.keys())
    
        # Get the selected entity's word cloud
        entity_wordcloud, entity_data = wordclouds[wordcloud_entity]
        words = list(entity_wordcloud.words_.keys())
    
        word_frequencies = entity_wordcloud.words_
        words_f = list(word_frequencies.keys())
    
        # Create a list of tuples containing (word, frequency)
        word_frequency_list = [(word, frequency) for word, frequency in word_frequencies.items()]
    
        # Create a selectbox for the words in the word cloud
        selected_word = st.selectbox("Select a word from the word cloud", words)
    
        # Find rows where the selected word appears
        matching_rows = entity_data[entity_data['Headline'].str.contains(selected_word, case=False, na=False)]
    
        # Display the matching rows or a message if no matches are found
        if not matching_rows.empty:
            st.subheader(f"Matching Rows for '{selected_word}':")
            # Function to highlight the selected word
            def highlight_word(text, word):                
                return re.sub(f'\\b{word}\\b', f'**{word}**', text, flags=re.IGNORECASE)
    
            # Apply the highlight_word function to the Headline column
            matching_rows['Headline'] = matching_rows.apply(lambda row: highlight_word(row['Headline'], selected_word), axis=1)
    
            # Display the formatted dataframe
            st.dataframe(matching_rows)
        else:
            st.warning(f"No matching rows found for '{selected_word}'")
        st.write("**Click the below button**")
        if "show_interaction" not in st.session_state:
            st.session_state.show_interaction = False
        if st.button(f"View Word Cloud Interaction for Entity: {wordcloud_entity}"):
            st.session_state.show_interaction = not st.session_state.show_interaction
        if st.session_state.show_interaction:
            st.header(f"Word Cloud Interaction for Entity: {wordcloud_entity}")
            st.write("Entities in entities list:", entities)
            st.write("Keys in wordclouds dictionary:", list(wordclouds.keys()))
            st.write("Words and their frequencies:")
            st.write(word_frequency_list)
# Similar News section
elif file1 and not file:
    st.sidebar.write("File Uploaded Successfully!")
    
    # Importing Dataset for Similar News
    data = pd.read_excel(file1)
    if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
        data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
    if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
        data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
    if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
        data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
    
    # Define a function to clean the text
    def clean(text):
        text = text.lower()
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        return text

    # Cleaning the text in the Headline column
    data['Cleaned_Headline'] = (
    data['Headline']
    .fillna('')                    # Replace NaN with empty string
    .astype(str)                   # Ensure all are strings
    .apply(clean)
)
    # Full version with str.cat()
    if {'Headline', 'Opening Text', 'Hit Sentence'}.issubset(data.columns):
        data['Text'] = (
            data['Headline']
            .fillna('')
            .astype(str)
            .str.cat(
                [data['Opening Text'].fillna('').astype(str),
                 data['Hit Sentence'].fillna('').astype(str)],
                sep=' '
            )
        )

    # Fallback version
    else:
        data['Text'] = (
            data['Headline']
            .fillna('')
            .astype(str)
            .str.cat(
                data.get('Summary', pd.Series('')).fillna('').astype(str),
                sep=' '
            )
        )

    data['Text'] = data['Text'].str.replace(r'\s+', ' ', regex=True).str.strip()

    # Define a function to clean the text (more thorough)
    def cleaned(text):
        # Removes all special characters and numericals leaving the alphabets
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub(r'[[0-9]*]', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        text = re.sub('[\\n]', ' ', text)
        #text = re.sub(r'\b\w{1,3}\b', '', text)
        # removing apostrophes
        text = re.sub("'s", '', str(text))
        # removing hyphens
        text = re.sub("-", ' ', str(text))
        text = re.sub("— ", '', str(text))
        # removing quotation marks
        text = re.sub('\"', '', str(text))
        # removing any reference to outside text
        text = re.sub("[\(\[].*?[\)\]]", "", str(text))
        tokens = text.split()
        filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
        return " ".join(filtered_tokens)
        
    # Cleaning the text in the review column
    data['Text'] = data['Text'].apply(cleaned)
    data.head()    

    # Define the 'entities' variable for Similar News
    entities = list(data['Entity'].unique())

    # Define an empty 'wordclouds' dictionary (if needed)
    wordclouds = {}
    st.header("Similar News")
    st.sidebar.subheader("Similarity News Parameters")
    # Place your parameters for Similar News here

    selected_column = st.sidebar.selectbox(
        "Select column for Similar News Analysis", 
        options=["Headline", "Text"]
    )

    # Set column names based on the user's selection.
    if selected_column == "Headline":
        sim_column = "Cleaned_Headline"      # Use cleaned headlines for similarity computations
        classification_column = "Headline"     # Use the original headline for classification
        output_column = "Similar_Headline"
    else:
        sim_column = "Text"
        classification_column = "Text"
        output_column = "Similar_Text"
    
    # A slider for similarity percentage threshold
    sim_per = st.slider("Select Percentage for Similarity", 5, 100, 70)
    
    # Assume you have a list of unique entities
    entities = list(data['Entity'].unique())
    
    # Create an in-memory buffer for the Excel workbook
    output = io.BytesIO()
    updated_workbook = pd.ExcelWriter(output, engine='openpyxl')
    
    # Example dictionaries for keywords and topic mapping (adjust as needed)
    entity_keywords = {
        # 'Nothing Tech': ['nothing'],
        # 'Asian Paints': ['asian', 'keyword2', 'keyword3'],
    }
    
    topic_mapping = {
        'Merger': ['merger', 'merges'],
        'Acquire': ['acquire', 'acquisition', 'acquires'],
        'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
        'Partnership': ['IPO','ipo'],
        'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                 'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
        'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
        'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
        'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
        'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
        'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
        'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
        'Awards & Recognition': ['award', 'awards'],
        'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
        'Sale - Offers - Discounts': ['sale','offers','discount','discounts','discounted']
    }
    # Dictionary to store WordClouds (if needed)
    wordclouds = {}
    
    # Process each entity
    for entity in entities:
        # Filter data for the current entity
        entity_data = data[data['Entity'] == entity].copy()
        for val in entity_data[sim_column].unique():
            # Compute similarity using fuzz.ratio
            entity_data[val] = entity_data[sim_column].apply(lambda x: fuzz.ratio(x, val) >= sim_per)
            # Choose a group name. Here we use the lexicographically smallest value in the group.
            m = np.min(entity_data[entity_data[val] == True][sim_column])
            # Assign the group name to the new output column
            entity_data.loc[entity_data[sim_column] == val, output_column] = m
    
        # Sort the DataFrame by the new grouping column
        entity_data.sort_values(output_column, ascending=True, inplace=True)
        
        # Optionally, keep only columns up to (and including) the grouping column
        col_index = entity_data.columns.get_loc(output_column)
        entity_data = entity_data.iloc[:, :col_index + 1]
        
        # Optionally, drop the column immediately preceding the grouping column
        column_to_delete = entity_data.columns[entity_data.columns.get_loc(output_column) - 1]
        entity_data = entity_data.drop(column_to_delete, axis=1)
        
        # Define a function to classify news as "Exclusive" or "Not Exclusive"
        def classify_exclusivity(row):
            entity_name = entity_data['Entity'].iloc[0]
            if entity_name.lower() in row[classification_column].lower() or entity_name.lower() in row[output_column].lower():
                return "Exclusive"
            else:
                return "Not Exclusive"
        
        entity_data['Exclusivity'] = entity_data.apply(classify_exclusivity, axis=1)
        
        # Define a function to qualify the entity based on keyword matching
        def qualify_entity(row):
            entity_name = row['Entity']
            text = row[classification_column]
            if entity_name in entity_keywords:
                keywords = entity_keywords[entity_name]
                if any(keyword in text for keyword in keywords):
                    return "Qualified"
            return "Not Qualified"
        
        entity_data['Qualification'] = entity_data.apply(qualify_entity, axis=1)
        
        # Define a function to classify topics based on keywords
        def classify_topic(headline):
            lowercase_headline = headline.lower()
            for topic, words in topic_mapping.items():
                for word in words:
                    if word in lowercase_headline:
                        return topic
            return 'Other'
        
        entity_data['Topic'] = entity_data[classification_column].apply(classify_topic)
        
        # Insert a serial number column (optional)
        entity_data.insert(0, "sr no.", range(1, len(entity_data) + 1))
        
        # Write the processed entity data to a separate sheet in the workbook
        entity_data.to_excel(updated_workbook, sheet_name=entity, index=False, startrow=0)
        
        # Create a WordCloud for the entity based on the chosen analysis column
        wordcloud = WordCloud(width=550, height=400, background_color='white').generate(' '.join(entity_data[sim_column]))
        wordclouds[entity] = (wordcloud, entity_data)
    
    # Close the workbook and prepare the in-memory file for download
    updated_workbook.close()
    output.seek(0)
    excel_bytes = output.getvalue()

    st.markdown("### Download Grouped Data")
    st.markdown(
        f"Download the grouped data as an Excel file: [Similar_News_Grouped.xlsx](data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(excel_bytes).decode()})"
    )
    
    data.insert(0, "sr no.", range(1, len(data) + 1))
    data_csv = data.to_csv(index=False)
    st.markdown(
        f"Download the original data as a CSV file: [Original_Data.csv](data:text/csv;base64,{base64.b64encode(data_csv.encode()).decode()})"
    )
    
    grouped_data = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None)
    
    st.sidebar.subheader("Data Preview")
    entities_list = list(wordclouds.keys())
    selected_entities = st.sidebar.multiselect("Select Entities to Preview", entities_list)
    
    if selected_entities:
        for entity in selected_entities:
            st.header(f"Preview for Entity: {entity}")
            entity_data = grouped_data[entity]
            st.write(entity_data)

elif file and file1:
    latest = st.session_state.get("latest", None)
    if latest == "wordcloud":
        data = pd.read_excel(file)
        def clean(text):
            text = text.lower()
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            return text

        data['Cleaned_Headline'] = (
    data['Headline']
    .fillna('')                    # Replace NaN with empty string
    .astype(str)                   # Ensure all are strings
    .apply(clean)
)
    
        # Define a function to clean the text (more thorough)
        def cleaned(text):
            # Removes all special characters and numericals leaving the alphabets
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub(r'[[0-9]*]', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            text = re.sub('[\\n]', ' ', text)
            #text = re.sub(r'\b\w{1,3}\b', '', text)
            # removing apostrophes
            text = re.sub("'s", '', str(text))
            # removing hyphens
            text = re.sub("-", ' ', str(text))
            text = re.sub("— ", '', str(text))
            # removing quotation marks
            text = re.sub('\"', '', str(text))
            # removing any reference to outside text
            text = re.sub("[\(\[].*?[\)\]]", "", str(text))
            tokens = text.split()
            filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
            return " ".join(filtered_tokens)
    
        # Cleaning the text in the review column
        # Full version with str.cat()
        if {'Headline', 'Opening Text', 'Hit Sentence'}.issubset(data.columns):
            data['Text'] = (
                data['Headline']
                .fillna('')
                .astype(str)
                .str.cat(
                    [data['Opening Text'].fillna('').astype(str),
                     data['Hit Sentence'].fillna('').astype(str)],
                    sep=' '
                )
            )

        # Fallback version
        else:
            data['Text'] = (
                data['Headline']
                .fillna('')
                .astype(str)
                .str.cat(
                    data.get('Summary', pd.Series('')).fillna('').astype(str),
                    sep=' '
                )
            )

        data['Text'] = data['Text'].str.replace(r'\s+', ' ', regex=True).str.strip()
        data['Text'] = data['Text'].apply(cleaned)
        data.head()    
    
        # Define the 'entities' variable outside of the conditional blocks
        entities = list(data['Entity'].unique())
    
        # Define an empty 'wordclouds' dictionary
        wordclouds = {}
        st.sidebar.subheader("Word Cloud Parameters")
        st.sidebar.title("Word Clouds")
        wordcloud_entity = st.sidebar.selectbox("Select Entity for Word Cloud", entities)
    
        # Custom Stop Words Section
        st.sidebar.title("Custom Stop Words")
            
        custom_stopwords = st.sidebar.text_area("Enter custom stop words (comma-separated)", "")
        custom_stopwords = [word.strip() for word in custom_stopwords.split(',')]
        
        # Widget to adjust word cloud parameters
        wordcloud_size_height = st.slider("Select Word Cloud Size Height", 100, 1000, 400, step=50, key="wordcloud_height")
        wordcloud_size_width = st.slider("Select Word Cloud Size Width", 100, 1000, 400, step=50, key="wordcloud_width")
        wordcloud_max_words = st.slider("Select Max Number of Words", 10, 500, 50)
        
        if wordcloud_entity:
            st.header(f"Word Cloud for Entity: {wordcloud_entity}")
            # Generate Word Cloud with custom stop words removed
            cleaned_headlines = ' '.join(data[data['Entity'] == wordcloud_entity]['Text'])
        
            if custom_stopwords:
                for word in custom_stopwords:
                    cleaned_headlines = cleaned_headlines.replace(word, '')
        
            wordcloud_image = WordCloud(
                background_color="white", 
                width=wordcloud_size_width, 
                height=wordcloud_size_height, 
                max_font_size=80, 
                max_words=wordcloud_max_words,
                colormap='Set1', 
                contour_color='black', 
                contour_width=2, 
                collocations=False
            ).generate(cleaned_headlines)
            
            # Create entity_data for the selected entity
            entity_data = data[data['Entity'] == wordcloud_entity]
        
            # Resize the word cloud image using PIL
            img = Image.fromarray(np.array(wordcloud_image))
            img = img.resize((wordcloud_size_width, wordcloud_size_height))
            
            # Add the entity to the wordclouds dictionary
            wordclouds[wordcloud_entity] = (wordcloud_image, entity_data)
        
            # Display the resized word cloud image in Streamlit
            st.image(img, caption=f"Word Cloud for Entity: {wordcloud_entity}")
        
        # Word Cloud Interaction
        if wordcloud_entity:
            # Get the selected entity's word cloud
            entity_wordcloud, entity_data = wordclouds.get(wordcloud_entity, (None, None))  # Use .get() to handle missing keys gracefully
            if entity_wordcloud is None:
                st.warning(f"No word cloud found for '{wordcloud_entity}'")
            else:
                words = list(entity_wordcloud.words_.keys())
        
            # Get the selected entity's word cloud
            entity_wordcloud, entity_data = wordclouds[wordcloud_entity]
            words = list(entity_wordcloud.words_.keys())
        
            word_frequencies = entity_wordcloud.words_
            words_f = list(word_frequencies.keys())
        
            # Create a list of tuples containing (word, frequency)
            word_frequency_list = [(word, frequency) for word, frequency in word_frequencies.items()]
        
            # Create a selectbox for the words in the word cloud
            selected_word = st.selectbox("Select a word from the word cloud", words)
        
            # Find rows where the selected word appears
            matching_rows = entity_data[entity_data['Headline'].str.contains(selected_word, case=False, na=False)]
        
            # Display the matching rows or a message if no matches are found
            if not matching_rows.empty:
                st.subheader(f"Matching Rows for '{selected_word}':")
                # Function to highlight the selected word
                def highlight_word(text, word):                
                    return re.sub(f'\\b{word}\\b', f'**{word}**', text, flags=re.IGNORECASE)
        
                # Apply the highlight_word function to the Headline column
                matching_rows['Headline'] = matching_rows.apply(lambda row: highlight_word(row['Headline'], selected_word), axis=1)
        
                # Display the formatted dataframe
                st.dataframe(matching_rows)
            else:
                st.warning(f"No matching rows found for '{selected_word}'")
                
        st.write("**Click the below button**")
        if "show_interaction" not in st.session_state:
            st.session_state.show_interaction = False
        if st.button(f"View Word Cloud Interaction for Entity: {wordcloud_entity}"):
            st.session_state.show_interaction = not st.session_state.show_interaction
        if st.session_state.show_interaction:
            st.header(f"Word Cloud Interaction for Entity: {wordcloud_entity}")
            st.write("Entities in entities list:", entities)
            st.write("Keys in wordclouds dictionary:", list(wordclouds.keys()))
            st.write("Words and their frequencies:")
            st.write(word_frequency_list)
    elif latest == "similarity":
        data = pd.read_excel(file1)
        if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
        if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        
        # Define a function to clean the text
        def clean(text):
            text = text.lower()
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            return text
    
        # Cleaning the text in the Headline column
        data['Cleaned_Headline'] = (
    data['Headline']
    .fillna('')                    # Replace NaN with empty string
    .astype(str)                   # Ensure all are strings
    .apply(clean)
)
        # Full version with str.cat()
        if {'Headline', 'Opening Text', 'Hit Sentence'}.issubset(data.columns):
            data['Text'] = (
                data['Headline']
                .fillna('')
                .astype(str)
                .str.cat(
                    [data['Opening Text'].fillna('').astype(str),
                     data['Hit Sentence'].fillna('').astype(str)],
                    sep=' '
                )
            )

        # Fallback version
        else:
            data['Text'] = (
                data['Headline']
                .fillna('')
                .astype(str)
                .str.cat(
                    data.get('Summary', pd.Series('')).fillna('').astype(str),
                    sep=' '
                )
            )

        data['Text'] = data['Text'].str.replace(r'\s+', ' ', regex=True).str.strip()
    
        # Define a function to clean the text (more thorough)
        def cleaned(text):
            # Removes all special characters and numericals leaving the alphabets
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub(r'[[0-9]*]', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            text = re.sub('[\\n]', ' ', text)
            #text = re.sub(r'\b\w{1,3}\b', '', text)
            # removing apostrophes
            text = re.sub("'s", '', str(text))
            # removing hyphens
            text = re.sub("-", ' ', str(text))
            text = re.sub("— ", '', str(text))
            # removing quotation marks
            text = re.sub('\"', '', str(text))
            # removing any reference to outside text
            text = re.sub("[\(\[].*?[\)\]]", "", str(text))
            tokens = text.split()
            filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
            return " ".join(filtered_tokens)
            
        # Cleaning the text in the review column
        data['Text'] = data['Text'].apply(cleaned)
        data.head()    
    
        # Define the 'entities' variable for Similar News
        entities = list(data['Entity'].unique())
    
        # Define an empty 'wordclouds' dictionary (if needed)
        wordclouds = {}
        st.header("Similar News")
        st.sidebar.subheader("Similarity News Parameters")
        # Place your parameters for Similar News here
    
        selected_column = st.sidebar.selectbox(
            "Select column for Similar News Analysis", 
            options=["Headline", "Text"]
        )
    
        # Set column names based on the user's selection.
        if selected_column == "Headline":
            sim_column = "Cleaned_Headline"      # Use cleaned headlines for similarity computations
            classification_column = "Headline"     # Use the original headline for classification
            output_column = "Similar_Headline"
        else:
            sim_column = "Text"
            classification_column = "Text"
            output_column = "Similar_Text"
        
        # A slider for similarity percentage threshold
        sim_per = st.slider("Select Percentage for Similarity", 5, 100, 50)
        
        # Assume you have a list of unique entities
        entities = list(data['Entity'].unique())
        
        # Create an in-memory buffer for the Excel workbook
        output = io.BytesIO()
        updated_workbook = pd.ExcelWriter(output, engine='openpyxl')
        
        # Example dictionaries for keywords and topic mapping (adjust as needed)
        entity_keywords = {
            # 'Nothing Tech': ['nothing'],
            # 'Asian Paints': ['asian', 'keyword2', 'keyword3'],
        }
        
        topic_mapping = {
            'Merger': ['merger', 'merges'],
            'Acquire': ['acquire', 'acquisition', 'acquires'],
            'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
            'Partnership': ['IPO','ipo'],
            'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                     'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
            'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
            'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
            'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
            'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
            'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
            'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
            'Awards & Recognition': ['award', 'awards'],
            'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
            'Sale - Offers - Discounts': ['sale','offers','discount','discounts','discounted']
        }
        # Dictionary to store WordClouds (if needed)
        wordclouds = {}
        
        # Process each entity
        for entity in entities:
            # Filter data for the current entity
            entity_data = data[data['Entity'] == entity].copy()
            for val in entity_data[sim_column].unique():
                # Compute similarity using fuzz.ratio
                entity_data[val] = entity_data[sim_column].apply(lambda x: fuzz.ratio(x, val) >= sim_per)
                # Choose a group name. Here we use the lexicographically smallest value in the group.
                m = np.min(entity_data[entity_data[val] == True][sim_column])
                # Assign the group name to the new output column
                entity_data.loc[entity_data[sim_column] == val, output_column] = m
        
            # Sort the DataFrame by the new grouping column
            entity_data.sort_values(output_column, ascending=True, inplace=True)
            
            # Optionally, keep only columns up to (and including) the grouping column
            col_index = entity_data.columns.get_loc(output_column)
            entity_data = entity_data.iloc[:, :col_index + 1]
            
            # Optionally, drop the column immediately preceding the grouping column
            column_to_delete = entity_data.columns[entity_data.columns.get_loc(output_column) - 1]
            entity_data = entity_data.drop(column_to_delete, axis=1)
            
            # Define a function to classify news as "Exclusive" or "Not Exclusive"
            def classify_exclusivity(row):
                entity_name = entity_data['Entity'].iloc[0]
                if entity_name.lower() in row[classification_column].lower() or entity_name.lower() in row[output_column].lower():
                    return "Exclusive"
                else:
                    return "Not Exclusive"
            
            entity_data['Exclusivity'] = entity_data.apply(classify_exclusivity, axis=1)
            
            # Define a function to qualify the entity based on keyword matching
            def qualify_entity(row):
                entity_name = row['Entity']
                text = row[classification_column]
                if entity_name in entity_keywords:
                    keywords = entity_keywords[entity_name]
                    if any(keyword in text for keyword in keywords):
                        return "Qualified"
                return "Not Qualified"
            
            entity_data['Qualification'] = entity_data.apply(qualify_entity, axis=1)
            
            # Define a function to classify topics based on keywords
            def classify_topic(headline):
                lowercase_headline = headline.lower()
                for topic, words in topic_mapping.items():
                    for word in words:
                        if word in lowercase_headline:
                            return topic
                return 'Other'
            
            entity_data['Topic'] = entity_data[classification_column].apply(classify_topic)
            
            # Insert a serial number column (optional)
            entity_data.insert(0, "sr no.", range(1, len(entity_data) + 1))
            
            # Write the processed entity data to a separate sheet in the workbook
            entity_data.to_excel(updated_workbook, sheet_name=entity, index=False, startrow=0)
            
            # Create a WordCloud for the entity based on the chosen analysis column
            wordcloud = WordCloud(width=550, height=400, background_color='white').generate(' '.join(entity_data[sim_column]))
            wordclouds[entity] = (wordcloud, entity_data)
        
        # Close the workbook and prepare the in-memory file for download
        updated_workbook.close()
        output.seek(0)
        excel_bytes = output.getvalue()
    
        st.markdown("### Download Grouped Data")
        st.markdown(
            f"Download the grouped data as an Excel file: [Similar_News_Grouped.xlsx](data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(excel_bytes).decode()})"
        )
        
        data.insert(0, "sr no.", range(1, len(data) + 1))
        data_csv = data.to_csv(index=False)
        st.markdown(
            f"Download the original data as a CSV file: [Original_Data.csv](data:text/csv;base64,{base64.b64encode(data_csv.encode()).decode()})"
        )
        
        grouped_data = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None)
        
        st.sidebar.subheader("Data Preview")
        entities_list = list(wordclouds.keys())
        selected_entities = st.sidebar.multiselect("Select Entities to Preview", entities_list)
        
        if selected_entities:
            for entity in selected_entities:
                st.header(f"Preview for Entity: {entity}")
                entity_data = grouped_data[entity]
                st.write(entity_data)
