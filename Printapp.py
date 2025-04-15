import streamlit as st
import pandas as pd
import base64
import io
import numpy as np
import re
from PIL import Image
import matplotlib.pyplot as plt
# import spacy
import logging
import warnings
from nltk.corpus import stopwords
import nltk
import os
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Border, Side, Alignment, Font,PatternFill
from openpyxl.utils.dataframe import dataframe_to_rows # Add these imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches
from io import BytesIO
from matplotlib.ticker import PercentFormatter
# from wordcloud import WordCloud
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator


# Streamlit app with a sidebar layout
st.set_page_config(layout="wide")

# Function to process the Excel file
def process_excel(file):
    # Initialize Excel writer
    output = BytesIO()
    excel_writer = pd.ExcelWriter(output, engine='xlsxwriter')
    all_dframes = []

#     # Iterate through each sheet in the uploaded file
#     for sheet_name in pd.ExcelFile(file).sheet_names:
#         data = pd.read_excel(file, sheet_name=sheet_name)

#         # Convert 'unnamed 2' column to numeric and sort by 'unnamed 0' and 'unnamed 2'
#         data['unnamed 2'] = pd.to_numeric(data['unnamed 2'], errors='coerce')
#         sorted_data = data.sort_values(by=['unnamed 0', 'unnamed 2'], kind='mergesort')
#         sorted_data.drop("unnamed 2", axis=1, inplace=True)
#         sorted_data['Source'] = ""

#         # Process different subsets of data
#         df1 = sorted_data[sorted_data['unnamed 0'] == 'c'].drop(columns=["unnamed 0"] + sorted_data.columns[2:].tolist())
#         df2 = sorted_data[sorted_data['unnamed 0'] == 'd'].drop(columns=["unnamed 0"] + sorted_data.columns[2:].tolist())
#         df3 = sorted_data[sorted_data['unnamed 0'] == 'b'].drop(columns=sorted_data.columns[:2].tolist() + ['Source', 'unnamed 4'])

#         # Reset indexes
#         df1.reset_index(drop=True, inplace=True)
#         df2.reset_index(drop=True, inplace=True)
#         df3.reset_index(drop=True, inplace=True)

#         # Combine dataframes
#         result_1 = pd.concat([df3, df2, df1], axis=1, join='outer')
#         result_1.rename({'unnamed 3': 'Headline', 'unnamed 1': 'Summary'}, axis=1, inplace=True)

#         # Replace the column names
#         s = result_1.columns.to_series()
#         s.iloc[2] = 'Source'
#         result_1.columns = s

#         # Split 'Source' column
#         split_data = result_1['Source'].str.split(',', expand=True)
#         dframe = pd.concat([result_1, split_data], axis=1)
#         dframe.drop('Source', axis=1, inplace=True)
#         dframe.rename({0: 'Source', 1: 'Date', 2: 'Words', 3: 'Journalists'}, axis=1, inplace=True)
#         dframe['Headline'] = dframe['Headline'].str.replace("Factiva Licensed Content", "").str.strip()

#         # Add 'Entity' column
#         dframe.insert(dframe.columns.get_loc('Headline'), 'Entity', sheet_name)

#         # Replace specific words in 'Journalists' column with 'Bureau News'
#         words_to_replace = ['Hans News Service', 'IANS', 'DH Web Desk', 'HT Entertainment Desk', 'Livemint', 
#                             'Business Reporter', 'HT Brand Studio', 'Outlook Entertainment Desk', 'Outlook Sports Desk',
#                             'DHNS', 'Express News Service', 'TIMES NEWS NETWORK', 'Staff Reporter', 'Affiliate Desk', 
#                             'Best Buy', 'FE Bureau', 'HT News Desk', 'Mint SnapView', 'Our Bureau', 'TOI Sports Desk',
#                             'express news service', '(English)', 'HT Correspondent', 'DC Correspondent', 'TOI Business Desk',
#                             'India Today Bureau', 'HT Education Desk', 'PNS', 'Our Editorial', 'Sports Reporter',
#                             'TOI News Desk', 'Legal Correspondent', 'The Quint', 'District Correspondent', 'etpanache',
#                             'ens economic bureau', 'Team Herald', 'Equitymaster']
#         dframe['Journalists'] = dframe['Journalists'].replace(words_to_replace, 'Bureau News', regex=True)
        
#         additional_replacements = ['@timesgroup.com', 'TNN']
#         dframe['Journalists'] = dframe['Journalists'].replace(additional_replacements, '', regex=True)

#         # Fill NaN or spaces in 'Journalists' column
#         dframe['Journalists'] = dframe['Journalists'].apply(lambda x: 'Bureau News' if pd.isna(x) or x.isspace() else x)
#         dframe['Journalists'] = dframe['Journalists'].str.lstrip()

#         # Read additional data for merging
#         data2 = pd.read_excel(r"FActiva Publications.xlsx")
        
#         # Merge the current dataframe with additional data
#         merged = pd.merge(dframe, data2, how='left', left_on=['Source'], right_on=['Source'])

#         # Save the merged data to Excel with the sheet name
#         merged.to_excel(excel_writer, sheet_name=sheet_name, index=False)
        
#         # Append DataFrame to the list
#         all_dframes.append(merged)
    
#     # Combine all DataFrames into a single DataFrame
#     combined_data = pd.concat(all_dframes, ignore_index=True)

#     # Add a serial number column
#     combined_data['sr no'] = combined_data.reset_index().index + 1

#     # Rearrange columns to have 'sr no' before 'Entity'
#     combined_data = combined_data[['sr no', 'Entity'] + [col for col in combined_data.columns if col not in ['sr no', 'Entity']]]

#     # Save the combined data to a new sheet
#     combined_data.to_excel(excel_writer, sheet_name='Combined_All_Sheets', index=False)
    
#     # Show the processed dataframe in the web app
#     st.write(combined_data)


#     # Save and return the Excel file
#     excel_writer.close()
#     output.seek(0)
#     return output
    
    
# # Streamlit app setup
# st.title("Print Excel File Processor & Merger")

# Upload file
# uploaded_file = st.file_uploader("Upload an Excel file", type=["xlsx"])

# # Process the file if uploaded
# if uploaded_file is not None:
#     processed_file = process_excel(uploaded_file)
    
#     # Download button
#     st.download_button(
#         label="Download Processed Excel",
#         data=processed_file,
#         file_name="Processed_Excel.xlsx",

#     )

# Function to extract entity name from file path
def extract_entity_name(file_path):
    base_name = os.path.basename(file_path)
    entity_name = base_name.split('_or_')[0].replace("_", " ").split('-')[0].strip()
    return entity_name

# # Web app title
# st.title('Online Excel File Merger & Entity Extractor')

# # File uploader
# uploaded_files = st.file_uploader("Upload your Excel files", accept_multiple_files=True, type=['xlsx'])

# if uploaded_files:
#     final_df = pd.DataFrame()
    
#     # Loop through each uploaded file
#     for uploaded_file in uploaded_files:
#         df = pd.read_excel(uploaded_file)
        
#         # Extract the entity name and add it as a new column
#         entity_name = extract_entity_name(uploaded_file.name)
#         df['Entity'] = entity_name
        
#         # Concatenate all the dataframes
#         final_df = pd.concat([final_df, df], ignore_index=True)
    
#     # Process columns as required
#     existing_columns = final_df.columns.tolist()
#     influencer_index = existing_columns.index('Influencer')
#     country_index = existing_columns.index('Country')
    
#     new_order = (
#         existing_columns[:influencer_index + 1] +  # All columns up to and including 'Influencer'
#         ['Entity', 'Reach', 'Sentiment', 'Keywords', 'State', 'City', 'Engagement'] +  # Adding new columns
#         existing_columns[influencer_index + 1:country_index + 1]  # All columns between 'Influencer' and 'Country'
#     )
    
    
#     # Fill missing values in 'Influencer' column with 'Bureau News'
#     final_df['Influencer'] = final_df['Influencer'].fillna('Bureau News')
#     final_df['Date'] = pd.to_datetime(final_df['Date']).dt.normalize()
    
#     # Reorder the DataFrame
#     final_df = final_df[new_order]
    
#     # Show the processed dataframe in the web app
#     st.write(final_df)
    
#     # Prepare Excel file in memory
#     output = BytesIO()
#     with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
#         final_df.to_excel(writer, index=False)
    
#     # Convert buffer to bytes
#     processed_data = output.getvalue()

#     # Option to download the merged file
#     st.download_button(
#         label="Download Merged Excel",
#         data=processed_data,
#         file_name='merged_excel_with_entity.xlsx',
#         mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
#     )
# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Data preprocessing function (You can include your data preprocessing here)

# Function to create separate Excel sheets by Entity
def create_entity_sheets(data, writer):
    # Define a format with text wrap
    wrap_format = writer.book.add_format({'text_wrap': True})

    for Entity in data['Entity'].unique():
        entity_df = data[data['Entity'] == Entity]
        entity_df.to_excel(writer, sheet_name=Entity, index=False)
        worksheet = writer.sheets[Entity]
        worksheet.set_column(1, 4, 48, cell_format=wrap_format)
        # Calculate column widths based on the maximum content length in each column except columns 1 to 4
        max_col_widths = [
            max(len(str(value)) for value in entity_df[column])
            for column in entity_df.columns[5:]  # Exclude columns 1 to 4
        ]

        # Set the column widths dynamically for columns 5 onwards
        for col_num, max_width in enumerate(max_col_widths):
            worksheet.set_column(col_num + 5, col_num + 5, max_width + 2)  # Adding extra padding for readability       


def add_entity_info(ws, entity_info, start_row):
    for i, line in enumerate(entity_info.split('\n'), start=1):
        cell = ws.cell(row=start_row + i - 1, column=1)
        cell.value = line
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))
#         cell.alignment = Alignment(horizontal='center')  # Merge and center for all lines
#         ws.merge_cells(start_row=start_row + i - 1, start_column=1, end_row=start_row + i, end_column=5)
        
        # Apply specific formatting for Source line
        if line.startswith('Source:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=3, start_column=1, end_row=3, end_column=5)
            cell.font = Font(color="000000",name="Gill Sans")
            
        # Apply specific formatting for Source line
        if line.startswith('Entity:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans", bold=True )
            cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
            
        # Apply specific formatting for Source line
        if line.startswith('Time Period of analysis:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans")
            
        # Apply specific formatting for Source line
        if line.startswith('News search:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=4, start_column=1, end_row=4, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans")
            

def add_styling_to_worksheet(ws, df, start_row, comment):
    # Apply table heading as comment
    cell = ws.cell(row=start_row, column=1)
    cell.value = comment
    cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
    cell.font = Font(color="000000", bold=True, name="Gill Sans")
    cell.alignment = Alignment(horizontal='center')
    ws.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=len(df.columns))
    
    # Increment the start row
    start_row += 1

    # Apply styling to column headers
    for col_idx, col_name in enumerate(df.columns, start=1):
        cell = ws.cell(row=start_row, column=col_idx)
        cell.value = col_name
        cell.font = Font(color="000000", bold=True ,name="Gill Sans")
        cell.alignment = Alignment(horizontal='center')
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))  
        
    start_row += 1

    # Write DataFrame values with styling
    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=False), start=start_row):
        for c_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx)
            if isinstance(value, pd.Period):
                cell.value = value.strftime('%Y-%m') 
            else:
                cell.value = value
            cell.font = Font(name="Gill Sans")    
            cell.alignment = Alignment(horizontal='center')
    
    # Apply borders to all cells
    for row in ws.iter_rows(min_row=start_row, max_row=start_row+len(df), min_col=1, max_col=len(df.columns)):
        for cell in row:
            cell.border = Border(left=Side(border_style="thin", color="000000"),
                                 right=Side(border_style="thin", color="000000"),
                                 top=Side(border_style="thin", color="000000"),
                                 bottom=Side(border_style="thin", color="000000"))
            
def multiple_dfs(df_list, sheet_name, file_name, comments, entity_info):
    wb = Workbook()
    ws = wb.active
    current_row = 1
    
    # Add entity information to the first 4 rows
    add_entity_info(ws, entity_info, current_row)
    current_row += 6
    
    for df, comment in zip(df_list, comments):
        add_styling_to_worksheet(ws, df, current_row, comment)
        current_row += len(df) + 4
    
    wb.save(file_name)


def add_table_to_slide(slide, df, title, textbox_text):
    rows, cols = df.shape
    left = Inches(0.8)
    top = Inches(2.8)
    width = Inches(14)
    max_table_height = Inches(5)
    total_height_needed = Inches(0.8 * (rows + 1))
    height = max_table_height if total_height_needed > max_table_height else total_height_needed

    # Add title shape (above the table)
    title_shape = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.2))
    title_frame = title_shape.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add the table
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = df.columns[i]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Gill Sans'
                run.font.size = Pt(15)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        # # Add borders to header cell
        # for border in ["top", "bottom", "left", "right"]:
        #     cell.border.__getattribute__(border).color.rgb = RGBColor(0, 0, 0)
        #     cell.border.__getattribute__(border).width = Pt(1)
            # # Add borders to data cell
            # for border in ["top", "bottom", "left", "right"]:
            #     cell.border.__getattribute__(border).color.rgb = RGBColor(0, 0, 0)
            #     cell.border.__getattribute__(border).width = Pt(0.75)

    

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.values[i, j])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Gill Sans'
                    run.font.size = Pt(15)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Add a text box above the table (shared across all DataFrame slides)
    textbox_left = Inches(0.25)  # Adjust left positioning as needed
    textbox_right = Inches(0.25)
    textbox_top = Inches(0.8)  # Adjust top positioning as needed
    textbox_width = Inches(15.5)  # Adjust width
    textbox_height = Inches(2.1)  # Adjust height

    text_box = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = textbox_text  # The custom text box content for each slide
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(17)  # Adjust the font size as needed
#             run.font.bold = True
            run.font.name = 'Gill Sans'
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text

    # Add the image (footer logo) at the bottom of the slide
    left = Inches(0.0)
    top = prs.slide_height - Inches(1)
    slide.shapes.add_picture( img_path,left, top, height=Inches(1))  # Adjust as needed




# # Function to save multiple DataFrames in a single Excel sheet
# def multiple_dfs(df_list, sheets, file_name, spaces, comments):
#     writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
#     row = 2
#     for dataframe, comment in zip(df_list, comments):
#         pd.Series(comment).to_excel(writer, sheet_name=sheets, startrow=row,
#                                     startcol=1, index=False, header=False)
#         dataframe.to_excel(writer, sheet_name=sheets, startrow=row + 1, startcol=0)
#         row = row + len(dataframe.index) + spaces + 2
#     writer.close()

# Function to add image to slide
def add_image_to_slide(slide, img_path3):
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    slide.shapes.add_picture(img_path3, left, top, width=width)

# Generate an image from the bar chart
def generate_bar_chart(df):
    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.bar(
        df["Entity"], 
        df["News Count"], 
        color="skyblue", 
        edgecolor="black"
    )
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width() / 2, height, f"{height}", ha="center", va="bottom", fontsize=10)
    ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Entity", fontsize=12)
    ax.set_ylabel("News Count", fontsize=12)
    ax.grid(axis="y", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path3 = "bar_chart.png"
    fig.savefig(img_path3, dpi=300)
    plt.close(fig)
    return img_path3

def generate_bar_chart1(df):
    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.bar(
        df["Entity"], 
        df["News Count"], 
        color="skyblue", 
        edgecolor="black"
    )
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width() / 2, height, f"{height}", ha="center", va="bottom", fontsize=10)
    ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Entity", fontsize=12)
    ax.set_ylabel("News Count", fontsize=12)
    ax.grid(axis="y", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path3 = "bar_chart.png"
    fig.savefig(img_path3, dpi=300)
    plt.close(fig)
    return img_path3



def generate_line_graph(df):
    fig, ax = plt.subplots(figsize=(15, 5.6))
    
    # Exclude the 'Total' column and row for the graph
    filtered_df = df.loc[df['Date'] != 'Total'].copy()
    filtered_df = filtered_df.drop(columns=['Total'], errors='ignore')

    for entity in filtered_df.columns[1:]:  # Exclude the first column (Date)
        ax.plot(filtered_df['Date'].astype(str), filtered_df[entity], marker='o', label=entity)
        for x, y in zip(df['Date'].astype(str), df[entity]):
            ax.text(x, y, str(y), fontsize=10, ha='right', va='bottom')

    
    # ax.set_title("Month-on-Month Trends", fontsize=14)
    ax.set_xlabel("Month", fontsize=12)
    ax.set_ylabel("News Count", fontsize=12)
    ax.legend(title="Entities", fontsize=10)
    ax.grid(axis='y', linestyle='--', alpha=0.7)
    plt.xticks(rotation=45)

    # Save plot as image
    img_path6 = "line_graph.png"
    fig.savefig(img_path6, dpi=300)
    plt.close(fig)
    return img_path6


# Function to add image to slide
def add_image_to_slide(slide, img_path6):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(15)  # Specify exact width
    height = Inches(5.6)  # Specify exact height
    slide.shapes.add_picture(img_path6, left, top, width=width, height=height)

def generate_horizontal_bar_chart(df):
    # Remove 'Client-' prefix from 'Entity' column
    # df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
    
    # Filter out unwanted rows
    # df = df[df["City"] != "GrandTotal"]
    
    # Create the bar chart
    fig, ax = plt.subplots(figsize=(12, 6))  # Increase figure width for better label visibility
    x = range(len(df["Publication Name"]))  # Define x positions for the bars
    bars = ax.bar(
        x, 
        df["Industry"], 
        color="skyblue", 
        edgecolor="black"
    )
    
    # Add data labels on top of the bars without decimal
    for bar in bars:
        height = int(bar.get_height())  # Convert height to integer
        ax.text(
            bar.get_x() + bar.get_width() / 2, 
            height, 
            f"{height}", 
            ha="center", 
            va="bottom", 
            fontsize=12,
            fontweight="bold"
        )
    
    # Set chart title and axis labels
    # ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Publication Name", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")
    
    # Customize x-axis ticks and labels for better visibility
    ax.set_xticks(x)
    ax.set_xticklabels(df["Publication Name"], rotation=45, ha="right", fontsize=12,fontweight="bold")

    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black", which="major", width=1, labelrotation=0)
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")
    
    # Add gridlines for better readability
    ax.grid(axis="y", linestyle="--", alpha=0.7)

# def generate_horizontal_bar_chart(df):
#     df_sorted = df.sort_values(by="Industry", ascending=False)
#     fig, ax = plt.subplots(figsize=(10, 6))
#     bars = ax.barh(
#         df_sorted["Publication Name"], 
#         df_sorted["Industry"], 
#         color="skyblue", 
#         edgecolor="black"
#     )
    
#     for bar in bars:
#         width = bar.get_width()
#         ax.text(
#             width, 
#             bar.get_y() + bar.get_height() / 2, 
#             f"{width}", 
#             ha="left", 
#             va="center", 
#             fontsize=10
#         )
    
#     # ax.set_title("Publication Name", fontsize=14)
#     ax.set_xlabel("News Count", fontsize=12)
#     ax.set_ylabel("Publication Name", fontsize=12)
#     ax.grid(axis="x", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path7 = "horizontal_bar_chart.png"
    fig.savefig(img_path7, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path7

# Function to add image to slide
def add_image_to_slide1(slide, img_path7):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path7, left, top, width=width, height=height)


def generate_grouped_bar_chartbt(df, img_name="grouped_bar_chart.png"):
    """
    Generates a grouped bar chart for Publication Name, showing % of articles by Bureaus and Journalists.

    Parameters:
        df (pd.DataFrame): A DataFrame containing:
                          - 'Publication Name': Categories for the x-axis.
                          - '% of articles by Bureaus': Percentage of articles by Bureaus (as strings with '%').
                          - '% of articles by Journalists': Percentage of articles by Journalists (as strings with '%').
        img_name (str): Name of the image file to save. Default is 'grouped_bar_chart.png'.

    Returns:
        str: Path to the saved grouped bar chart image.
    """
    # Convert percentage strings to numeric values
    df["% of articles by Bureaus"] = df["% of articles by Bureaus"].replace({"%": ""}, regex=True).astype(float)
    df["% of articles by Journalists"] = df["% of articles by Journalists"].replace({"%": ""}, regex=True).astype(float)
    
    # Extract data
    publication_names = df["Publication Name"]
    bureaus = df["% of articles by Bureaus"]
    journalists = df["% of articles by Journalists"]
    
    # Define positions on the x-axis
    x = np.arange(len(publication_names))  # Numeric positions for each category
    bar_width = 0.4  # Width of each bar
    
    # Create the plot
    fig, ax = plt.subplots(figsize=(12, 6))
    bars1 = ax.bar(x - bar_width/2, bureaus, bar_width, label="% of articles by Bureaus", color="skyblue", edgecolor="black")
    bars2 = ax.bar(x + bar_width/2, journalists, bar_width, label="% of articles by Journalists", color="orange", edgecolor="black")
    
    # Add data labels
    for bars in (bars1, bars2):
        for bar in bars:
            height = bar.get_height()
            ax.text(
                bar.get_x() + bar.get_width() / 2, 
                height + 1,  # Add a slight offset above the bar
                f"{int(height)}%",  # Correctly format the percentage display
                ha="center", 
                va="bottom", 
                fontsize=10
            )

    # Chart aesthetics
    ax.set_xlabel("Publication Name", fontsize=12)
    ax.set_ylabel("Percentage", fontsize=12)
    # ax.set_title("Percentage of Articles by Bureaus and Journalists", fontsize=14)
    ax.set_xticks(x)
    ax.set_xticklabels(publication_names, rotation=45, ha="right")
    ax.legend(title="Source", fontsize=10)
    ax.grid(axis="y", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path8 = "horizontal_bar_chart.png"
    fig.savefig(img_path8, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path8

    # Function to add image to slide
def add_image_to_slide2(slide, img_path8):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path8, left, top, width=width, height=height)

def generate_horizontal_bar_chartj(df):
    # Remove 'Client-' prefix from 'Entity' column
    # df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
    
    # Filter out unwanted rows
    # df = df[df["City"] != "GrandTotal"]
    
    # Create the bar chart
    fig, ax = plt.subplots(figsize=(12, 6))  # Increase figure width for better label visibility
    x = range(len(df["Journalist"]))  # Define x positions for the bars
    bars = ax.bar(
        x, 
        df["Industry"], 
        color="skyblue", 
        edgecolor="black"
    )
    
    # Add data labels on top of the bars without decimal
    for bar in bars:
        height = int(bar.get_height())  # Convert height to integer
        ax.text(
            bar.get_x() + bar.get_width() / 2, 
            height, 
            f"{height}", 
            ha="center", 
            va="bottom", 
            fontsize=12,
            fontweight="bold"
        )
    
    # Set chart title and axis labels
    # ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Journalist", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")
    
    # Customize x-axis ticks and labels for better visibility
    ax.set_xticks(x)
    ax.set_xticklabels(df["Journalist"], rotation=45, ha="right", fontsize=12,fontweight="bold")

    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black", which="major", width=1, labelrotation=0)
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")
    
    # Add gridlines for better readability
    ax.grid(axis="y", linestyle="--", alpha=0.7)



# def generate_horizontal_bar_chartj(df):
#     df["Industry"] = pd.to_numeric(df["Industry"], errors="coerce")
#     df_sorted = df.sort_values(by="Industry", ascending=False)
#     fig, ax = plt.subplots(figsize=(10, 6))
#     bars = ax.barh(
#         df_sorted["Journalist"], 
#         df_sorted["Industry"], 
#         color="skyblue", 
#         edgecolor="black"
#     )
    
#     for bar in bars:
#         width = bar.get_width()
#         ax.text(
#             width, 
#             bar.get_y() + bar.get_height() / 2, 
#             f"{width}", 
#             ha="left", 
#             va="center", 
#             fontsize=10
#         )
    
#     # ax.set_title("Publication Name", fontsize=14)
#     ax.set_xlabel("News Count", fontsize=12)
#     ax.set_ylabel("Journalist", fontsize=12)
#     ax.grid(axis="x", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path9 = "horizontal_bar_chart.png"
    fig.savefig(img_path9, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path9

# Function to add image to slide
def add_image_to_slide3(slide, img_path9):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path9, left, top, width=width, height=height)


def generate_horizontal_bar_chartpt(df):
    # Remove 'Client-' prefix from 'Entity' column
    # df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
    
    # Filter out unwanted rows
    df = df[df["Publication Type"] != "GrandTotal"]
    
    # Create the bar chart
    fig, ax = plt.subplots(figsize=(12, 6))  # Increase figure width for better label visibility
    x = range(len(df["Publication Type"]))  # Define x positions for the bars
    bars = ax.bar(
        x, 
        df["Industry"], 
        color="skyblue", 
        edgecolor="black"
    )
    
    # Add data labels on top of the bars without decimal
    for bar in bars:
        height = int(bar.get_height())  # Convert height to integer
        ax.text(
            bar.get_x() + bar.get_width() / 2, 
            height, 
            f"{height}", 
            ha="center", 
            va="bottom", 
            fontsize=12,
            fontweight="bold"
        )
    
    # Set chart title and axis labels
    # ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Publication Type", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")
    
    # Customize x-axis ticks and labels for better visibility
    ax.set_xticks(x)
    ax.set_xticklabels(df["Publication Type"], rotation=45, ha="right", fontsize=12,fontweight="bold")

    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black", which="major", width=1, labelrotation=0)
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")
    
    # Add gridlines for better readability
    ax.grid(axis="y", linestyle="--", alpha=0.7)

# def generate_horizontal_bar_chartpt(df):
#     """
#     Generates a horizontal bar chart and sorts the bars after plotting by their height (values).
#     """
#     # Ensure 'Industry' column is numeric (coerce errors to NaN)
#     df["Industry"] = pd.to_numeric(df["Industry"], errors="coerce")

#     # Remove rows where 'Industry' is NaN or 'GrandTotal'
#     df_filtered = df[df["Publication Type"] != "GrandTotal"]

#     # Create a horizontal bar chart
#     fig, ax = plt.subplots(figsize=(10, 6))
#     bars = ax.barh(
#         df_filtered["Publication Type"], 
#         df_filtered["Industry"], 
#         color="skyblue", 
#         edgecolor="black"
#     )

#     # Add data labels on the bars
#     for bar in bars:
#         width = bar.get_width()
#         ax.text(
#             width, 
#             bar.get_y() + bar.get_height() / 2, 
#             f"{width}", 
#             ha="left", 
#             va="center", 
#             fontsize=10
#         )

#     # Sort the bars by width (Industry values)
#     sorted_bars = sorted(bars, key=lambda bar: bar.get_width(), reverse=True)

#     # Remove original bars
#     for bar in bars:
#         bar.remove()

#     # Re-add the bars in sorted order
#     for bar in sorted_bars:
#         ax.barh(
#             bar.get_y(),
#             bar.get_width(),
#             color="skyblue", 
#             edgecolor="black"
#         )

#     # Adjust axis labels and grid
#     ax.set_xlabel("News Count", fontsize=12)
#     ax.set_ylabel("Publication Type", fontsize=12)
#     ax.grid(axis="x", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path17 = "horizontal_bar_chart.png"
    fig.savefig(img_path17, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path17

# Function to add image to slide
def add_image_to_slide4(slide, img_path17):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path17, left, top, width=width, height=height)


# def generate_horizontal_bar_chartst(df):
#     df_filtered = df[df["Sentiment"] != "GrandTotal"]
    
#     # Sort the data
#     df_sorted = df_filtered.sort_values(by="Industry", ascending=False)
#     # df_sorted["Industry"] = pd.to_numeric(df_sorted["Industry"], errors="coerce")
#     # df_sorted = df.sort_values(by="Industry", ascending=False)
#     fig, ax = plt.subplots(figsize=(10, 6))
#     bars = ax.barh(
#         df_sorted["Sentiment"], 
#         df_sorted["Industry"], 
#         color="skyblue", 
#         edgecolor="black"
#     )
    
#     for bar in bars:
#         width = bar.get_width()
#         ax.text(
#             width, 
#             bar.get_y() + bar.get_height() / 2, 
#             f"{width}", 
#             ha="left", 
#             va="center", 
#             fontsize=10
#         )
    
#     # ax.set_title("Publication Name", fontsize=14)
#     ax.set_xlabel("News Count", fontsize=12)
#     ax.set_ylabel("Sentiment", fontsize=12)
#     ax.grid(axis="x", linestyle="--", alpha=0.7)

def generate_horizontal_bar_chartst(df):
    # Ensure the sentiment column contains expected categories
    df_filtered = df[df["Sentiment"] != "GrandTotal"]
    # Summarize the sentiment counts
    sentiment_counts = df_filtered['Industry'].values  # Extract sentiment counts
    sentiment_labels = df_filtered['Sentiment'].values  # Extract sentiment labels
    
    # Define colors for the sentiments: Negative = Red, Neutral = Blue, Positive = Green
    colors = ['salmon','skyblue', 'lightgreen']

     
    fig, ax = plt.subplots(figsize=(5, 5))
    wedges, texts, autotexts = ax.pie(
        sentiment_counts,
        labels=sentiment_labels,
        autopct='%1.1f%%',
        colors=colors,
        startangle=90,
        textprops={'fontsize': 12},  # Change font size of labels
        pctdistance=1.3,  # Position percentage labels outside the pie
    )
    # Customize the percentage text appearance
    for autotext in autotexts:
        autotext.set_fontsize(10)  # Set font size for percentages
        autotext.set_color('black')  # Set text color for percentages
    
    # Add a title
    # ax.set_title("Sentiment Distribution")
    
    # Save plot as image
    img_path11 = "horizontal_bar_chart.png"
    fig.savefig(img_path11, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path11

# Function to add image to slide
def add_image_to_slide5(slide, img_path11):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(7.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path11, left, top, width=width, height=height)

def generate_horizontal_bar_chartct(df):
    # Remove 'Client-' prefix from 'Entity' column
    # df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
    
    # Filter out unwanted rows
    df = df[df["City"] != "GrandTotal"]
    
    # Create the bar chart
    fig, ax = plt.subplots(figsize=(12, 6))  # Increase figure width for better label visibility
    x = range(len(df["City"]))  # Define x positions for the bars
    bars = ax.bar(
        x, 
        df["Industry"], 
        color="skyblue", 
        edgecolor="black"
    )
    
    # Add data labels on top of the bars without decimal
    for bar in bars:
        height = int(bar.get_height())  # Convert height to integer
        ax.text(
            bar.get_x() + bar.get_width() / 2, 
            height, 
            f"{height}", 
            ha="center", 
            va="bottom", 
            fontsize=12,
            fontweight="bold"
        )
    
    # Set chart title and axis labels
    # ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("City", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")
    
    # Customize x-axis ticks and labels for better visibility
    ax.set_xticks(x)
    ax.set_xticklabels(df["City"], rotation=45, ha="right", fontsize=12,fontweight="bold")

    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black", which="major", width=1, labelrotation=0)
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")
    
    # Add gridlines for better readability
    ax.grid(axis="y", linestyle="--", alpha=0.7)

# def generate_horizontal_bar_chartct(df):
#     df_filtered = df[df["City"] != "GrandTotal"]
    
#     # Sort the data
#     df_sorted = df_filtered.sort_values(by="Industry", ascending=False)
#     # df_sorted["Industry"] = pd.to_numeric(df_sorted["Industry"], errors="coerce")
#     # df_sorted = df.sort_values(by="Industry", ascending=False)
#     fig, ax = plt.subplots(figsize=(10, 6))
#     bars = ax.barh(
#         df_sorted["City"], 
#         df_sorted["Industry"], 
#         color="skyblue", 
#         edgecolor="black"
#     )
    
#     for bar in bars:
#         width = bar.get_width()
#         ax.text(
#             width, 
#             bar.get_y() + bar.get_height() / 2, 
#             f"{width}", 
#             ha="left", 
#             va="center", 
#             fontsize=10
#         )
    
#     # ax.set_title("Publication Name", fontsize=14)
#     ax.set_xlabel("News Count", fontsize=12)
#     ax.set_ylabel("City", fontsize=12)
#     ax.grid(axis="x", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path12 = "horizontal_bar_chart.png"
    fig.savefig(img_path12, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path12

# Function to add image to slide
def add_image_to_slide6(slide, img_path12):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path12, left, top, width=width, height=height)


def generate_grouped_bar_chartpt(df):
    """
    Generates a grouped bar chart for Publication Type, showing No of Publications and AVG News Count.

    Parameters:
        df (pd.DataFrame): A DataFrame containing:
                          - 'Publication Type': Categories for the x-axis.
                          - 'No of Publications': Number of publications.
                          - 'AVG News Count': Average news count.
        img_name (str): Name of the image file to save. Default is 'publication_bar_chart.png'.

    Returns:
        str: Path to the saved grouped bar chart image.
    """
    # Extract data
    publication_types = df["Publication Type"]
    num_publications = df["No of Publications"]
    avg_news_count = df["AVG News Count"]
    
    # Define positions on the x-axis
    x = np.arange(len(publication_types))  # Numeric positions for each category
    bar_width = 0.4  # Width of each bar
    
    # Create the plot
    fig, ax = plt.subplots(figsize=(10, 6))
    bars1 = ax.bar(x - bar_width/2, num_publications, bar_width, label="No of Publications", color="skyblue", edgecolor="black")
    bars2 = ax.bar(x + bar_width/2, avg_news_count, bar_width, label="AVG News Count", color="orange", edgecolor="black")
    
    # Add data labels
    for bars in (bars1, bars2):
        for bar in bars:
            height = bar.get_height()
            ax.text(
                bar.get_x() + bar.get_width() / 2, 
                height + 0.5,  # Add a slight offset above the bar
                f"{int(height)}", 
                ha="center", 
                va="bottom", 
                fontsize=10
            )

    # Chart aesthetics
    ax.set_xlabel("Publication Type", fontsize=12)
    ax.set_ylabel("Count", fontsize=12)
    # ax.set_title("Publication Type with Total Publications and AVG News Count", fontsize=14)
    ax.set_xticks(x)
    ax.set_xticklabels(publication_types, rotation=45, ha="right")
    ax.legend(title="Metrics", fontsize=10)
    ax.grid(axis="y", linestyle="--", alpha=0.7)

    # Save plot as image
    img_path13 = "horizontal_bar_chart.png"
    fig.savefig(img_path13, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path13

    # Function to add image to slide
def add_image_to_slide9(slide, img_path13):
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path13, left, top, width=width, height=height)
    
# Function to clean text
def clean_text(text):
    text = text.lower()  # Convert to lowercase
    text = re.sub(r"http\S+|www\S+|https\S+", '', text, flags=re.MULTILINE)  # Remove URLs
    text = re.sub(r'\@\w+|\#', '', text)  # Remove mentions and hashtags
    text = re.sub(r'[^a-zA-Z\s]', '', text)  # Remove non-alphabetic characters
    text = re.sub(r'\s+', ' ', text).strip()  # Remove extra whitespace
    return text
    
# Function to generate word cloud
def generate_word_cloud(df):
    text = ' '.join(df['Headline'].astype(str))
    text = clean_text(text)  # Clean the text
    stopwords = set(STOPWORDS)
    wordcloud = WordCloud(stopwords=stopwords, background_color="white" ,width=550,
        height=450,max_font_size=90, max_words=120,colormap='Set1',collocations=False).generate(text)
    
    # Plotting the word cloud
    fig, ax = plt.subplots(figsize=(6, 6), facecolor = 'black', edgecolor='black')
    ax.imshow(wordcloud, interpolation='bilinear')
    # ax.tight_layout(pad = 0) 
    ax.axis('off')
    
    # Save plot as image
    img_path11 = "wordcloud.png"
    fig.savefig(img_path11, dpi=300, bbox_inches='tight')
    plt.close(fig)
    
    return img_path11

    # # Example usage
    # img_path11 = generate_word_cloud(df)
    # print(f"Word cloud saved at: {img_path11}")

# Function to add image to slide (similar to the example you shared)
def add_image_to_slide11(slide, img_path11):
    from pptx.util import Inches
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(10)  # Adjust width
    height = Inches(6)  # Adjust height
    slide.shapes.add_picture(img_path11, left, top, width=width, height=height)
    
def top_10_dfs(df_list, file_name, comments, top_11_flags):
    writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
    row = 2
    for dataframe, comment, top_11_flag in zip(df_list, comments, top_11_flags):
        if top_11_flag:
            top_df = dataframe.head(50)  # Select the top 11 rows for specific DataFrames
        else:
            top_df = dataframe  # Leave other DataFrames unchanged

        top_df.to_excel(writer, sheet_name="Top 10 Data", startrow=row, index=True)
        row += len(top_df) + 2  # Move the starting row down by len(top_df) + 2 rows

    # Create a "Report" sheet with all the DataFrames
    for dataframe, comment in zip(df_list, comments):
        dataframe.to_excel(writer, sheet_name="Report", startrow=row, index=True, header=True)
        row += len(dataframe) + 2  # Move the starting row down by len(dataframe) + 2 rows

    writer.close()    
    

# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("Industry Data Insights/Tables Dashboard")

# Sidebar for file upload and download options
st.sidebar.title("Upload an Industry data file for tables")

# File Upload Section
file = st.sidebar.file_uploader("Upload Data File (Excel or CSV)", type=["xlsx", "csv"])

if file:
    st.sidebar.write("File Uploaded Successfully!")

    # Load data
    data = load_data(file)

    if data is not None:
        # Data Preview Section (optional)
        # st.write("## Data Preview")
        # st.write(data)

        # Data preprocessing
        # data.drop(columns=data.columns[20:], axis=1, inplace=True)
        # Process columns as required
     #    existing_columns = data.columns.tolist()
     #    influencer_index = existing_columns.index('Influencer')
     #    country_index = existing_columns.index('Country')


        
     #    new_order = (
     #     existing_columns[:influencer_index + 1] +  # All columns up to and including 'Influencer'
     #    ['Entity', 'Reach', 'Sentiment', 'Keywords', 'State', 'City', 'Engagement'] +  # Adding new columns
     #    existing_columns[influencer_index + 1:country_index + 1]  # All columns between 'Influencer' and 'Country'
     # )
     #    data = data[new_order]
        data.drop(columns=data.columns[20:], axis=1, inplace=True)
        data['Influencer'] = data['Influencer'].fillna('Bureau News')
        data = data.rename({'Influencer': 'Journalist'}, axis=1)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        # Check if specific columns exist before dropping duplicates
        if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)

        if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
            
        if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)

        finaldata = data
        finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()

        # Share of Voice (SOV) Calculation
        En_sov = pd.crosstab(finaldata['Entity'], columns='News Count', values=finaldata['Entity'], aggfunc='count').round(0)
        En_sov.sort_values('News Count', ascending=False)
        En_sov['% '] = ((En_sov['News Count'] / En_sov['News Count'].sum()) * 100).round(2)
        Sov_table = En_sov.sort_values(by='News Count', ascending=False)
        Sov_table.loc['Total'] = Sov_table.sum(numeric_only=True, axis=0)
        Entity_SOV1 = Sov_table
        Entity_SOV3 = pd.DataFrame(Entity_SOV1.to_records()).round()
        Entity_SOV3['% '] = Entity_SOV3['% '].astype(int)
        Entity_SOV3['% '] = Entity_SOV3['% '].astype(str) + '%'
        # Entity_SOV3 = pd.DataFrame(Entity_SOV3.to_records())
        total_news_count = Entity_SOV3.loc[Entity_SOV3["Entity"] == "Total", "News Count"].values[0]

        # Additional MOM DataFrames
        finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()
        sov_dt = pd.crosstab((finaldata['Date'].dt.to_period('M')), finaldata['Entity'], margins=True, margins_name='Total')
        sov_dt1 = pd.DataFrame(sov_dt.to_records())
        selected_columndt = sov_dt1

        selected_columndt = selected_columndt.sort_values(by="Industry", ascending=False)

        # Extract the top 3 publications and their counts
        topdt_1 = selected_columndt.iloc[1:2]  # First publication
        topc_2 = selected_columndt.iloc[2:3]  # Second publication
        # topc_3 = sov_dt1.iloc[2:3]  # Third publication

        # Save them in separate DataFrames
        df_topdt1 = topdt_1.reset_index(drop=True)
        df_topc2 = topc_2.reset_index(drop=True)
        # df_topc3 = topc_3.reset_index(drop=True)

        # Extract publication name and count for the top 3
        topdt_1_name = df_topdt1.iloc[0]["Date"]
        topdt_1_count = df_topdt1.iloc[0]["Industry"]

        topdt_2_name = df_topc2.iloc[0]["Date"]
        topdt_2_count = df_topc2.iloc[0]["Industry"]
        
        
        #Publication Name
        pub_table = pd.crosstab(finaldata['Publication Name'], finaldata['Entity'])
        pubs_table=pub_table.sort_values('Industry',ascending=False).round()
        pubs_table['% '] = ((pubs_table['Industry'] / pubs_table['Industry'].sum())*100).round(2)#.astype(str) + '%'
        pubs_table.loc['GrandTotal']= pubs_table.sum(numeric_only=True,axis=0)
        pubs_table['Industry'] = pubs_table['Industry'].astype(int)
        pubs_table['% '] = pubs_table['% '].astype(int)
        pubs_table = pd.DataFrame(pubs_table.to_records())
        pubs_table['% '] = pubs_table['% '].astype(str) + '%'
        
        pubs_table1 = pubs_table.head(10)
        pubs_table1.loc['Total'] = pubs_table1.sum(numeric_only=True, axis=0)
        # pubs_table1['% '] = pubs_table1['% '].astype(int)
        pubs_table1 = pd.DataFrame(pubs_table1.to_records())
        # pubs_table1['% '] = pubs_table1['% '].astype(str) + '%'

        # Extract the top 3 publications and their counts
        top_1 = pubs_table1.iloc[0:1]  # First publication
        top_2 = pubs_table1.iloc[1:2]  # Second publication
        # Check if a third publication exists
        if len(pubs_table1) > 2:
            top_3 = pubs_table1.iloc[2:3]  # Third publication
            df_top3 = top_3.reset_index(drop=True)
            top_3_name = df_top3.iloc[0]["Publication Name"]
            top_3_count = df_top3.iloc[0]["Industry"]
        else:
            top_3_name = ""
            top_3_count = 0  # You can assign any default value for count

        # Save the first two publications in separate DataFrames
        df_top1 = top_1.reset_index(drop=True)
        df_top2 = top_2.reset_index(drop=True)

        # Extract publication name and count for the top 2
        top_1_name = df_top1.iloc[0]["Publication Name"]
        top_1_count = df_top1.iloc[0]["Industry"]

        top_2_name = df_top2.iloc[0]["Publication Name"]
        top_2_count = df_top2.iloc[0]["Industry"]


        # # Extract the top 3 publications and their counts
        # top_1 = pubs_table1.iloc[0:1]  # First publication
        # top_2 = pubs_table1.iloc[1:2]  # Second publication
        # top_3 = pubs_table1.iloc[2:3]  # Third publication
        # top_4 = pubs_table1.iloc[10:11]  # fourth publication

        # # Save them in separate DataFrames
        # df_top1 = top_1.reset_index(drop=True)
        # df_top2 = top_2.reset_index(drop=True)
        # df_top3 = top_3.reset_index(drop=True)
        # df_top4 = top_4.reset_index(drop=True)

        # # Extract publication name and count for the top 3
        # top_1_name = df_top1.iloc[0]["Publication Name"]
        # top_1_count = df_top1.iloc[0]["Industry"]

        # top_2_name = df_top2.iloc[0]["Publication Name"]
        # top_2_count = df_top2.iloc[0]["Industry"]

        # top_3_name = df_top3.iloc[0]["Publication Name"]
        # top_3_count = df_top3.iloc[0]["Industry"]

        top_4_name = "Publication Name"
        top_4_per = "Count"

        
        # pub_table['Total'] = pub_table.sum(axis=1)
        # pubs_table = pub_table.sort_values('Total', ascending=False).round()
        # pubs_table.loc['GrandTotal'] = pubs_table.sum(numeric_only=True, axis=0)
        # pubs_table = pd.DataFrame(pubs_table.to_records())
        

        # PP = pd.crosstab(finaldata['Publication Name'], finaldata['Publication Type'])
        # PP['Total'] = PP.sum(axis=1)
        # PP_table = PP.sort_values('Total', ascending=False).round()
        # PP_table.loc['GrandTotal'] = PP_table.sum(numeric_only=True, axis=0)
        
        #Publication Name & Entity Table
        PT_Entity = pd.crosstab(finaldata['Publication Type'], finaldata['Entity'])
        PType_Entity=PT_Entity.sort_values('Industry',ascending=False).round()
        PType_Entity['% '] = ((PType_Entity['Industry'] / PType_Entity['Industry'].sum())*100).round()#.astype(str) + '%'
        PType_Entity.loc['GrandTotal']= PType_Entity.sum(numeric_only=True,axis=0)
        PType_Entity = pd.DataFrame(PType_Entity.to_records())
        PType_Entity['Industry'] = PType_Entity['Industry'].astype(int)
        PType_Entity['% '] = PType_Entity['% '].astype(int)
        PType_Entity['% '] = PType_Entity['% '].astype(str) + '%'
        # PT_Entity['Total'] = PT_Entity.sum(axis=1)
        # PType_Entity = PT_Entity.sort_values('Total', ascending=False).round()
        # PType_Entity.loc['GrandTotal'] = PType_Entity.sum(numeric_only=True, axis=0)
        # PType_Entity = pd.DataFrame(PType_Entity.to_records())

        # Extract the top 3 publications and their counts
        topt_1 = PType_Entity.iloc[0:1]  # First publication
        topt_2 = PType_Entity.iloc[1:2]  # Second publication
        # topt_3 = PType_Entity.iloc[2:3]  # Third publication

        # Save them in separate DataFrames
        df_topt1 = topt_1.reset_index(drop=True)
        df_topt2 = topt_2.reset_index(drop=True)
        # df_topt3 = topt_3.reset_index(drop=True)

        # Extract publication name and count for the top 3
        topt_1_name = df_topt1.iloc[0]["Publication Type"]
        topt_1_count = df_topt1.iloc[0]["Industry"]
        topt_1_per = df_topt1.iloc[0]["% "]

        topt_2_name = df_topt2.iloc[0]["Publication Type"]
        topt_2_count = df_topt2.iloc[0]["Industry"]
        topt_2_per = df_topt2.iloc[0]["% "]

        # topt_3_name = df_topt3.iloc[0]["Publication Type"]
        # topt_3_count = df_topt3.iloc[0]["Total"]

        # Journalist Table
        finaldata['Journalist'] = finaldata['Journalist'].str.split(',')
        finaldata = finaldata.explode('Journalist')
        jr_tab = pd.crosstab(finaldata['Journalist'], finaldata['Entity'])
        jr_tab = jr_tab.reset_index(level=0)
        newdata = finaldata[['Journalist', 'Publication Name']]
        Journalist_Table = pd.merge(jr_tab, newdata, how='inner', left_on=['Journalist'], right_on=['Journalist'])
        Journalist_Table.drop_duplicates(subset=['Journalist'], keep='first', inplace=True)
        valid_columns = Journalist_Table.select_dtypes(include='number').columns
        Journalist_Table['Total'] = Journalist_Table[valid_columns].sum(axis=1)
        Jour_table = Journalist_Table.sort_values('Total', ascending=False).round()
        bn_row = Jour_table.loc[Jour_table['Journalist'] == 'Bureau News']
        Jour_table = Jour_table[Jour_table['Journalist'] != 'Bureau News']
        Jour_table = pd.concat([Jour_table, bn_row], ignore_index=True)
#         Jour_table = Journalist_Table.reset_index(drop=True)
        Jour_table.loc['GrandTotal'] = Jour_table.sum(numeric_only=True, axis=0)
        columns_to_convert = Jour_table.columns.difference(['Journalist', 'Publication Name'])
        Jour_table[columns_to_convert] = Jour_table[columns_to_convert].astype(int)
        Jour_table.insert(1, 'Publication Name', Jour_table.pop('Publication Name'))
        Jour_table1 = Jour_table.head(10)

        Jour_table1.loc['Total'] = Jour_table1.sum(numeric_only=True, axis=0)
        # pubs_table1['% '] = pubs_table1['% '].astype(int)
        Jour_table1 = pd.DataFrame(Jour_table1.to_records())

        # Extract the top 3 publications and their counts
        topj_1 = Jour_table1.iloc[0:1]  # First publication
        topj_2 = Jour_table1.iloc[1:2]  # Second publication
        # Check if a third publication exists
        if len(Jour_table1) > 2:
            topj_3 = Jour_table1.iloc[2:3]  # Third publication
            df_topj3 = topj_3.reset_index(drop=True)
            topj_3_name = df_topj3.iloc[0]["Journalist"]
            topj_3_count = df_topj3.iloc[0]["Total"]
        else:
            topj_3_name = ""
            topj_3_count = 0  # You can assign any default value for count

        # Save the first two publications in separate DataFrames
        df_topj1 = topj_1.reset_index(drop=True)
        df_topj2 = topj_2.reset_index(drop=True)

        # Extract publication name and count for the top 2
        topj_1_name = df_topj1.iloc[0]["Journalist"]
        topj_1_count = df_topj1.iloc[0]["Total"]

        topj_2_name = df_topj2.iloc[0]["Journalist"]
        topj_2_count = df_topj2.iloc[0]["Total"]


        # # Extract the top 3 publications and their counts
        # topj_1 = Jour_table1.iloc[0:1]  # First publication
        # topj_2 = Jour_table1.iloc[1:2]  # Second publication
        # topj_3 = Jour_table1.iloc[2:3]  # Third publication
        # topj_4 = Jour_table1.iloc[10:11]  # Third publication

        # # Save them in separate DataFrames
        # df_topj1 = topj_1.reset_index(drop=True)
        # df_topj2 = topj_2.reset_index(drop=True)
        # df_topj3 = topj_3.reset_index(drop=True)
        # df_topj4 = topj_4.reset_index(drop=True)

        # # Extract publication name and count for the top 3
        # topj_1_name = df_topj1.iloc[0]["Journalist"]
        # topj_1_count = df_topj1.iloc[0]["Total"]

        # topj_2_name = df_topj2.iloc[0]["Journalist"]
        # topj_2_count = df_topj2.iloc[0]["Total"]

        # topj_3_name = df_topj3.iloc[0]["Journalist"]
        # topj_3_count = df_topj3.iloc[0]["Total"]

        topj_4_name = "Journalist Name"
        topj_4_count = "Count"

        # Extract the top 3 publications and their counts
        topjt_1 = Jour_table1.iloc[0:1]  # First publication
        topjt_2 = Jour_table1.iloc[1:2]  # Second publication
        # Check if a third publication exists
        if len(Jour_table1) > 2:
            topjt_3 = Jour_table1.iloc[2:3]  # Third publication
            df_topjt3 = topjt_3.reset_index(drop=True)
            topjt_3_name = df_topjt3.iloc[0]["Publication Name"]
            topjt_3_count = df_topjt3.iloc[0]["Total"]
        else:
            topjt_3_name = ""
            # topjt_3_count = 0  # You can assign any default value for count

        # Save the first two publications in separate DataFrames
        df_topjt1 = topjt_1.reset_index(drop=True)
        df_topjt2 = topjt_2.reset_index(drop=True)

        # Extract publication name and count for the top 2
        topjt_1_name = df_topjt1.iloc[0]["Publication Name"]
        # topjt_1_count = df_topjt1.iloc[0]["Total"]

        topjt_2_name = df_topjt2.iloc[0]["Publication Name"]
        # topjt_2_count = df_topjt2.iloc[0]["Total"]

        # # Extract the top 3 publications and their counts
        # topjt_1 = Jour_table1.iloc[0:1]  # First publication
        # topjt_2 = Jour_table1.iloc[1:2]  # Second publication
        # topjt_3 = Jour_table1.iloc[2:3]  # Third publication

        # # Save them in separate DataFrames
        # df_topjt1 = topjt_1.reset_index(drop=True)
        # df_topjt2 = topjt_2.reset_index(drop=True)
        # df_topjt3 = topjt_3.reset_index(drop=True)

        # # Extract publication name and count for the top 3
        # topjt_1_name = df_topjt1.iloc[0]["Publication Name"]
        # # top_1_count = df_topjt1.iloc[0]["Total"]

        # topjt_2_name = df_topjt2.iloc[0]["Publication Name"]
        # # top_2_count = df_topjt2.iloc[0]["Total"]

        # topjt_3_name = df_topjt3.iloc[0]["Publication Name"]
        # # top_3_count = df_topjt3.iloc[0]["Total"]

        # Group by Publication Name and calculate the News Count
        news_count = finaldata['Publication Name'].value_counts().reset_index()
        news_count.columns = ['Publication Name', 'News Count']

        # Identify articles written by Bureau and by Journalists
        finaldata['Is Bureau'] = finaldata['Journalist'].str.contains('Bureau', case=False, na=False)
        finaldata['Is Journalist'] = ~finaldata['Is Bureau']

        # Calculate the percentage of articles by Bureaus and Journalists
        bureau_count = finaldata.groupby('Publication Name')['Is Bureau'].sum().reset_index()
        journalist_count = finaldata.groupby('Publication Name')['Is Journalist'].sum().reset_index()

        # Merge the counts with the news_count dataframe
        final_df1 = news_count.merge(bureau_count, on='Publication Name')
        final_df1 = final_df1.merge(journalist_count, on='Publication Name')

        # Calculate percentages
        final_df1['% of articles by Bureaus'] = ((final_df1['Is Bureau'] / final_df1['News Count']) * 100).round()
        final_df1['% of articles by Journalists'] = ((final_df1['Is Journalist'] / final_df1['News Count']) * 100).round()

        # Select and rename columns
        final_df1 = final_df1[['Publication Name', 'News Count', '% of articles by Bureaus', '% of articles by Journalists']]
        final_df1 = final_df1.round(2)

        final_df1['% of articles by Bureaus'] = final_df1['% of articles by Bureaus'].astype(int)
        final_df1['% of articles by Journalists'] = final_df1['% of articles by Journalists'].astype(int)
        
        final_df1['% of articles by Bureaus'] = final_df1['% of articles by Bureaus'].astype(str) + '%'
        final_df1['% of articles by Journalists'] = final_df1['% of articles by Journalists'].astype(str) + '%'

        final_df11 = final_df1.head(10)

        # Extract the top 3 publications and their counts
        topb_1 = final_df11.iloc[0:1]  # First publication
        topb_2 = final_df11.iloc[1:2]  # Second publication
        topb_3 = final_df11.iloc[2:3]  # Third publication
        # top_4 = final_df11.iloc[10:11]  # fourth publication

        # Save them in separate DataFrames
        df_topb1 = topb_1.reset_index(drop=True)
        df_topb2 = topb_2.reset_index(drop=True)
        df_topb3 = topb_3.reset_index(drop=True)
        # df_top4 = top_4.reset_index(drop=True)

        # Extract publication name and count for the top 3
        topb_1_name = df_topb1.iloc[0]["Publication Name"]
        topb_1_count = df_topb1.iloc[0]["News Count"]
        topb_1_br = df_topb1.iloc[0]["% of articles by Bureaus"]
        topb_1_jr = df_topb1.iloc[0]["% of articles by Journalists"]

        topb_2_name = df_topb2.iloc[0]["Publication Name"]
        topb_2_count = df_topb2.iloc[0]["News Count"]
        topb_2_br = df_topb2.iloc[0]["% of articles by Bureaus"]
        topb_2_jr = df_topb2.iloc[0]["% of articles by Journalists"]

        topb_3_name = df_topb3.iloc[0]["Publication Name"]
        topb_3_count = df_topb3.iloc[0]["News Count"]
        topb_3_br = df_topb3.iloc[0]["% of articles by Bureaus"]
        topb_3_jr = df_topb3.iloc[0]["% of articles by Journalists"]

        # top_4_name = df_top4.iloc[0]["Industry"]
        # top_4_per = df_top4.iloc[0]["% "]

        # Calculate News Count for each Publication Type
        news_count = finaldata.groupby('Publication Type').size().reset_index(name='News Count')

        # Calculate the total News Count for percentage calculation
        total_news_count = news_count['News Count'].sum()

        # Calculate the percentage of total articles for each Publication Type
        news_count['%'] = (news_count['News Count'] / total_news_count) * 100

        # Count the number of unique Publications for each Publication Type
        no_of_publications = finaldata.groupby('Publication Type')['Publication Name'].nunique().reset_index(name='No of Publications')

        # Merge the counts with the news_count dataframe
        final_df7 = news_count.merge(no_of_publications, on='Publication Type')

        # Calculate the average number of articles per publication
        final_df7['AVG News Count'] = final_df7['News Count'] / final_df7['No of Publications']

        # Round percentage and AVG News Count to 2 decimal places
        final_df7['%'] = final_df7['%'].round(2)
        final_df7['AVG News Count'] = final_df7['AVG News Count'].round(2)

        # Add a 'Grand Total' row at the bottom
        grand_total = {
    'Publication Type': 'Grand Total',
    'News Count': final_df7['News Count'].sum(),
    '%': final_df7['%'].sum(),
    'No of Publications': final_df7['No of Publications'].sum(),
    'AVG News Count': (final_df7['News Count'].sum() / final_df7['No of Publications'].sum()).round(2)
}

        # final_df = final_df.append(grand_total, ignore_index=True)

        # Print or save the final table
        print(final_df7)
        # final_df.to_csv('output_table23.csv', index=False)  # Uncomment to save as CSV
        final_df7 = final_df7.sort_values("News Count", ascending=False)
        final_df7['AVG News Count'] = final_df7['AVG News Count'].astype(int)
        final_df7['%'] = final_df7['%'].astype(int)
        final_df7['%'] = final_df7['%'].astype(str) + '%'

        # Extract the top 3 publications and their counts
        topav_1 = final_df7.iloc[0:1]  # First publication
        topav_2 = final_df7.iloc[1:2]  # Second publication
        topav_3 = final_df7.iloc[2:3]  # Third publication
        # top_4 = final_df7.iloc[10:11]  # fourth publication

        # Save them in separate DataFrames
        df_topav1 = topav_1.reset_index(drop=True)
        df_topav2 = topav_2.reset_index(drop=True)
        df_topav3 = topav_3.reset_index(drop=True)
        # df_topav4 = topav_4.reset_index(drop=True)

        # Extract publication name and count for the top 3
        topav_1_name = df_topav1.iloc[0]["Publication Type"]
        topav_1_count = df_topav1.iloc[0]["News Count"]
        topav_1_br = df_topav1.iloc[0]["%"]
        topav_1_jr = df_topav1.iloc[0]["No of Publications"]
        topav_1_av = df_topav1.iloc[0]["AVG News Count"]

        topav_2_name = df_topav2.iloc[0]["Publication Type"]
        topav_2_count = df_topav2.iloc[0]["News Count"]
        topav_2_br = df_topav2.iloc[0]["%"]
        topav_2_jr = df_topav2.iloc[0]["No of Publications"]
        topav_2_av = df_topav2.iloc[0]["AVG News Count"]

        topav_3_name = df_topav3.iloc[0]["Publication Type"]
        topav_3_count = df_topav3.iloc[0]["News Count"]
        topav_3_br = df_topav3.iloc[0]["%"]
        topav_3_jr = df_topav3.iloc[0]["No of Publications"]
        topav_3_av = df_topav3.iloc[0]["AVG News Count"]



        Set_Entity = pd.crosstab(finaldata['Sentiment'],finaldata['Entity'])
        Senti_Entity=Set_Entity.sort_values('Industry',ascending=False).round()
        Senti_Entity['% '] = ((Senti_Entity['Industry'] / Senti_Entity['Industry'].sum())*100).round(2)#.astype(str) + '%'
        Senti_Entity.loc['GrandTotal']= Senti_Entity.sum(numeric_only=True,axis=0)
        Senti_Entity = pd.DataFrame(Senti_Entity.to_records())
        Senti_Entity['Industry'] = Senti_Entity['Industry'].astype(int)
        Senti_Entity['% '] = Senti_Entity['% '].astype(int)
        Senti_Entity['% '] = Senti_Entity['% '].astype(str) + '%'

        # Extract the top 3 publications and their counts
        tops_1 = Senti_Entity.iloc[0:1]  # First publication
        tops_2 = Senti_Entity.iloc[1:2]  # Second publication
        tops_3 = Senti_Entity.iloc[2:3]  # Third publication
        # top_4 = pubs_table1.iloc[10:11]  # fourth publication

        # Save them in separate DataFrames
        df_tops1 = tops_1.reset_index(drop=True)
        df_tops2 = tops_2.reset_index(drop=True)
        df_tops3 = tops_3.reset_index(drop=True)
        # df_top4 = top_4.reset_index(drop=True)

        # Extract publication name and count for the top 3
        tops_1_name = df_tops1.iloc[0]["Sentiment"]
        tops_1_count = df_tops1.iloc[0]["Industry"]
        tops_1_per = df_tops1.iloc[0]["% "]

        tops_2_name = df_tops2.iloc[0]["Sentiment"]
        tops_2_count = df_tops2.iloc[0]["Industry"]
        tops_2_per = df_tops2.iloc[0]["% "]

        tops_3_name = df_tops3.iloc[0]["Sentiment"]
        tops_3_count = df_tops3.iloc[0]["Industry"]
        tops_3_per = df_tops3.iloc[0]["% "]
        

        Ct_Entity = pd.crosstab(finaldata['City'],finaldata['Entity'])
        City_Entity=Ct_Entity.sort_values('Industry',ascending=False).round()
        City_Entity['% '] = ((City_Entity['Industry'] / City_Entity['Industry'].sum())*100).round(2)#.astype(str) + '%'
        City_Entity.loc['GrandTotal']= City_Entity.sum(numeric_only=True,axis=0)
        City_Entity = pd.DataFrame(City_Entity.to_records())
        City_Entity['Industry'] = City_Entity['Industry'].astype(int)
        City_Entity['% '] = City_Entity['% '].astype(int)
        City_Entity['% '] = City_Entity['% '].astype(str) + '%'

        # Extract the top 3 publications and their counts
        topct_1 = City_Entity.iloc[0:1]  # First publication
        topct_2 = City_Entity.iloc[1:2]  # Second publication
        topct_3 = City_Entity.iloc[2:3]  # Third publication
        topct_4 = City_Entity.iloc[3:4]  # fourth publication

        # Save them in separate DataFrames
        df_topct1 = topct_1.reset_index(drop=True)
        df_topct2 = topct_2.reset_index(drop=True)
        df_topct3 = topct_3.reset_index(drop=True)
        df_topct4 = topct_4.reset_index(drop=True)

        # Extract publication name and count for the top 3
        topct_1_name = df_topct1.iloc[0]["City"]
        topct_1_count = df_topct1.iloc[0]["Industry"]
        topct_1_per = df_topct1.iloc[0]["% "]

        topct_2_name = df_topct2.iloc[0]["City"]
        topct_2_count = df_topct2.iloc[0]["Industry"]
        topct_2_per = df_topct2.iloc[0]["% "]

        topct_3_name = df_topct3.iloc[0]["City"]
        topct_3_count = df_topct3.iloc[0]["Industry"]
        topct_3_per = df_topct3.iloc[0]["% "]

        topct_4_name = df_topct4.iloc[0]["City"]
        topct_4_count = df_topct4.iloc[0]["Industry"]
        topct_4_per = df_topct4.iloc[0]["% "]

        
        
        # Remove square brackets and single quotes from the 'Journalist' column
        data['Journalist'] = data['Journalist'].str.replace(r"^\['(.+)'\]$", r"\1", regex=True)
        # Fill missing values in 'Influencer' column with 'Bureau News'
        # data['Journalist'] = data['Journalist'].fillna('Bureau News')

        # Function to classify news exclusivity and topic
        def classify_exclusivity(row):
            entity_name = row['Entity']
            if entity_name.lower() in row['Headline'].lower():
                return "Exclusive"
            else:
                return "Not Exclusive"

        finaldata['Exclusivity'] = finaldata.apply(classify_exclusivity, axis=1)

        # # Define a dictionary of keywords for each entity
        # entity_keywords = {
        #     'Amazon': ['Amazon', 'Amazons', 'amazon'],
        #     # Add other entities and keywords here
        # }

        # def qualify_entity(row):
        #     entity_name = row['Entity']
        #     text = row['Headline']
        #     if entity_name in entity_keywords:
        #         keywords = entity_keywords[entity_name]
        #         if any(keyword in text for keyword in keywords):
        #             return "Qualified"
        #     return "Not Qualified"

        # finaldata['Qualification'] = finaldata.apply(qualify_entity, axis=1)

        # Topic classification
        topic_mapping = {
              'Merger': ['merger', 'merges'],
                
              'Acquire': ['acquire', 'acquisition', 'acquires'],
                
              'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
                'Partnership': ['IPO','ipo'],
               'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                    'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
                
               'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
                
              'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
                
              'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
            'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
                
               'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
                
               'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
                
                'Awards & Recognition': ['award', 'awards'],
                
                'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
            
            'Sale - Offers - Discounts' : ['sale','offers','discount','discounts','discounted']
        }

        def classify_topic(headline):
            for topic, words in topic_mapping.items():
                if any(word in headline.lower() for word in words):
                    return topic
            return 'Other'

        finaldata['Topic'] = finaldata['Headline'].apply(classify_topic)

        dfs = [Entity_SOV3, sov_dt1, pubs_table,final_df11 , Jour_table, PType_Entity, final_df7,Senti_Entity,City_Entity]
        comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Publication Name with Bureau and Journalist Percentages', 'Journalist Table','PubType Entity Table',
                    'Publication Type with Total Publications and Avg news count','Sentiment Table','Citywise Table']

        # Sidebar for download options
        st.sidebar.write("## Download Options")
        download_formats = st.sidebar.selectbox("Select format:", ["Excel", "CSV", "Excel (Entity Sheets)"])

        if st.sidebar.button("Download Data"):
            if download_formats == "Excel":
                # Download all DataFrames as a single Excel file
                excel_io = io.BytesIO()
                with pd.ExcelWriter(excel_io, engine="xlsxwriter") as writer:
                    for df, comment in zip(dfs, comments):
                        df.to_excel(writer, sheet_name=comment, index=False)
                excel_io.seek(0)
                b64_data = base64.b64encode(excel_io.read()).decode()
                href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="data.xlsx">Download Excel</a>'
                st.sidebar.markdown(href_data, unsafe_allow_html=True)

            elif download_formats == "CSV":
                # Download all DataFrames as CSV
                csv_io = io.StringIO()
                for df in dfs:
                    df.to_csv(csv_io, index=False)
                csv_io.seek(0)
                b64_data = base64.b64encode(csv_io.read().encode()).decode()
                href_data = f'<a href="data:text/csv;base64,{b64_data}" download="data.csv">Download CSV</a>'
                st.sidebar.markdown(href_data, unsafe_allow_html=True)

            elif download_formats == "Excel (Entity Sheets)":
                # Download DataFrames as Excel with separate sheets by entity
                excel_io = io.BytesIO()
                with pd.ExcelWriter(excel_io, engine="xlsxwriter") as writer:
                    create_entity_sheets(finaldata, writer)
                excel_io.seek(0)
                b64_data = base64.b64encode(excel_io.read()).decode()
                href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="entity_sheets.xlsx">Download Entity Sheets</a>'
                st.sidebar.markdown(href_data, unsafe_allow_html=True)
                
         
        # Download selected DataFrame
        st.sidebar.write("## Download Selected DataFrame")
        
        dataframes_to_download = {
            "Entity_SOV1": Entity_SOV3,
            "Data": data,
            "Finaldata": finaldata,
            "Month-on-Month":sov_dt1,
            "Publication Table":pubs_table,
            "Journalist Table":Jour_table,
            # "Publication Type and Name Table":PP_table,
            "Publication Type Table with Entity":PType_Entity,
            # "Publication type,Publication Name and Entity Table":ppe1,
            "Entity-wise Sheets": finaldata,  # Add this option to download entity-wise sheets
            "Publication Name with Bureau and Journalist Percentages" : final_df11,
            "Publication Type with Total Publications and Avg news count" : final_df7,
            "Sentiment table":Senti_Entity,
           "Citywise table" : City_Entity
        }
        selected_dataframe = st.sidebar.selectbox("Select DataFrame:", list(dataframes_to_download.keys()))
        
        if st.sidebar.button("Download Selected DataFrame"):
            if selected_dataframe in dataframes_to_download:
                # Create a link to download the selected DataFrame in Excel
                selected_df = dataframes_to_download[selected_dataframe]
                excel_io_selected = io.BytesIO()
                with pd.ExcelWriter(excel_io_selected, engine="xlsxwriter", mode="xlsx") as writer:
                    selected_df.to_excel(writer, index=True)
                excel_io_selected.seek(0)
                b64_selected = base64.b64encode(excel_io_selected.read()).decode()
                href_selected = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_selected}" download="{selected_dataframe}.xlsx">Download {selected_dataframe} Excel</a>'
                st.sidebar.markdown(href_selected, unsafe_allow_html=True)
                
                
        # Download All DataFrames as a Single Excel Sheet
        st.sidebar.write("## Download All DataFrames as a Single Excel Sheet")
        file_name_all = st.sidebar.text_input("Enter file name for all DataFrames", "all_dataframes.xlsx")
        
        if st.sidebar.button("Download All DataFrames"):
            # List of DataFrames to save
            dfs = [Entity_SOV3, sov_dt1, pubs_table,final_df11 , Jour_table, PType_Entity, final_df7,Senti_Entity,City_Entity]
            comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Publication Name with Bureau and Journalist Percentages', 'Journalist Table','PubType Entity Table',
                    'Publication Type with Total Publications and Avg news count','Sentiment Table','Citywise Table']
            
            entity_info = """Entity:
Time Period of analysis: 19th April 2023 to 18th April 2024
Source: (Online) Meltwater, Select 100 online publications, which include General mainlines, Business and Financial publications, news age media, technology publications.
News search: All Articles: entity mentioned at least once in the article"""
            excel_io_all = io.BytesIO()
            multiple_dfs(dfs, 'Tables', excel_io_all, comments, entity_info)
            excel_io_all.seek(0)
            b64_all = base64.b64encode(excel_io_all.read()).decode()
            href_all = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_all}" download="{file_name_all}">Download All DataFrames Excel</a>'
            st.sidebar.markdown(href_all, unsafe_allow_html=True)

        st.write("## Preview Selected DataFrame")
        selected_dataframe = st.selectbox("Select DataFrame to Preview:", list(dataframes_to_download.keys()))
        st.dataframe(dataframes_to_download[selected_dataframe])

    # Load the image files
    img_path = r"New logo snip.PNG"
    img_path1 = r"New Templete main slide.PNG"

    # Create a new PowerPoint presentation with widescreen dimensions
    prs = Presentation()               
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add the first slide with the image
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture( img_path1,Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Add the text box above the image
    textbox_left = Inches(0.5)  # Adjust the left position as needed
    textbox_top = Inches(5)   # Adjust the top position as needed
    textbox_width = Inches(15)  # Adjust the width as needed
    textbox_height = Inches(1)  # Adjust the height as needed

    text_box = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = "Industry"
        
    # Set font size to 30 and make the text bold and white
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(50)
            run.font.bold = True
#           run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(255, 255, 255)  # White color
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0.0)  # Adjust the left position as needed
    top = prs.slide_height - Inches(1)  # Adjust the top position as needed
    slide.shapes.add_picture(img_path, left, top, height=Inches(1))  # Adjust the height as needed 

        
    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames

    # Set title text and format for Parameters slide
    header_text = "Parameters"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.7))
    header_frame = header_shape.text_frame
    header_frame.text = header_text

    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    # Add Time Period text
    time_period_text = "Time Period : 19th April 2023 to 18th April 2024"
    time_period_shape = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(14), Inches(0.5))
    time_period_frame = time_period_shape.text_frame
    time_period_frame.text = time_period_text
    # time_period_frame.paragraphs[0].font.bold = True
    time_period_frame.paragraphs[0].font.size = Pt(24)
    time_period_frame.paragraphs[0].font.name = 'Gill Sans'


    # Add Source text
    source_text = "Source : (Online)Meltwater, Select 100 online publications, which include General mainlines, Business and Financial publications, news age media, technology publications."
    source_shape = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))  # Adjusted width
    source_frame = source_shape.text_frame
    source_frame.word_wrap = True  # Enable text wrapping
    p = source_frame.add_paragraph()  # Create a paragraph for text
    p.text = source_text  # Set the text

    p.font.size = Pt(24)
    p.font.name = 'Gill Sans'  # Changed to Arial for compatibility

    # Add News Search text
    news_search_text = "News Search : All Articles: entity mentioned at least once in the article "
    news_search_shape = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(10), Inches(0.75))  # Adjusted width and height
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True  # Enable text wrapping
    p2 = news_search_frame.add_paragraph()  # Create a paragraph for text
    p2.text = news_search_text  # Set the text

    # Set font properties after text is added
    # p2.font.bold = True
    p2.font.size = Pt(24)
    p2.font.name = 'Gill Sans'  # Changed to Arial for compatibility
        
    # Add the first slide with the image
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture( img_path1,Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Add the text box above the image
    textbox_left = Inches(0.5)  # Adjust the left position as needed
    textbox_top = Inches(5)   # Adjust the top position as needed
    textbox_width = Inches(15)  # Adjust the width as needed
    textbox_height = Inches(1)  # Adjust the height as needed

    text_box = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = ("Online Media\n"
                      "News Analysis\n"
                    "By Media Research & Analytics Team")

    # Set font size to 30 and make the text bold and white
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(50)
            run.font.bold = True
#           run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(255, 255, 255)  # White color
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0.0)  # Adjust the left position as needed
    top = prs.slide_height - Inches(1)  # Adjust the top position as needed
    slide.shapes.add_picture(img_path, left, top, height=Inches(1))  # Adjust the height as needed 
         
    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames

    # Set title text and format for Parameters slide
    header_text = "Inferences and Recommendations"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(14), Inches(0.7))
    header_frame = header_shape.text_frame
    header_frame.text = header_text
    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  


    # Add SOV text
    sov_text = ("Share of Voice :")
    sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
    sov_text_frame = sov_text_shape.text_frame
    sov_text_frame.word_wrap = True
    sov_text_frame.clear()  # Clear any default paragraph

    p = sov_text_frame.add_paragraph()
    p.text = "Share of Voice :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True

    sov_text = (
    f"The Steel Industry sector has a total of {total_news_count} prominent mentions.\n"
    "This represents the overall volume of online media coverage related to Steel Industry.\n",)
    # "IIT Delhi follows IIT Madras, capturing 21% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
    # "IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 20%, 17%, and 6% of the SOV respectively.\n"
    # "IIT Ropar holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e. last in the SOV.\n"
    # "Despite ranking lower in terms of SOV, IIT Ropar's presence indicates some level of visibility and recognition within the online media landscape.\n"
    # "Given the relatively lower SOV compared to peers like IIT Delhi, IIT Madras, and others, there are opportunities for IIT Ropar to enhance its online presence and visibility through strategic communications efforts.\n"
    # "IIT Ropar has received 239 all mentions and 44 prominent articles in online media and stands last in both the SOVs.\n"
    #     )
    sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
    sov_text_frame = sov_text_shape.text_frame
    sov_text_frame.word_wrap = True
    sov_text_frame.clear()  # Clear any default paragraph


    p = sov_text_frame.add_paragraph()
    p.text = (
    f"The Steel Industry sector has a total of {total_news_count} prominent mentions.\n"
    "This represents the overall volume of online media coverage related to Steel Industry.\n")
    # "IIT Delhi follows IIT Madras, capturing 21% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
    # "IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 20%, 17%, and 6% of the SOV respectively.\n"
    # "IIT Ropar holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e. last in the SOV.\n"
    # "Despite ranking lower in terms of SOV, IIT Ropar's presence indicates some level of visibility and recognition within the online media landscape.\n"
    # "Given the relatively lower SOV compared to peers like IIT Delhi, IIT Madras, and others, there are opportunities for IIT Ropar to enhance its online presence and visibility through strategic communications efforts.\n"
    # "IIT Ropar has received 239 all mentions and 44 prominent articles in online media and stands last in both the SOVs.\n"
    # )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'

    # Add Source text
    source_text = ("Publications :")
    source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.8), Inches(14), Inches(1))
    source_frame = source_shape.text_frame
    source_frame.word_wrap = True
    source_frame.clear()  # Clear any default paragraph
    p = source_frame.add_paragraph()
    p.text = "Publications :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True


    source_text = (
    f"Of the {total_news_count} articles having prominent mentions of Steel Industry, top 10 publications, among the 100 publications, contributed {top_4_name} articles ({top_4_per} of all mentions).\n"
f"{top_1_name} is the most prominent publication covering Steel Industry, with {top_1_count} news articles. {top_2_name} and {top_3_name} follow, with {top_2_count} and {top_3_count} news articles, respectively.\n"
# "The top 10 publications writing articles on IIT Ropar contribute 86% of the total 44 articles.\n" 
)
    source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.1), Inches(14), Inches(1))
    source_frame = source_shape.text_frame
    source_frame.word_wrap = True
    source_frame.clear()  # Clear any default paragraph
    p = source_frame.add_paragraph()
    p.text = (
    f"Of the {total_news_count} articles having prominent mentions of Steel Industry, top 10 publications, among the 100 publications, contributed {top_4_name}  articles ({top_4_per}  of all mentions).\n"
f"{top_1_name} is the most prominent publication covering Steel Industry, with {top_1_count} news articles. {top_2_name} and {top_3_name} follow, with {top_2_count} and {top_3_count} news articles, respectively.\n"
# "The top 10 publications writing articles on IIT Ropar contribute 86% of the total 44 articles.\n" 
)
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'

    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)


    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames


    # Set title text and format for Parameters slide
    header_text = "Inferences and Recommendations"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
    header_frame = header_shape.text_frame
    header_frame.text = header_text 
    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


    # Add News Search text
    news_search_text = ("Journalists :")
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = "Journalists :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True

    # Add News Search text
    news_search_text = (f"{topj_1_name} from {topjt_1_name} leads with {topj_1_count} articles, indicating a strong focus on Steel Industry topics. {topj_2_name} , from {topjt_2_name}, follows closely with {topj_2_count} articles and {topj_3_name} from {topjt_3_name} with {topj_3_count} articles.\n"
                    "The Economic Times is prominently featured, with two journalists in the top 10: Nikita Periwal and Twesh Mishra.\n"
                    f"This reinforces The Economic Times as a major publication in the Steel Industry sector.The top 10 journalists have written {topj_4_count} prominent articles on Steel Industry out of {total_news_count} i.e only 8% . Majority of the articles i.e 73% of them were filed by Bureaus.\n")
                   #  "A total of 387 journalists have written 1155 articles covering IIT Ropar and its competitors.\n"
                   #  "Out of which, 5 journalists have specifically written 5 articles mentioning IIT Ropar i.e of the total journalists writing on IIT Ropar and its competitors only 1% of them have mentioned IIT Ropar in their articles.\n"
                   #  "While this constitutes a very less number, there is an large opportunity for IIT Ropar to engage with the remaining 882 journalists to enhance its news coverage and reach.\n"
                   # )
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = (f"{topj_1_name} from {topjt_1_name} leads with {topj_1_count} articles, indicating a strong focus on Steel Industry topics. {topj_2_name} , from {topjt_2_name}, follows closely with {topj_2_count} articles and {topj_3_name} from {topjt_3_name} with {topj_3_count} articles.\n"
                    "The Economic Times is prominently featured, with two journalists in the top 10: Nikita Periwal and Twesh Mishra.\n"
                    f"This reinforces The Economic Times as a major publication in the Steel Industry sector.The top 10 journalists have written {topj_4_count} prominent articles on Steel Industry out of {total_news_count} i.e only 8% . Majority of the articles i.e 73% of them were filed by Bureaus.\n")
                   #   "A total of 387 journalists have written 1155 articles covering IIT Ropar and its competitors.\n"
                   #  "Out of which, 5 journalists have specifically written 5 articles mentioning IIT Ropar i.e of the total journalists writing on IIT Ropar and its competitors only 1% of them have mentioned IIT Ropar in their articles.\n"
                   #  "While this constitutes a very less number, there is an large opportunity for IIT Ropar to engage with the remaining 882 journalists to enhance its news coverage and reach.\n"
                   # )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'

    # Add News Search text
    news_search_text = ("Publication Types :" )
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.6), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = "Publication Type :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True

    news_search_text = (f"{topt_1_name} publications account for the majority of news coverage with {topt_1_count} articles, which is approximately {topt_1_per} of the total.\n"
f"{topt_2_name}  publications contribute {topt_2_count} articles, making up about {topt_2_per} of the total coverage.Advertising/PR/Media, Digital First, and Technology publications collectively contribute only 112 articles (approximately 2% of the total).\n"
# "IIT Ropar may find value in engaging more with General and Business along with technology, and digital-first publications to expand her reach and visibility among broader audiences.\n"
                   )
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.0), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = (f"{topt_1_name} publications account for the majority of news coverage with {topt_1_count} articles, which is approximately {topt_1_per} of the total.\n"
f"{topt_2_name}  publications contribute {topt_2_count} articles, making up about {topt_2_per} of the total coverage.Advertising/PR/Media, Digital First, and Technology publications collectively contribute only 112 articles (approximately 2% of the total).\n"
# "IIT Ropar may find value in engaging more with General and Business along with technology, and digital-first publications to expand her reach and visibility among broader audiences.\n"
                   )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'
        
    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames
        
    # Set title text and format for Parameters slide
    header_text = "Inferences and Recommendations"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
    header_frame = header_shape.text_frame
    header_frame.text = header_text
    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


    # # Add Time Period text
    time_period_text = ("Monthly Coverage , Peak and Topics :")
    time_period_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
    time_period_frame = time_period_shape.text_frame
    time_period_frame.text = time_period_text
    time_period_frame.word_wrap = True
    time_period_frame.clear() 

    p = time_period_frame.add_paragraph()
    p.text = "Monthly Coverage , Peak and Topics :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True


    time_period_text = (f"The total number of news articles over the period is {total_news_count}. There is noticeable fluctuation in coverage from month to month, with periods of both increase and decline.\n"
f"{topdt_1_name} saw the highest number of articles, with {topdt_1_count} mentions. This spike suggests a significant event or increased media focus on Steel Industry during that month.\n"
f"{topdt_2_name} has the 2nd highest peak coverage, with {topdt_2_count} mentions.\n"
f"There was  peak in {topdt_1_name} due to following news:The increase in volume is due to Prime Minister Narendra Modi's visit to Chhattisgarh and the launch of multiple development projects, including the dedication of the NMDC Steel Plant in Bastar. The projects are expected to provide employment opportunities and contribute to the reduction of debt burden. The political controversy surrounding the Nagarnar Steel Plant and the statements made by Chief Minister Bhupesh Baghel have also contributed to the spike in volume\n"
                   )
    time_period_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(14), Inches(0.5))
    time_period_frame = time_period_shape.text_frame
    time_period_frame.text = time_period_text
    time_period_frame.word_wrap = True
    time_period_frame.clear() 

    p = time_period_frame.add_paragraph()
    p.text = (f"The total number of news articles over the period is {total_news_count}. There is noticeable fluctuation in coverage from month to month, with periods of both increase and decline.\n"
f"{topdt_1_name} saw the highest number of articles, with {topdt_1_count} mentions. This spike suggests a significant event or increased media focus on Steel Industry during that month.\n"
f"{topdt_2_name} has the 2nd highest peak coverage, with {topdt_2_count} mentions.\n"
f"There was  peak in {topdt_1_name} due to following news:The increase in volume is due to Prime Minister Narendra Modi's visit to Chhattisgarh and the launch of multiple development projects, including the dedication of the NMDC Steel Plant in Bastar. The projects are expected to provide employment opportunities and contribute to the reduction of debt burden. The political controversy surrounding the Nagarnar Steel Plant and the statements made by Chief Minister Bhupesh Baghel have also contributed to the spike in volume\n"
                   )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'


    # Sidebar for PowerPoint download settings
    st.sidebar.write("## Download All DataFrames as a PowerPoint File")
    pptx_file_name = st.sidebar.text_input("Enter file name for PowerPoint", "dataframes_presentation.pptx")

    # df = pd.DataFrame(finaldata)

    if st.sidebar.button("Download PowerPoint"):
        # List of DataFrames to save
        pubs_table1 = pubs_table.head(10)
        Jour_table1 = Jour_table.head(10)
        dfs = [Entity_SOV3, sov_dt1, pubs_table1,final_df11 , Jour_table1, PType_Entity, final_df7,Senti_Entity,City_Entity]
        table_titles = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Publication Name with Bureau and Journalist Percentages', 'Journalist Table','PubType Entity Table',
                    'Publication Type with Total Publications and Avg news count','Sentiment Table','Citywise Table'     ]

        
        textbox_text = [ 
             f"The total number of news articles over the period is {total_news_count}. ",
   #     "Among these, IIT Madras dominates the conversation with 28% of the total SOV, indicating significant media coverage and visibility.\n"
#     "IIT Delhi follows IIT Madras, capturing 25% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
#     "IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 21%, 17%, and 7% of the SOV respectively.\n"
#     "IIT Ropar holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e., last in the SOV.\n"
#     "Despite ranking lower in terms of SOV, IIT Ropar's presence indicates some level of visibility and recognition within the online media landscape.",
#         "IIT Ropar witnessed its highest news coverage in Sept -2023, with 7 articles. The news during this period mainly revolved around topics such as:\n"
#     "1.IIT Ropar Placements: Average salary, placed students increase despite Covid slowdown\n"
#     "2.Purohit allows IIT-Ropar to set up campus in Edu City\n"
#     "3.UPES Runway Incubator Signs MoU With IIT Ropars Ihub  Awadh\n"
#     "4.SKUAST-K, IIT Ropar hold 2-day event"
#     , "The leading publications reporting on IIT Ropar and its competitors are Times of India, contributing 561 articles, followed by Economic Times with 467 articles, and The Indian Express with 455 articles.\n"
#     "Among these, publications covering news on IIT Ropar specifically are The Indian Express with 9 articles, followed by Tribune with 9 articles, and Times of India with 7 articles.\n"
#     "The top 10 publications writing articles on IIT Ropar contribute 86% of the total 44 articles.",

# "The Economic Times is the most prominent publication covering Steel Industry, with 1088 news articles , of which 81% of articles where field by Bureaus and remaining 19% were filed by 56 Journalists.MoneyControl ranks second with 413 articles on Steel Industry. A significant 70% of these articles were filed by Bureaus, with the remaining 30% filed by 31 journalists Business Standard is the 3rd top publication covering Steel Industry, with 384 news articles , of which 79% of articles where field by Bureaus and remaining 21% were filed by 27 Journalists.Almost , across all  publications, there is a clear trend of bureau dominance in Steel Industry news coverage.ABP Live published 160 articles on Steel Industry.A dominant 99% of these articles were filed by Bureaus. Only 1% of the articles were contributed by journalists.The Hindu Business Line is the only publication having a higher percentage of journalist-driven articles (65%), indicating a strong focus on original reportingBureau-written articles often includes announcements, press releases, and reports from news agencies, while journalist-written articles often includes detailed analyses, specialized content and exclusive stories.Publications with more journalist-driven content typically deliver more detailed and specialized analysis",
   
    
    f"The total number of news articles over the period is {total_news_count}. There is noticeable fluctuation in coverage from month to month, with periods of both increase and decline.\n"
f"{topdt_1_name} saw the highest number of articles, with {topdt_1_count} mentions. This spike suggests a significant event or increased media focus on Steel Industry during that month.\n"
f"{topdt_2_name} has the 2nd highest peak coverage, with {topdt_2_count} mentions.\n"
f"There was  peak in {topdt_1_name} due to following news:The increase in volume is due to Prime Minister Narendra Modi's visit to Chhattisgarh and the launch of multiple development projects, including the dedication of the NMDC Steel Plant in Bastar. The projects are expected to provide employment opportunities and contribute to the reduction of debt burden. The political controversy surrounding the Nagarnar Steel Plant and the statements made by Chief Minister Bhupesh Baghel have also contributed to the spike in volume\n",
 
                     f"{top_1_name} is the most prominent publication covering Steel Industry, with {top_1_count} news articles. {top_2_name} and {top_3_name} follow, with {top_2_count}  and {top_3_count}  news articles, respectively. Of the {total_news_count} articles having prominent mentions of Steel Industry, top 10 publications, among the 100 publications, contributed {top_4_name} articles ({top_4_per} of all mentions)",
    
     f"{topb_1_name} is the most prominent publication covering Steel Industry, with {topb_1_count} news articles , of which {topb_1_br} of articles where field by Bureaus and remaining {topb_1_jr} were filed by Journalists. {topb_2_name} ranks second with {topb_2_count} articles on Steel Industry. A significant {topb_2_br} of these articles were filed by Bureaus, with the remaining {topb_2_jr} filed by  journalists. {topb_3_name} is the 3rd top publication covering Steel Industry, with {topb_3_count} news articles , of which {topb_3_br} of articles where field by Bureaus and remaining {topb_3_jr} were filed by Journalists. Almost , across all  publications, there is a clear trend of bureau dominance in Steel Industry news coverage. Bureau-written articles often includes announcements, press releases, and reports from news agencies, while journalist-written articles often includes detailed analyses, specialized content and exclusive stories.Publications with more journalist-driven content typically deliver more detailed and specialized analysis",
    
     f"{topj_1_name} from {topjt_1_name} leads with {topj_1_count} articles, indicating a strong focus on Steel Industry topics. {topj_2_name} , from {topjt_2_name}, follows closely with {topj_2_count} articles and {topj_3_name} from {topjt_3_name} with {topj_3_count} articles.\n"
     f"The top 10 journalists have written {topj_4_count} prominent articles on Steel Industry out of {total_news_count} i.e only 8% . Majority of the articles i.e 73% of them were filed by Bureaus.\n",
    
f"{topt_1_name} publications account for the majority of news coverage with {topt_1_count} articles, which is approximately {topt_1_per} of the total. This dominance indicates that the Steel Industry sector is primarily of interest to financial and business-oriented audiences, reflecting the sector's strong ties to economic trends, market conditions, and financial health.{topt_2_name} contribute {topt_2_count} articles, making up about {topt_2_per} of the total coverage.Advertising/PR/Media, Digital First, and Technology publications collectively contribute only 112 articles (approximately 2% of the total).Advertising/PR/Media publications account for 33 articles, reflecting some interest in how the Steel Industry sector is marketed or publicized .Digital First publications contribute 42 articles, indicating limited but notable coverage in online platforms.Technology publications have the least coverage, with only 37 articles. This minimal coverage suggests that while there might be technological innovations in the Steel Industry sector, they have not yet become a major focus for tech media.",
    
f"Dominance of {topav_1_name} News: Despite having only {topav_1_jr} publications, this category accounts for {topav_1_br} of all news articles. It has the highest average news count per publication ({topav_1_av}), significantly higher than other categories. {topav_2_name} News Publications: While they have the highest number of publications ({topav_2_jr}), they account for only {topav_2_br} of news articles. Their average news count ({topav_2_av}) is much lower than {topav_1_name} publications. Other Categories:Advertising/PR/Media, Digital First, and Technology categories each account for only 11% & 5% of total news articles. They have fewer publications and lower average news counts. Overall AverageThe 65 publications collectively produced 5,311 news articles, with an overall average of 369 per publication. However, this average is  skewed by the high output of Business & Financial & General publications.",
    f"{tops_1_name} sentiment dominates the news coverage with {tops_1_count} news items, making up approximately {tops_1_per} of the total. This suggests that the majority of the news in the Steel Industry sector is either informational or factual, lacking any strong positive or negative tone. This might include regular updates, policy changes, or reports.{tops_2_name} sentiment is present in {tops_2_count} news articles, accounting for about {tops_2_per} of the total news count. This is a significant portion, indicating that nearly a quarter of the news in the Steel Industry sector is positive. These articles might include stories about falling interest rates, successful housing policies, increased homeownership, or the introduction of beneficial loan schemes.{tops_3_name} sentiment is the least represented, with only {tops_3_name} news items, making up a mere {tops_3_per} of the total coverage. The low percentage of negative news suggests that the Steel Industry sector has relatively fewer controversies, issues, or negative events reported. Negative sentiment might be related to rising interest rates, fraud cases, or unfavorable changes in housing policies.",
     f"{topct_1_name} leads the news coverage in the Steel Industry sector with {topct_1_count} news articles, accounting for approximately {topct_1_per} of the total news count. This significant concentration indicates that {topct_1_name} is a major city for discussions and developments in the Steel Industry market. {topct_2_name} follows with {topct_2_count} news items, representing around {topct_2_per} of the total. {topct_3_name} and {topct_4_name} also show noteworthy activity, with {topct_3_count} and {topct_4_count} news articles respectively, highlighting their importance in the sector.",
                      ]


        # Create a new PowerPoint presentation
        # prs = Presentation()

        # Loop through each DataFrame and create a new slide with a table
        for i, (df, title) in enumerate(zip(dfs, table_titles)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_table_to_slide(slide, df, title, textbox_text[i])
            # Add image only to the first slide
            if i == 0:  
                img_path4 = generate_bar_chart(dfs[0])  # Generate chart from first DataFrame
                add_image_to_slide(slide, img_path4)
            if i == 1:  
                line_graph_path = generate_line_graph(sov_dt1)  # Generate chart from first DataFrame
                add_image_to_slide(slide, line_graph_path)

            if i == 2:  
                hor_graph_path = generate_horizontal_bar_chart(dfs[2])  # Generate chart from first DataFrame
                add_image_to_slide1(slide, hor_graph_path)

            if i == 3:  
                phor_graph_path = generate_grouped_bar_chartbt(dfs[3])  # Generate chart from first DataFrame
                add_image_to_slide2(slide, phor_graph_path)

            if i == 4:  
                jhor_graph_path = generate_horizontal_bar_chartj(dfs[4])  # Generate chart from first DataFrame
                add_image_to_slide3(slide, jhor_graph_path)

            if i == 5:  
                pthor_graph_path = generate_horizontal_bar_chartpt(dfs[5])  # Generate chart from first DataFrame
                add_image_to_slide4(slide, pthor_graph_path)

            if i == 6:  
                pthor_graph_path = generate_grouped_bar_chartpt(dfs[6])  # Generate chart from first DataFrame
                add_image_to_slide9(slide, pthor_graph_path)

            if i == 7:  
                sthor_graph_path = generate_horizontal_bar_chartst(dfs[7])  # Generate chart from first DataFrame
                add_image_to_slide5(slide, sthor_graph_path)

            if i == 8:  
                cthor_graph_path = generate_horizontal_bar_chartct(dfs[8])  # Generate chart from first DataFrame
                add_image_to_slide6(slide, cthor_graph_path)

            if i == 8:
                wordcloud_path = generate_word_cloud(finaldata)  # Generate word cloud from DataFrame
                add_image_to_slide11(slide, wordcloud_path)



        # Save presentation to BytesIO for download
        pptx_output = io.BytesIO()
        prs.save(pptx_output)
        pptx_output.seek(0)

        # Provide download button
        st.sidebar.download_button(
            label="Download PowerPoint Presentation",
            data=pptx_output,
            file_name=pptx_file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
#         # Download All DataFrames as a Single Excel Sheet
#         st.sidebar.write("## Download All DataFrames as a Single Excel Sheet")
#         file_name_all = st.sidebar.text_input("Enter file name for all DataFrames", "all_dataframes.xlsx")
# #         download_options = st.sidebar.selectbox("Select Download Option:", [ "Complete Dataframes"])
        
#         if st.sidebar.button("Download All DataFrames"):
#             # List of DataFrames to save
#             dfs = [Entity_SOV1, sov_dt, pubs_table, Jour_table, PType_Entity, PP_table, ppe1]
#             comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
#                         'Pub Type and Entity Table', 'Pub Type and Pub Name Table',
#                         'PubType PubName and Entity Table']
            
#             excel_path_all = os.path.join(download_path, file_name_all)
#             multiple_dfs(dfs, 'Tables', excel_path_all, 2, comments)
#             st.sidebar.write(f"All DataFrames saved at {excel_path_all}")

#         # Loop through each dataframe and create a new slide for each one
#         for i, (df, title) in enumerate(zip(dfrs, table_titles)):
#             slide = prs.slides.add_slide(prs.slide_layouts[6])
#             add_table_to_slide(slide, df, title, textbox_text[i])

else:
    st.sidebar.write("No file uploaded yet.")


from wordcloud import WordCloud
from PIL import Image
from fuzzywuzzy import fuzz
import matplotlib.pyplot as plt
# import gensim
# import spacy
# import pyLDAvis.gensim_models
# from gensim.utils import simple_preprocess
# from gensim.models import CoherenceModel
from pprint import pprint
import logging
import warnings
from nltk.corpus import stopwords
# import gensim.corpora as corpora
from io import BytesIO
import nltk

# Download NLTK stopwords
nltk.download('stopwords')

# Set up logging
logging.basicConfig(format='%(asctime)s : %(levelname)s : %(message)s', level=logging.ERROR)

# Ignore warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

# Initialize Spacy 'en' model
# nlp = spacy.load('en_core_web_sm', disable=['parser', 'ner'])


# Define a function to clean the text
def clean(text):
    text = text.lower()
    text = re.sub('[^A-Za-z]+', ' ', text)
    text = re.sub('[,\.!?]', ' ', text)
    return text

# Streamlit app with a sidebar layout
# st.set_page_config(layout="wide")

# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("SimilarNews , Wordcloud and Topic Explorer")

# Sidebar for file upload
st.sidebar.title("Upload a file for Data Analysis")

file = st.sidebar.file_uploader("Upload Excel File", type=["xlsx"])

if file:
    st.sidebar.write("File Uploaded Successfully!")

    # Importing Dataset
    data = pd.read_excel(file)
    
    # data['Text'] = (data['Headline'].astype(str) + data['Opening Text'].astype(str) + data['Hit Sentence'].astype(str))
    # data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
    # data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
    # data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)

    # Define a function to clean the text
    def clean(text):
        text = text.lower()
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        return text

    # Cleaning the text in the Headline column
    data['Cleaned_Headline'] = data['Headline'].apply(clean)

    # Define a function to clean the text
    def cleaned(text):
        # Removes all special characters and numericals leaving the alphabets
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub(r'[[0-9]*]', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        text = re.sub('[\\n]', ' ', text)
        text = re.sub(r'\b\w{1,3}\b', '', text)
        # removing apostrophes
        text = re.sub("'s", '', str(text))
        # removing hyphens
        text = re.sub("-", ' ', str(text))
        text = re.sub(" ", '', str(text))
        # removing quotation marks
        text = re.sub('\"', '', str(text))
        # removing any reference to outside text
        text = re.sub("[\(\[].*?[\)\]]", "", str(text))
        return text

    # Cleaning the text in the review column
    data['Text'] = data['Text'].apply(cleaned)
    data.head()    
    
    st.sidebar.header("Select An Analysis you want to Work On")
    analysis_option = st.sidebar.selectbox(" ", ["Similarity News", "Word Cloud"])

    # Define the 'entities' variable outside of the conditional blocks
    entities = list(data['Entity'].unique())

    # Define an empty 'wordclouds' dictionary
    wordclouds = {}

    if analysis_option == "Similarity News":
        st.header("Similar News")
        st.sidebar.subheader("Similarity News Parameters")
        # Place your parameters for Similarity News here

        # Create a new workbook to store the updated sheets
        updated_workbook = pd.ExcelWriter('Similar_News_Grouped.xlsx', engine='xlsxwriter')

        # Iterate over unique entities
        for entity in entities:
            # Filter data for the current entity
            entity_data = data[data['Entity'] == entity].copy()

            # for each unique value in Cleaned_Headline within the entity
            for headline in entity_data['Cleaned_Headline'].unique():
                # Compute Levenshtein distance and set to True if >= a limit
                entity_data[headline] = entity_data['Cleaned_Headline'].apply(lambda x: fuzz.ratio(x, headline) >= 70)

                # Set a name for the group (the shortest headline)
                m = np.min(entity_data[entity_data[headline] == True]['Cleaned_Headline'])

                # Assign the group
                entity_data.loc[entity_data['Cleaned_Headline'] == headline, 'Similar_Headline'] = m

            # Drop unnecessary columns
            # entity_data.drop(entity_data.columns[36:], axis=1, inplace=True)

            # Sort the dataframe based on the 'Similar_Headline' column
            entity_data.sort_values('Similar_Headline', ascending=True, inplace=True)

            headline_index = entity_data.columns.get_loc('Similar_Headline')

            entity_data = entity_data.iloc[:, :headline_index + 1]

            column_to_delete = entity_data.columns[entity_data.columns.get_loc('Similar_Headline') - 1]  # Get the column name before 'group'
            entity_data = entity_data.drop(column_to_delete, axis=1)
            
            # Define a function to classify news as "Exclusive" or "Not Exclusive" for the current entity
            def classify_exclusivity(row):
                
                entity_name = entity_data['Entity'].iloc[0]  # Get the entity name for the current sheet
                # Check if the entity name is mentioned in either 'Headline' or 'Similar_Headline'
                if entity_name.lower() in row['Headline'].lower() or entity_name.lower() in row['Similar_Headline'].lower():                
                    return "Exclusive"
                else:
                    return "Not Exclusive"
                    
            # Apply the classify_exclusivity function to each row in the current entity's data
            entity_data['Exclusivity'] = entity_data.apply(classify_exclusivity, axis=1)    
            
            
            # Define a dictionary of keywords for each entity
            entity_keywords = {
                        # 'Nothing Tech': ['Nothing','nothing'],
#                         'Asian Paints': ['asian', 'keyword2', 'keyword3'],
            }
            
            # Define a function to qualify entity based on keyword matching
            def qualify_entity(row):                
                entity_name = row['Entity']
                text = row['Headline']   
                
                if entity_name in entity_keywords:
                    keywords = entity_keywords[entity_name]
                    # Check if at least one keyword appears in the text
                    if any(keyword in text for keyword in keywords):
                        return "Qualified"
                
                return "Not Qualified"
            
            # Apply the qualify_entity function to each row in the current entity's data
            entity_data['Qualification'] = entity_data.apply(qualify_entity, axis=1)
            
            # Define a dictionary to map predefined words to topics
            topic_mapping = {
              'Merger': ['merger', 'merges'],
                
              'Acquire': ['acquire', 'acquisition', 'acquires'],
                
              'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
                
               'Business Strategy': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                    'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
                
               'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
                
              'Employee Engagement': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo'],
                
              'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
                
               'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
                
               'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo'], 
                
               'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
                
                'Awards & Recognition': ['award', 'awards'],
                
                'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
            # Add more topics and corresponding words as needed
}
           
             # Define a function to classify headlines into topics
#             def classify_topic(headline):
#                 for topic, words in topic_mapping.items():
#                     for word in words:
#                         if word.lower() in headline.lower():
#                             return topic
#                 return 'Other'  # If none of the predefined words are found, assign 'Other'
            
            # Define a function to classify headlines into topics
            def classify_topic(headline):
                lowercase_headline = headline.lower()
                for topic, words in topic_mapping.items():
                    for word in words:
                        if word in lowercase_headline:
                            return topic
                return 'Other'  # If none of the predefined words are found, assign 'Other'

            
            
            
            # Apply the classify_topic function to each row in the dataframe
            entity_data['Topic'] = entity_data['Headline'].apply(classify_topic)            
            
            # Save the updated sheet to the new workbook
            entity_data.to_excel(updated_workbook, sheet_name=entity, index=False, startrow=0)

            # Create a Word Cloud for the entity based on the Headline column
            wordcloud = WordCloud(width=550, height=400, background_color='white').generate(' '.join(entity_data['Cleaned_Headline']))
            wordclouds[entity] = (wordcloud, entity_data)

        # Save the new workbook
        updated_workbook.close()

        # Provide a download link for the grouped data
        st.markdown("### Download Grouped Data")
        st.markdown(
            f"Download the grouped data as an Excel file: [Similar_News_Grouped.xlsx](data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(open('Similar_News_Grouped.xlsx', 'rb').read()).decode()})"
        )

        # Provide a download link for the grouped data in CSV format
        data_csv = data.to_csv(index=False)
        st.markdown(
            f"Download the original data as a CSV file: [Original_Data.csv](data:text/csv;base64,{base64.b64encode(data_csv.encode()).decode()})"
        )

        # Load the grouped data from the "Similar_News_Grouped.xlsx" file
        grouped_data = pd.read_excel("Similar_News_Grouped.xlsx", sheet_name=None)

        # Data Preview Section
        st.sidebar.subheader("Data Preview")
        entities = list(wordclouds.keys())

        selected_entities = st.sidebar.multiselect("Select Entities to Preview", entities)

        if selected_entities:
            for entity in selected_entities:
                st.header(f"Preview for Entity: {entity}")
                entity_data = grouped_data[entity]
                st.write(entity_data)

    elif analysis_option == "Word Cloud":
        st.header("WordCloud")
        st.sidebar.subheader("Word Cloud Parameters")
        # Place your parameters for Word Cloud here

        # Generate and display word clouds for selected entities
        st.sidebar.title("Word Clouds")

        wordcloud_entity = st.sidebar.selectbox("Select Entity for Word Cloud", entities)

        # Custom Stop Words Section
        st.sidebar.title("Custom Stop Words")

        custom_stopwords = st.sidebar.text_area("Enter custom stop words (comma-separated)", "")
        custom_stopwords = [word.strip() for word in custom_stopwords.split(',')]

        # Widget to adjust word cloud parameters
        wordcloud_size_height = st.slider("Select Word Cloud Size Height", 100, 1000, 400, step=50, key="wordcloud_height")
        wordcloud_size_width = st.slider("Select Word Cloud Size Width", 100, 1000, 400, step=50, key="wordcloud_width")
        wordcloud_max_words = st.slider("Select Max Number of Words", 10, 500, 50)

        if wordcloud_entity:
            st.header(f"Word Cloud for Entity: {wordcloud_entity}")
            # Generate Word Cloud with custom stop words removed
            cleaned_headlines = ' '.join(data[data['Entity'] == wordcloud_entity]['Text'])

            if custom_stopwords:
                for word in custom_stopwords:
                    cleaned_headlines = cleaned_headlines.replace(word, '')

            wordcloud_image = WordCloud(background_color="white", width=wordcloud_size_width, height=wordcloud_size_height, max_font_size=80, max_words=wordcloud_max_words,
                                        colormap='Set1', contour_color='black', contour_width=2, collocations=False).generate(cleaned_headlines)
            
            
            # Create entity_data for the selected entity
            entity_data = data[data['Entity'] == wordcloud_entity]

            # Resize the word cloud image using PIL
            img = Image.fromarray(np.array(wordcloud_image))
            img = img.resize((wordcloud_size_width, wordcloud_size_height))
            
            # Add the entity to the wordclouds dictionary
            wordclouds[wordcloud_entity] = (wordcloud_image, entity_data)

            # Display the resized word cloud image in Streamlit
            st.image(img, caption=f"Word Cloud for Entity: {wordcloud_entity}")

        # Word Cloud Interaction
        if wordcloud_entity:
            st.header(f"Word Cloud Interaction for Entity: {wordcloud_entity}")
            
            # Debugging statements
            st.write("Entities in entities list:", entities)
            st.write("Keys in wordclouds dictionary:", list(wordclouds.keys()))
            
            # Get the selected entity's word cloud
            entity_wordcloud, entity_data = wordclouds.get(wordcloud_entity, (None, None))  # Use .get() to handle missing keys gracefully
            if entity_wordcloud is None:
                st.warning(f"No word cloud found for '{wordcloud_entity}'")
            else:
                words = list(entity_wordcloud.words_.keys())
            

            # Get the selected entity's word cloud
            entity_wordcloud, entity_data = wordclouds[wordcloud_entity]
            words = list(entity_wordcloud.words_.keys())

            word_frequencies = entity_wordcloud.words_
            words_f = list(word_frequencies.keys())

            # Create a list of tuples containing (word, frequency)
            word_frequency_list = [(word, frequency) for word, frequency in word_frequencies.items()]

            # Add this line to preview the words and their frequencies
            st.write("Words and their frequencies:")
            st.write(word_frequency_list)

            # Filter out bigrams from the list of words
            # individual_words = [word for word in words if ' ' not in word]

            # Create a selectbox for the words in the word cloud
            selected_word = st.selectbox("Select a word from the word cloud", words)

            # Find rows where the selected word appears
            matching_rows = entity_data[entity_data['Headline'].str.contains(selected_word, case=False, na=False)]

            # Display the matching rows or a message if no matches are found
            if not matching_rows.empty:
                st.subheader(f"Matching Rows for '{selected_word}':")
                # Function to highlight the selected word
                def highlight_word(text, word):                
                    return re.sub(f'\\b{word}\\b', f'**{word}**', text, flags=re.IGNORECASE)

                # Apply the highlight_word function to the Cleaned_Headline column
                matching_rows['Headline'] = matching_rows.apply(lambda row: highlight_word(row['Headline'], selected_word), axis=1)

                # Display the formatted dataframe
                st.dataframe(matching_rows)

            else:
                st.warning(f"No matching rows found for '{selected_word}'")
    
    
    elif analysis_option == "LDA":
        st.header("Topic Modelling (LDA)")
        selected_entity = st.sidebar.selectbox("Select Entity for LDA", entities)
        
        # Apply LDA to entity_data['Text']
        # Create entity_data for the selected entity
#         entity_data = data[data['Entity'] == wordcloud_entity]

        num_of_topics = st.sidebar.slider("Number of Topics", min_value=1, max_value=20, value=10)
#         per_word_topics = st.sidebar.checkbox("Per Word Topics",
#                                      help="If True, the model also computes a list of topics, sorted in descending order of most likely topics for each word, along with their phi values multiplied by the feature length (i.e. word count).")
        iterations = st.sidebar.number_input("Iterations", min_value=1, value=50,
                                    help="Maximum number of iterations through the corpus when inferring the topic distribution of a corpus.")
        no_words = st.sidebar.number_input("No of Words", min_value=10, value=30,
                                help="Number of words to be displayed in the LDA graph.")       


        cut_off_percentage = st.sidebar.slider("Topic Cutoff Percentage", min_value=0.0, max_value=1.0, value=0.25, step=0.01)
        
        # Run LDA for the selected entity
        if selected_entity:
            st.header(f"Running LDA for entity: {selected_entity}")
            
            # Filter data for the selected entity
            entity_data = data[data['Entity'] == selected_entity]
#             data_text = entity_data['Text'].apply(cleaned).values.tolist()
        

#           uploaded_file = st.sidebar.file_uploader("Upload Excel File", type=["xlsx"])
            file_name = st.sidebar.text_input("Output File Name", f"{selected_entity} new_topics.xlsx") 
            entity_data['Text'] = entity_data['Text'].apply(cleaned)  # Use your text cleaning function if needed
            st.write(entity_data)
            data1 = (entity_data.Text.values.tolist())
#             st.write(data1)
            stop_words = stopwords.words('english')

            # Tokenize and preprocess text data
        
        
            def tokenize_and_preprocess(data1):
                data_words = list(sent_to_words(data1))
                data_words_nostops = remove_stopwords(data_words)
                data_words_bigrams = make_bigrams(data_words_nostops)
                data_lemmatized = lemmatization(data_words_bigrams, allowed_postags=['NOUN', 'ADJ', 'VERB', 'ADV'])
                return data_lemmatized

            def sent_to_words(sentences):
                for sentence in sentences:
                    yield(gensim.utils.simple_preprocess(str(sentence).encode('utf-8'), deacc=True))
            
            data_words = list(sent_to_words(data1))        

            def remove_stopwords(texts):
                return [[word for word in simple_preprocess(str(doc)) if word not in stop_words] for doc in texts]

            def make_bigrams(texts):
                bigram = gensim.models.Phrases(data_words, min_count=5, threshold=100) # higher threshold fewer phrases.
                trigram = gensim.models.Phrases(bigram[data_words], threshold=100)  

                # Faster way to get a sentence clubbed as a trigram/bigram
                bigram_mod = gensim.models.phrases.Phraser(bigram)
                trigram_mod = gensim.models.phrases.Phraser(trigram)
                return [bigram_mod[doc] for doc in texts]

            def lemmatization(texts, allowed_postags=['NOUN', 'ADJ', 'VERB', 'ADV']):
                texts_out = []
                for sent in texts:
                    doc = nlp(" ".join(sent))
                    texts_out.append([token.lemma_ for token in doc if token.pos_ in allowed_postags])
                return texts_out

            data_lemmatized = tokenize_and_preprocess(data1)

            # Build the LDA model
#           @st.cache_data 
            def build_lda_model(corpus, _id2word, num_topics):
                lda_model = gensim.models.ldamodel.LdaModel(
                corpus=corpus,
                id2word=_id2word,
                num_topics=num_topics,
                random_state=100,
                update_every=1,
                chunksize=100,
                passes=10,
                alpha='auto',
                per_word_topics=True,
                iterations=iterations    
        )
                return lda_model

            # Create Dictionary and Corpus
            id2word = corpora.Dictionary(data_lemmatized)
            texts = data_lemmatized
            corpus = [id2word.doc2bow(text) for text in texts]

            # Build the LDA model
            lda_model = build_lda_model(corpus, id2word, num_of_topics)

            # Visualize the topics using pyLDAvis
#           pyLDAvis.enable_notebook()
            vis = pyLDAvis.gensim_models.prepare(lda_model, corpus, id2word,R=no_words)
#           st.title("Topic Modeling Visualization")
#           st.write(vis, use_container_width=True)

            # Create a list to store the assigned topic numbers for each document
            topic_assignments = []

            # Loop through each document in the corpus and assign a topic number
            for doc in corpus:
                # Get the topic probabilities for the document
                topic_probs = lda_model.get_document_topics(doc)

                # Sort the topic probabilities in descending order by probability score
                topic_probs.sort(key=lambda x: x[1], reverse=True)

                # Filter out topics with a contribution less than the cutoff percentage
                topic_probs = [topic for topic in topic_probs if topic[1] >= cut_off_percentage]

                # Get the topic number with the highest probability score
                top_topic_num = topic_probs[0][0] if topic_probs else -1

                # Append the topic number to the list of topic assignments
                topic_assignments.append(top_topic_num)

            # Add the topic assignments to the DataFrame
            entity_data['Topic'] = topic_assignments
    
            st.markdown(
    """
    <style>
    .stApp {
        width: 100%;
    }
    </style>
    """,
            unsafe_allow_html=True,
)

            # Streamlit UI
            st.header(f"LDA Topic Modeling Results {selected_entity}")
            st.dataframe(entity_data)
            st.header(f"Topic Modeling Visualization {selected_entity}")
            st.components.v1.html(pyLDAvis.prepared_data_to_html(vis),width=1500, height=800)

            # Download button for the modified DataFrame
            st.sidebar.markdown("## Download Results")
            st.sidebar.markdown("Click below to download the modified DataFrame:")
    
            # Save the DataFrame to an in-memory Excel file
            excel_buffer = BytesIO()
            with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                entity_data.to_excel(writer, index=False, sheet_name="Sheet1")
        
            # Convert in-memory Excel file to bytes
            excel_data = excel_buffer.getvalue()

            # Provide a download link for the Excel file
            st.sidebar.download_button(
               label="Download Excel",
               data=excel_data,
               file_name=file_name,
               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
)    
  

        else:
            st.write("Please select an entity to run LDA.")
            
            
        # Collect all the documents for each topic
        topic_docs = {}
        for i in range(num_of_topics):
            topic_docs[i] = ' '.join(entity_data[entity_data['Topic'] == i]['Text'].values)
            
        # Add a dropdown for selecting topics from the LDA results
        selected_topic_index = st.sidebar.selectbox("Select a topic", range(num_of_topics))
        
        st.header(f"Word Cloud for Topic: {selected_topic_index}")
        
        # Custom Stop Words Section
        st.sidebar.title("Custom Stop Words")

        custom_stopwords = st.sidebar.text_area("Enter custom stop words (comma-separated)", "")
        custom_stopwords = [word.strip() for word in custom_stopwords.split(',')]
        
        # Widget to adjust word cloud parameters
        wordcloud_size_height = st.sidebar.slider("Select Word Cloud Size Height", 100, 1000, 400, step=50, key="wordcloud_height")
        wordcloud_size_width = st.sidebar.slider("Select Word Cloud Size Width", 100, 1000, 400, step=50, key="wordcloud_width")
        wordcloud_max_words = st.slider("Select Max Number of Words", 10, 500, 100)
        
        # Display the word cloud for the selected topic
        if selected_topic_index is not None:
            selected_topic_text = topic_docs[selected_topic_index]
            wordcloud_image = WordCloud(font_path="D:\Akshay.Annaldasula\OneDrive - Adfactors PR Pvt Ltd\Downloads\Helvetica.ttf", background_color='white',colormap='Set1', contour_color='black', contour_width=2, collocations=False,max_font_size=80,width=wordcloud_size_width, height=wordcloud_size_height,stopwords=custom_stopwords, max_words=wordcloud_max_words).generate(selected_topic_text)
#             plt.figure(figsize=(5, 5))
#             plt.imshow(wordcloud, interpolation='bilinear')
#             plt.axis('off')
#             # Display the word cloud using Streamlit's st.pyplot() method
#             fig, ax = plt.subplots()
#             ax.imshow(wordcloud, interpolation='bilinear')
#             ax.axis('off')
#             st.pyplot(fig)
            # Resize the word cloud image using PIL
            img = Image.fromarray(np.array(wordcloud_image))
            img = img.resize((wordcloud_size_width, wordcloud_size_height))
            # Display the resized word cloud image in Streamlit
            st.image(img)
            
        # Iterate over unique topics
        # Assuming 'entity_data' DataFrame already has a column named 'Topics'
        topics = entity_data['Topic'].unique()
        
        for topic in topics:
            # Filter data for the current topic
            topic_data = entity_data[entity_data['Topic'] == topic].copy()

            # for each unique value in Cleaned_Headline within the topic
            for headline in topic_data['Cleaned_Headline'].unique():
                # Compute Levenshtein distance and set to True if >= a limit
                topic_data[headline] = topic_data['Cleaned_Headline'].apply(lambda x: fuzz.ratio(x, headline) >= 70)

                # Set a name for the group (the shortest headline)
                m = np.min(topic_data[topic_data[headline] == True]['Cleaned_Headline'])

                # Assign the group
                topic_data.loc[topic_data['Cleaned_Headline'] == headline, 'Similar_Headline'] = m

            # Drop unnecessary columns
            topic_data.drop(topic_data.columns[36:], axis=1, inplace=True)

            # Sort the dataframe based on the 'Similar_Headline' column
            topic_data.sort_values('Similar_Headline', ascending=True, inplace=True)

            headline_index = topic_data.columns.get_loc('Similar_Headline')

            topic_data = topic_data.iloc[:, :headline_index + 1]

            column_to_delete = topic_data.columns[topic_data.columns.get_loc('Similar_Headline') - 1]  # Get the column name before 'group'
            topic_data = topic_data.drop(column_to_delete, axis=1)

            # Display a preview of the data for the current topic
            if selected_topic_index == topic:
                st.subheader(f"Preview of Similar News for Topic {topic}")
                st.write(topic_data)      
                
                
                # Add a button to download the data as an Excel file
                if st.sidebar.button('Download Topics Data as Excel'):
                    file_name = f"topic_{topic}.xlsx"
#                     excel_data = BytesIO()
#                     topic_data.to_excel(file_name, index=False)
#                     st.markdown(get_download_link(file_name), unsafe_allow_html=True)
#                     # Create a link to download the Excel file for data
                    excel_path = os.path.join(download_path, file_name)
                    topic_data.to_excel(excel_path, index=False)
                    st.sidebar.write(f"Excel file saved at {excel_path}")
